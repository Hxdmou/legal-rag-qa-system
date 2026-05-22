from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

def fill_cell(cell, text, font_size=9, bold=False):
    cell.text_frame.clear()
    para = cell.text_frame.add_paragraph()
    run = para.add_run()
    run.text = str(text) if text else "---"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(255, 255, 255)
    run.font.name = '微软雅黑'
    para.alignment = PP_ALIGN.CENTER
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

def update_with_new_models(prs):
    # 更新项目5 - 加入DeepSeek V5.0和GPT 5.0
    if len(prs.slides) > 4:
        slide = prs.slides[4]
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                # 更新第7行第1列 - 加入新模型热点
                fill_cell(table.cell(7, 0), "新模型爆发")
                fill_cell(table.cell(7, 1), "高")
                fill_cell(table.cell(7, 2), "颠覆性")
                fill_cell(table.cell(7, 3), "DeepSeek V5.0+GPT 5.0")
                fill_cell(table.cell(7, 4), "极高")
                fill_cell(table.cell(7, 5), "高")
                fill_cell(table.cell(7, 6), "极热")
                fill_cell(table.cell(7, 7), "激烈")
                fill_cell(table.cell(7, 8), "8倍")
                fill_cell(table.cell(7, 9), "中")
                fill_cell(table.cell(7, 10), "1-2年")
                fill_cell(table.cell(7, 11), "紧缺")
                fill_cell(table.cell(7, 12), "抢先布局")
                fill_cell(table.cell(7, 13), "12个月回本")
                fill_cell(table.cell(7, 14), "技术领先")
                fill_cell(table.cell(7, 15), "P0")
                fill_cell(table.cell(7, 16), "10")
                break
    
    # 更新项目10 - 市场预测加入新模型影响
    if len(prs.slides) > 9:
        slide = prs.slides[9]
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                # 更新2027年技术突破加入GPT 5.0
                fill_cell(table.cell(2, 5), "GPT 5.0/DeepSeek V5.0")
                # 更新2028年技术突破
                fill_cell(table.cell(3, 5), "GPT 5.5/AGI")
                # 更新2029年加入DeepSeek V6.0
                fill_cell(table.cell(4, 5), "DeepSeek V6.0/神经连接")
                break
    
    # 更新项目11 - 技术评估加入新模型
    if len(prs.slides) > 10:
        slide = prs.slides[10]
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                # 更新推理能力测试环境加入新模型
                fill_cell(table.cell(1, 3), "GPT 5.0/DeepSeek V5.0")
                # 更新多模态测试环境
                fill_cell(table.cell(2, 3), "GPT 5.0/Gemini 2.5")
                break

def main():
    ppt_path = r"f:\个人作品\legal-rag-qa-system\A2A_PROTOCOL_AI_AGENT_V14_终极完整版.pptx"
    prs = Presentation(ppt_path)
    
    print("🔄 正在补充DeepSeek V5.0和GPT 5.0最新信息...")
    update_with_new_models(prs)
    
    output_path = r"f:\个人作品\legal-rag-qa-system\A2A_PROTOCOL_AI_AGENT_V14_终极完整版.pptx"
    prs.save(output_path)
    print(f"✅ 已完成！保存到: {output_path}")
    print("\n📌 补充内容:")
    print("   ✅ DeepSeek V5.0 - 2026年AI热点")
    print("   ✅ GPT 5.0 - 2027年技术突破")
    print("   ✅ GPT 5.5 - 2028年技术突破")
    print("   ✅ DeepSeek V6.0 - 2029年技术突破")

if __name__ == "__main__":
    main()