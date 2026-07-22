from pptx import Presentation

def check_all_projects():
    print("=" * 120)
    print("🔍 检查 docs/presentations/A2A_ENTERPRISE_PPT_V12.pptx 所有项目")
    print("=" * 120)
    
    ppt_path = r"f:\个人作品\legal-rag-qa-system\docs\presentations\A2A_ENTERPRISE_PPT_V12.pptx"
    prs = Presentation(ppt_path)
    
    for slide_idx, slide in enumerate(prs.slides):
        print(f"\n📄 幻灯片 {slide_idx + 1}:")
        
        # 获取幻灯片标题
        title = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                if shape.text.startswith("项目"):
                    title = shape.text[:20]
                    break
        
        # 查找表格
        table_found = False
        for shape in slide.shapes:
            if shape.has_table:
                table_found = True
                table = shape.table
                print(f"   标题: {title}")
                print(f"   表格: {len(table.rows)}行 × {len(table.columns)}列")
                
                # 打印第一列内容
                first_col = []
                for i in range(min(5, len(table.rows))):
                    cell_text = table.cell(i, 0).text.strip()
                    first_col.append(cell_text)
                print(f"   第一列: {first_col}...")
                break
        
        if not table_found:
            print(f"   标题: {title}")
            print("   └─ 无表格")

if __name__ == "__main__":
    check_all_projects()