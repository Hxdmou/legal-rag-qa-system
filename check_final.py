from pptx import Presentation

ppt_path = r"f:\个人作品\legal-rag-qa-system\A2A_PROTOCOL_AI_AGENT_V14_最终完整版.pptx"
prs = Presentation(ppt_path)

print("=" * 80)
print("验证最终完整版PPT")
print("=" * 80)

total_empty = 0

for slide_idx, slide in enumerate(prs.slides):
    print(f"\n📊 幻灯片 {slide_idx + 1}:")
    has_table = False
    for shape in slide.shapes:
        if shape.has_table:
            has_table = True
            table = shape.table
            rows, cols = len(table.rows), len(table.columns)
            print(f"  表格: {rows}行 x {cols}列")
            
            empty_count = 0
            for i in range(rows):
                for j in range(cols):
                    cell_text = table.cell(i, j).text.strip()
                    if not cell_text or cell_text == "------" or cell_text.startswith("........"):
                        empty_count += 1
            
            total_empty += empty_count
            if empty_count > 0:
                print(f"    ❌ 共 {empty_count} 个空单元格")
            else:
                print("    ✅ 所有单元格已填充")
    
    if not has_table:
        print("  (无表格)")

print(f"\n📊 总计: {total_empty} 个空单元格")
if total_empty == 0:
    print("🎉 所有单元格都已填充完成！")
else:
    print("⚠️ 还有空单元格需要处理")