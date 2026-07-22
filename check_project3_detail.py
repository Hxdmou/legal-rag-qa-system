from pptx import Presentation

def check_project3():
    ppt_path = r"f:\个人作品\legal-rag-qa-system\A2A_PROTOCOL_AI_AGENT_V14_终极完整版.pptx"
    prs = Presentation(ppt_path)
    
    if len(prs.slides) > 2:
        slide = prs.slides[2]
        print("=" * 120)
        print("📌 项目3当前内容检查")
        print("=" * 120)
        
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                print(f"\n表格尺寸: {len(table.rows)}行 × {len(table.columns)}列\n")
                
                for i in range(len(table.rows)):
                    row_data = []
                    for j in range(len(table.columns)):
                        cell_text = table.cell(i, j).text.strip()
                        row_data.append(cell_text)
                        if not cell_text:
                            print(f"⚠️ 警告: 第{i+1}行第{j+1}列为空！")
                    
                    print(f"第{i+1}行: {' | '.join(row_data[:5])}...")
                    if i == 0:
                        print(f"  列名: {', '.join(row_data)}")
                
                break
    
    print("\n" + "=" * 120)
    print("📌 项目3应该包含的内容（10层架构）:")
    print("=" * 120)
    print("L1: 传输层 - 消息路由+负载均衡")
    print("L2: 协议层 - 通信协议+认证")
    print("L3: 智能层 - AI决策+RAG")
    print("L4: 应用层 - 业务逻辑+工作流")
    print("L5: 表现层 - 用户交互+UI")
    print("L6: 数据层 - 存储计算+缓存")
    print("L7: 安全层 - 身份认证+审计")
    print("L8: 监控层 - 可观测性+告警")
    print("L9: 网关层 - 流量管理+限流")
    print("L10: 边缘层 - 边缘计算+AI")

if __name__ == "__main__":
    check_project3()