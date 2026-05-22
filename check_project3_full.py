from pptx import Presentation

def check_project3_display():
    ppt_path = r"f:\个人作品\legal-rag-qa-system\A2A_PROTOCOL_AI_AGENT_V14_终极完整版.pptx"
    prs = Presentation(ppt_path)
    
    print("=" * 120)
    print("📌 项目3完整内容检查")
    print("=" * 120)
    
    if len(prs.slides) > 2:
        slide = prs.slides[2]
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                print(f"\n表格尺寸: {len(table.rows)}行 × {len(table.columns)}列\n")
                
                # 显示每一行的完整内容
                for i in range(len(table.rows)):
                    row_data = []
                    for j in range(len(table.columns)):
                        cell_text = table.cell(i, j).text.strip()
                        row_data.append(cell_text)
                        if not cell_text:
                            print(f"⚠️ 警告: 第{i+1}行第{j+1}列为空！")
                    
                    # 显示整行内容
                    print(f"\n第{i+1}行:")
                    for j, col_name in enumerate(row_data):
                        print(f"  列{j+1}: {col_name}")
                
                # 检查是否有L10
                has_l10 = False
                for i in range(len(table.rows)):
                    if "L10" in table.cell(i, 0).text:
                        has_l10 = True
                        print(f"\n✅ 找到L10: {table.cell(i, 1).text}")
                        break
                
                if not has_l10:
                    print("\n❌ 未找到L10边缘层！")
                
                break

if __name__ == "__main__":
    check_project3_display()