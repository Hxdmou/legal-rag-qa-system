from pptx import Presentation

def check_v12():
    print("=" * 120)
    print("🔍 检查 A2A_ENTERPRISE_PPT_V12.pptx")
    print("=" * 120)
    
    ppt_path = r"f:\个人作品\legal-rag-qa-system\A2A_ENTERPRISE_PPT_V12.pptx"
    
    prs = Presentation(ppt_path)
    print(f"幻灯片总数: {len(prs.slides)}")
    
    # 检查幻灯片3（项目3）
    if len(prs.slides) > 2:
        slide = prs.slides[2]
        print(f"\n📄 幻灯片3 (项目3):")
        
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                print(f"   表格: {len(table.rows)}行 × {len(table.columns)}列")
                
                # 打印列标题
                headers = []
                for j in range(len(table.columns)):
                    headers.append(table.cell(0, j).text.strip())
                print(f"   列标题: {headers}")
                
                # 打印所有行的第一列
                rows = []
                for i in range(len(table.rows)):
                    rows.append(table.cell(i, 0).text.strip())
                print(f"   行内容: {rows}")
                
                break

if __name__ == "__main__":
    check_v12()