from pptx import Presentation

def compare_ppts():
    ppt1_path = r"f:\个人作品\legal-rag-qa-system\A2A_PROTOCOL_AI_AGENT_V14_完整填充版.pptx"
    ppt2_path = r"f:\个人作品\legal-rag-qa-system\A2A_PROTOCOL_AI_AGENT_V14_终极完整版.pptx"
    
    print("=" * 120)
    print("📊 两个PPT文件对比")
    print("=" * 120)
    
    prs1 = Presentation(ppt1_path)
    prs2 = Presentation(ppt2_path)
    
    print(f"\n文件1: {ppt1_path}")
    print(f"  幻灯片数: {len(prs1.slides)}")
    
    print(f"\n文件2: {ppt2_path}")
    print(f"  幻灯片数: {len(prs2.slides)}")
    
    # 检查项目3
    print("\n" + "=" * 120)
    print("📌 项目3对比")
    print("=" * 120)
    
    if len(prs1.slides) > 2:
        slide1 = prs1.slides[2]
        for shape in slide1.shapes:
            if shape.has_table:
                table1 = shape.table
                print(f"\n文件1项目3: {len(table1.rows)}行 × {len(table1.columns)}列")
                
                # 检查层级
                layers1 = []
                for i in range(1, min(12, len(table1.rows))):
                    layer_id = table1.cell(i, 0).text.strip()
                    layers1.append(layer_id)
                print(f"  层级: {', '.join(layers1)}")
                break
    
    if len(prs2.slides) > 2:
        slide2 = prs2.slides[2]
        for shape in slide2.shapes:
            if shape.has_table:
                table2 = shape.table
                print(f"\n文件2项目3: {len(table2.rows)}行 × {len(table2.columns)}列")
                
                # 检查层级
                layers2 = []
                for i in range(1, min(12, len(table2.rows))):
                    layer_id = table2.cell(i, 0).text.strip()
                    layers2.append(layer_id)
                print(f"  层级: {', '.join(layers2)}")
                break
    
    print("\n" + "=" * 120)
    print("💡 重要提示")
    print("=" * 120)
    print("请确保您打开的是: A2A_PROTOCOL_AI_AGENT_V14_终极完整版.pptx")
    print("而不是: A2A_PROTOCOL_AI_AGENT_V14_完整填充版.pptx")

if __name__ == "__main__":
    compare_ppts()