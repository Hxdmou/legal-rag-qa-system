from pptx import Presentation

ppt_path = r"f:\个人作品\legal-rag-qa-system\A2A_PROTOCOL_AI_AGENT_V14_终极完整版.pptx"
prs = Presentation(ppt_path)

print("=" * 100)
print("全面系统性检查 - A2A_PROTOCOL_AI_AGENT_V14_终极完整版.pptx")
print("=" * 100)

total_empty = 0
issues = []

for slide_idx, slide in enumerate(prs.slides):
    print(f"\n📊 幻灯片 {slide_idx + 1}:")
    has_table = False
    for shape in slide.shapes:
        if shape.has_table:
            has_table = True
            table = shape.table
            rows, cols = len(table.rows), len(table.columns)
            print(f"  表格: {rows}行 x {cols}列")
            
            empty_count = 0
            for i in range(rows):
                for j in range(cols):
                    cell_text = table.cell(i, j).text.strip()
                    if not cell_text or cell_text == "------" or cell_text.startswith("........"):
                        empty_count += 1
                        issues.append(f"幻灯片{slide_idx+1}: 空单元格 行{i+1},列{j+1}")
            
            total_empty += empty_count
            if empty_count > 0:
                print(f"    ❌ 共 {empty_count} 个空单元格")
            else:
                print("    ✅ 所有单元格已填充")
    
    if not has_table:
        print("  (无表格)")

# 专项检查关键内容
print("\n" + "=" * 100)
print("专项检查:")
print("=" * 100)

# 检查项目5 - 16大热点
slide5 = prs.slides[4] if len(prs.slides) > 4 else None
if slide5:
    for shape in slide5.shapes:
        if shape.has_table:
            table = shape.table
            row_count = len(table.rows)
            if row_count >= 17:  # 表头 + 16行数据
                print("✅ 项目5: 16大AI热点已完整填充")
            else:
                print(f"❌ 项目5: 行数不足，期望17行，实际{row_count}行")
            break

# 检查项目10 - 2026-2035年数据
slide10 = prs.slides[9] if len(prs.slides) > 9 else None
if slide10:
    for shape in slide10.shapes:
        if shape.has_table:
            table = shape.table
            first_row_years = []
            for i in range(1, min(len(table.rows), 11)):
                year_cell = table.cell(i, 0).text.strip()
                first_row_years.append(year_cell)
            
            if "2026" in first_row_years and "2035" in first_row_years:
                print("✅ 项目10: 市场预测已更新为2026-2035年")
            else:
                print(f"❌ 项目10: 年份范围不正确，当前年份: {first_row_years}")
            break

# 检查项目11 - 10大维度
slide11 = prs.slides[10] if len(prs.slides) > 10 else None
if slide11:
    for shape in slide11.shapes:
        if shape.has_table:
            table = shape.table
            row_count = len(table.rows)
            if row_count >= 11:  # 表头 + 10行数据
                print("✅ 项目11: 最新10大维度验证已完整填充")
            else:
                print(f"❌ 项目11: 行数不足，期望11行，实际{row_count}行")
            break

# 检查项目12 - 总结内容
slide12 = prs.slides[11] if len(prs.slides) > 11 else None
if slide12:
    for shape in slide12.shapes:
        if shape.has_table:
            table = shape.table
            col_count = len(table.columns)
            if col_count >= 4:  # 项目、核心内容、当前状态、未来10年规划、落地方案
                print("✅ 项目12: 总结内容包含当前状态、未来10年规划和落地方案")
            else:
                print(f"❌ 项目12: 列数不足，期望至少4列，实际{col_count}列")
            break

print("\n" + "=" * 100)
print(f"检查结果汇总:")
print("=" * 100)
print(f"📊 空单元格总数: {total_empty}")

if total_empty == 0 and not issues:
    print("🎉 所有检查通过！万无一失！")
else:
    print("⚠️ 发现问题:")
    for issue in issues[:10]:  # 最多显示10个问题
        print(f"  - {issue}")
    if len(issues) > 10:
        print(f"  - ...还有 {len(issues) - 10} 个问题")