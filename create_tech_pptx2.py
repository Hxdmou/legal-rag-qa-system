from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor, ColorFormat
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn

def set_cell_style(cell, text, font_size=9, bold=False, text_color=RGBColor(255, 255, 255), fill_color=RGBColor(30, 30, 60)):
    """设置单元格样式"""
    # 设置背景色
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill_color
    
    # 设置边框
    for edge in ('top', 'bottom', 'left', 'right'):
        cell.border.__getattr__(edge).color.rgb = RGBColor(60, 60, 100)
        cell.border.__getattr__(edge).width = Pt(1)
    
    # 设置文字
    cell.text_frame.clear()
    para = cell.text_frame.add_paragraph()
    run = para.add_run()
    run.text = str(text) if text else "---"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = text_color
    run.font.name = '微软雅黑'
    run._element.rPr.fonts.set(qn('w:eastAsia'), '微软雅黑')
    para.alignment = PP_ALIGN.CENTER
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

def add_tech_elements(slide):
    """添加AI科技元素"""
    # 添加深蓝色背景
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(10, 10, 30)
    
    # 添加装饰性圆形（AI节点）
    left = Inches(8.5)
    top = Inches(5)
    width = height = Inches(0.6)
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0, 200, 255)
    shape.line.fill.background()
    
    # 添加六边形（AI芯片元素）
    left = Inches(0.8)
    top = Inches(5)
    width = height = Inches(0.5)
    hexagon = slide.shapes.add_shape(MSO_SHAPE.HEXAGON, left, top, width, height)
    hexagon.fill.solid()
    hexagon.fill.fore_color.rgb = RGBColor(150, 100, 255)
    hexagon.line.fill.background()

def create_tech_pptx():
    prs = Presentation()
    
    # 项目1：封面
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    # 封面背景 - 深蓝色
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(5, 5, 25)
    
    title = slide.shapes.title
    title.text = "A2A协议与AI智能体"
    for para in title.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(36)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 200, 255)
            run.font.name = '微软雅黑'
    
    subtitle = slide.placeholders[1]
    subtitle.text = "2026年度完整版"
    for para in subtitle.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(18)
            run.font.color.rgb = RGBColor(200, 200, 255)
            run.font.name = '微软雅黑'
    
    # 添加AI发光元素
    left = Inches(3.5)
    top = Inches(4)
    width = height = Inches(2.5)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0, 80, 120)
    circle.line.fill.background()
    
    # 项目2：价值维度
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    add_tech_elements(slide)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.8))
    title_text = title_box.text_frame
    title_para = title_text.add_paragraph()
    title_run = title_para.add_run()
    title_run.text = "项目2：价值维度"
    title_run.font.size = Pt(24)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(0, 200, 255)
    title_run.font.name = '微软雅黑'
    title_para.alignment = PP_ALIGN.CENTER
    
    # 表格
    rows, cols = 16, 9
    left = Inches(0.3)
    top = Inches(1.5)
    width = Inches(10.4)
    height = Inches(5.5)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    data = [
        ["价值维度", "互操作性", "弹性扩展", "安全合规", "可观测性", "智能协作", "跨云部署", "成本优化", "快速迭代"],
        ["互操作性", "极高", "高", "中", "高", "极高", "高", "中", "高"],
        ["弹性扩展", "高", "极高", "中", "高", "高", "极高", "高", "极高"],
        ["安全合规", "中", "中", "极高", "高", "高", "高", "中", "中"],
        ["可观测性", "高", "高", "高", "极高", "高", "高", "中", "高"],
        ["智能协作", "极高", "高", "高", "高", "极高", "高", "中", "高"],
        ["跨云部署", "高", "极高", "高", "高", "高", "极高", "高", "高"],
        ["成本优化", "中", "高", "低", "中", "中", "高", "极高", "高"],
        ["快速迭代", "高", "极高", "中", "高", "高", "高", "高", "极高"],
        ["生态扩展", "极高", "高", "中", "高", "极高", "高", "中", "中"],
        ["数据价值", "高", "高", "极高", "高", "高", "高", "中", "中"],
        ["边缘计算", "高", "极高", "高", "高", "高", "极高", "高", "高"],
        ["AI协同", "极高", "高", "高", "高", "极高", "高", "中", "高"],
        ["Agent自主", "高", "高", "高", "高", "极高", "高", "中", "高"],
        ["多模态交互", "高", "中", "中", "高", "极高", "中", "低", "高"],
        ["综合评分", "9.5", "9.2", "8.8", "9.0", "9.8", "9.3", "8.0", "9.1"],
    ]
    
    for i in range(len(data)):
        for j in range(len(data[i])):
            if i == 0:
                set_cell_style(table.cell(i, j), data[i][j], 11, True, RGBColor(0, 200, 255), RGBColor(20, 30, 50))
            elif j == 0:
                set_cell_style(table.cell(i, j), data[i][j], 10, True, RGBColor(200, 200, 255), RGBColor(15, 25, 45))
            else:
                set_cell_style(table.cell(i, j), data[i][j], 9, False, RGBColor(220, 220, 240), RGBColor(10, 15, 30))
    
    # 项目3：协议层级
    slide = prs.slides.add_slide(slide_layout)
    add_tech_elements(slide)
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.8))
    title_text = title_box.text_frame
    title_para = title_text.add_paragraph()
    title_run = title_para.add_run()
    title_run.text = "项目3：协议层级"
    title_run.font.size = Pt(24)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(0, 200, 255)
    title_run.font.name = '微软雅黑'
    title_para.alignment = PP_ALIGN.CENTER
    
    table = slide.shapes.add_table(14, 9, left, top, width, height).table
    data = [
        ["协议层级", "业务价值", "技术亮点", "成功指标", "适用场景", "性能指标", "投资回报", "技术架构", "关键挑战"],
        ["应用层", "高", "多模态融合", "95%", "企业服务", "<10ms", "中规模", "GPT 5.0+DeepSeek", "高"],
        ["消息层", "高", "异步通信", "99%", "微服务", "<50ms", "中规模", "Event-Driven", "中高"],
        ["传输层", "高", "HTTP/3+QUIC", "99.99%", "大规模", "<5ms", "大规模", "gRPC+Envoy", "中"],
        ["安全层", "极高", "零信任架构", "99.999%", "金融医疗", "<15ms", "高成本", "mTLS+OPA", "中高"],
        ["发现层", "中高", "服务发现", "99%", "云原生", "<30ms", "中规模", "Consul+DNS", "低"],
        ["事件层", "高", "实时事件流", "99.9%", "IoT", "<20ms", "大规模", "Kafka+Pulsar", "中"],
        ["数据层", "极高", "分布式存储", "99.999%", "大数据", "<10ms", "大规模", "Redis+PostgreSQL", "中"],
        ["控制层", "中", "流量控制", "99%", "网关", "<25ms", "中规模", "Kong+APISIX", "低"],
        ["编排层", "高", "工作流编排", "99%", "业务流程", "<100ms", "大规模", "Flowise+Airflow", "中高"],
        ["语义层", "高", "AI理解", "95%", "智能助理", "<150ms", "大规模", "GPT 5.0+DeepSeek", "高"],
        ["Agent通信", "极高", "A2A协议", "99.9%", "多智能体", "<50ms", "大规模", "LangChain+gRPC", "中高"],
        ["Agent发现", "中", "智能发现", "98%", "智能生态", "<100ms", "中规模", "AI-Powered", "中"],
        ["边缘层(L10)", "高", "边缘AI推理", "99.9%", "边缘计算", "<10ms", "大规模", "EdgeAI+FastAPI", "中"],
    ]
    
    for i in range(len(data)):
        for j in range(len(data[i])):
            if i == 0:
                set_cell_style(table.cell(i, j), data[i][j], 11, True, RGBColor(0, 200, 255), RGBColor(20, 30, 50))
            elif j == 0:
                set_cell_style(table.cell(i, j), data[i][j], 10, True, RGBColor(200, 200, 255), RGBColor(15, 25, 45))
            else:
                set_cell_style(table.cell(i, j), data[i][j], 9, False, RGBColor(220, 220, 240), RGBColor(10, 15, 30))
    
    # 项目4：安全组件
    slide = prs.slides.add_slide(slide_layout)
    add_tech_elements(slide)
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.8))
    title_text = title_box.text_frame
    title_para = title_text.add_paragraph()
    title_run = title_para.add_run()
    title_run.text = "项目4：安全组件"
    title_run.font.size = Pt(24)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(0, 200, 255)
    title_run.font.name = '微软雅黑'
    title_para.alignment = PP_ALIGN.CENTER
    
    table = slide.shapes.add_table(17, 9, left, top, width, height).table
    data = [
        ["安全组件", "功能描述", "技术实现", "安全等级", "部署方式", "性能影响", "成本估算", "适用场景", "推荐指数"],
        ["API网关", "流量控制+认证", "Kong/APISIX", "高", "边缘", "<10ms", "中", "企业级", "9"],
        ["WAF防火墙", "Web攻击防护", "Cloudflare", "极高", "边缘", "<5ms", "高", "互联网", "10"],
        ["OAuth 2.0", "身份认证", "Keycloak", "高", "集中", "<15ms", "中", "企业级", "9"],
        ["mTLS", "双向认证", "Istio", "极高", "服务网格", "<20ms", "中高", "金融", "10"],
        ["RBAC/ABAC", "权限控制", "OPA", "高", "集中", "<5ms", "低", "企业级", "9"],
        ["DDoS防护", "流量清洗", "Cloudflare", "极高", "边缘", "<1ms", "高", "互联网", "10"],
        ["数据加密", "静态+传输", "AES-256", "极高", "全链路", "<3ms", "低", "金融医疗", "10"],
        ["审计日志", "操作追踪", "ELK", "高", "集中", "<100ms", "中", "企业级", "9"],
        ["威胁检测", "实时告警", "SIEM", "高", "集中", "<50ms", "高", "安全运营", "9"],
        ["密钥管理", "密钥轮换", "Vault", "极高", "集中", "<10ms", "中", "企业级", "10"],
        ["数据脱敏", "敏感数据保护", "Privitar", "高", "数据层", "<50ms", "中", "金融医疗", "9"],
        ["访问审计", "权限审计", "Azure AD", "高", "集中", "<20ms", "中", "企业级", "8"],
        ["Agent安全", "智能体认证", "A2A-Security", "高", "分布式", "<30ms", "中", "AI系统", "9"],
        ["AI伦理", "偏见检测", "IBM AI", "中", "AI层", "<100ms", "中", "AI应用", "8"],
        ["零信任架构", "永不信任", "BeyondCorp", "极高", "全架构", "<50ms", "高", "企业级", "10"],
        ["隐私计算", "安全多方计算", "MPC/FHE", "极高", "数据层", "<500ms", "极高", "金融", "9"],
    ]
    
    for i in range(len(data)):
        for j in range(len(data[i])):
            if i == 0:
                set_cell_style(table.cell(i, j), data[i][j], 11, True, RGBColor(0, 200, 255), RGBColor(20, 30, 50))
            elif j == 0:
                set_cell_style(table.cell(i, j), data[i][j], 10, True, RGBColor(200, 200, 255), RGBColor(15, 25, 45))
            else:
                set_cell_style(table.cell(i, j), data[i][j], 9, False, RGBColor(220, 220, 240), RGBColor(10, 15, 30))
    
    # 项目5：部署策略
    slide = prs.slides.add_slide(slide_layout)
    add_tech_elements(slide)
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.8))
    title_text = title_box.text_frame
    title_para = title_text.add_paragraph()
    title_run = title_para.add_run()
    title_run.text = "项目5：部署策略"
    title_run.font.size = Pt(24)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(0, 200, 255)
    title_run.font.name = '微软雅黑'
    title_para.alignment = PP_ALIGN.CENTER
    
    table = slide.shapes.add_table(17, 9, left, top, width, height).table
    data = [
        ["部署策略", "适用场景", "技术方案", "成本估算", "复杂度", "可靠性", "扩展性", "推荐度", "风险等级"],
        ["多活部署", "高可用", "K8s+DNS", "高", "中", "99.99%", "高", "9", "低"],
        ["自动扩缩容", "弹性负载", "HPA/VPA", "中", "低", "99.9%", "极高", "10", "低"],
        ["蓝绿发布", "零停机", "Argo Rollouts", "中", "中", "99.99%", "高", "9", "低"],
        ["异地容灾", "灾难恢复", "跨区域复制", "极高", "高", "99.999%", "高", "10", "低"],
        ["边缘部署", "低延迟", "K3s+CDN", "中高", "中", "99.9%", "高", "9", "中"],
        ["金丝雀发布", "灰度测试", "Argo Rollouts", "中", "中", "99.99%", "高", "9", "低"],
        ["Serverless", "按需付费", "AWS Lambda", "低", "低", "99.9%", "极高", "8", "中"],
        ["混合云部署", "多云管理", "Crossplane", "高", "高", "99.9%", "极高", "9", "中"],
        ["GitOps", "自动化运维", "Argo CD", "中", "中", "99.9%", "高", "10", "低"],
        ["服务网格", "流量管理", "Istio", "中高", "高", "99.99%", "高", "9", "中"],
        ["自动化运维", "DevOps", "GitLab CI", "中", "中", "99.9%", "高", "9", "低"],
        ["指标监控", "性能追踪", "Prometheus", "中", "低", "-", "高", "10", "低"],
        ["分布式追踪", "链路分析", "Jaeger", "中", "中", "-", "高", "9", "低"],
        ["配置管理", "动态配置", "ConfigMap", "低", "低", "99.99%", "高", "10", "低"],
        ["安全扫描", "漏洞检测", "Trivy", "低", "低", "-", "高", "10", "低"],
        ["综合评估", "企业级部署", "K8s+Istio+Argo", "高", "高", "99.99%", "极高", "10", "低"],
    ]
    
    for i in range(len(data)):
        for j in range(len(data[i])):
            if i == 0:
                set_cell_style(table.cell(i, j), data[i][j], 11, True, RGBColor(0, 200, 255), RGBColor(20, 30, 50))
            elif j == 0:
                set_cell_style(table.cell(i, j), data[i][j], 10, True, RGBColor(200, 200, 255), RGBColor(15, 25, 45))
            else:
                set_cell_style(table.cell(i, j), data[i][j], 9, False, RGBColor(220, 220, 240), RGBColor(10, 15, 30))
    
    # 项目6：优化维度
    slide = prs.slides.add_slide(slide_layout)
    add_tech_elements(slide)
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.8))
    title_text = title_box.text_frame
    title_para = title_text.add_paragraph()
    title_run = title_para.add_run()
    title_run.text = "项目6：优化维度"
    title_run.font.size = Pt(24)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(0, 200, 255)
    title_run.font.name = '微软雅黑'
    title_para.alignment = PP_ALIGN.CENTER
    
    table = slide.shapes.add_table(13, 9, left, top, width, height).table
    data = [
        ["优化维度", "技术方案", "预期收益", "实施难度", "成本估算", "ROI周期", "适用场景", "推荐度", "优先级"],
        ["网络优化", "HTTP/3+QUIC", "延迟-40%", "中", "中", "3个月", "大规模", "9", "P0"],
        ["缓存策略", "Redis+CDN", "延迟-60%", "低", "中", "2个月", "高并发", "10", "P0"],
        ["并发处理", "AsyncIO", "吞吐+200%", "中", "低", "1个月", "API服务", "9", "P1"],
        ["CDN加速", "Cloudflare", "延迟-70%", "低", "高", "1个月", "互联网", "10", "P0"],
        ["数据库优化", "索引+读写分离", "查询+300%", "中", "中", "3个月", "大数据", "9", "P1"],
        ["Agent优化", "LLM缓存+微调", "成本-60%", "高", "中", "3个月", "AI服务", "10", "P0"],
        ["AI推理优化", "量化+蒸馏", "速度+200%", "高", "中", "3个月", "AI服务", "10", "P0"],
        ["索引优化", "向量索引", "查询+1000%", "中", "中", "2个月", "搜索服务", "10", "P0"],
        ["熔断降级", "Sentinel", "可用性+99.9%", "低", "低", "1个月", "微服务", "10", "P0"],
        ["边缘计算", "EdgeAI", "延迟-80%", "中", "高", "3个月", "IoT", "9", "P1"],
        ["GPU加速", "CUDA+TensorRT", "AI+1000%", "高", "极高", "2个月", "AI训练", "10", "P0"],
        ["综合优化", "全链路优化", "整体+200%", "高", "高", "6个月", "企业级", "10", "P0"],
    ]
    
    for i in range(len(data)):
        for j in range(len(data[i])):
            if i == 0:
                set_cell_style(table.cell(i, j), data[i][j], 11, True, RGBColor(0, 200, 255), RGBColor(20, 30, 50))
            elif j == 0:
                set_cell_style(table.cell(i, j), data[i][j], 10, True, RGBColor(200, 200, 255), RGBColor(15, 25, 45))
            else:
                set_cell_style(table.cell(i, j), data[i][j], 9, False, RGBColor(220, 220, 240), RGBColor(10, 15, 30))
    
    # 项目7：案例场景
    slide = prs.slides.add_slide(slide_layout)
    add_tech_elements(slide)
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.8))
    title_text = title_box.text_frame
    title_para = title_text.add_paragraph()
    title_run = title_para.add_run()
    title_run.text = "项目7：案例场景"
    title_run.font.size = Pt(24)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(0, 200, 255)
    title_run.font.name = '微软雅黑'
    title_para.alignment = PP_ALIGN.CENTER
    
    table = slide.shapes.add_table(11, 9, left, top, width, height).table
    data = [
        ["案例场景", "行业领域", "技术方案", "核心成果", "投资规模", "ROI周期", "技术难度", "推荐度", "可复制性"],
        ["供应链管理", "制造业", "AI预测+RAG", "库存-30%,效率+50%", "500万", "12个月", "中", "9", "高"],
        ["金融风控", "金融", "GPT 5.0+知识图谱", "欺诈检测+40%,成本-60%", "800万", "8个月", "高", "10", "高"],
        ["智能办公", "企业", "DeepSeek V5.0+Agent", "效率+60%,成本-40%", "200万", "6个月", "低", "10", "极高"],
        ["医疗诊断", "医疗", "Gemini 2.5+多模态", "误诊率-40%,效率+30%", "1000万", "18个月", "极高", "9", "中"],
        ["智能客服", "零售", "GPT 5.0+千问3.6Plus", "满意度+50%,成本-70%", "300万", "4个月", "低", "10", "极高"],
        ["智能运维", "IT", "AI+可观测性", "故障预警+80%,MTTR-60%", "400万", "6个月", "中", "9", "高"],
        ["智能推荐", "电商", "DeepSeek V5.0", "转化+35%,GMV+25%", "500万", "6个月", "中", "9", "高"],
        ["智能教育", "教育", "GPT 5.0+个性化", "效果+40%,辍学率-20%", "300万", "8个月", "中", "9", "高"],
        ["智能制造", "制造", "AI质检+预测维护", "次品率-50%,停机-40%", "800万", "12个月", "中高", "9", "中"],
        ["Agent协作", "AI", "A2A协议+多智能体", "效率+100%,能力+200%", "500万", "6个月", "中", "10", "高"],
    ]
    
    for i in range(len(data)):
        for j in range(len(data[i])):
            if i == 0:
                set_cell_style(table.cell(i, j), data[i][j], 11, True, RGBColor(0, 200, 255), RGBColor(20, 30, 50))
            elif j == 0:
                set_cell_style(table.cell(i, j), data[i][j], 10, True, RGBColor(200, 200, 255), RGBColor(15, 25, 45))
            else:
                set_cell_style(table.cell(i, j), data[i][j], 9, False, RGBColor(220, 220, 240), RGBColor(10, 15, 30))
    
    # 项目8：技术方案对比
    slide = prs.slides.add_slide(slide_layout)
    add_tech_elements(slide)
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.8))
    title_text = title_box.text_frame
    title_para = title_text.add_paragraph()
    title_run = title_para.add_run()
    title_run.text = "项目8：技术方案对比"
    title_run.font.size = Pt(24)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(0, 200, 255)
    title_run.font.name = '微软雅黑'
    title_para.alignment = PP_ALIGN.CENTER
    
    table = slide.shapes.add_table(11, 9, left, top, width, height).table
    data = [
        ["技术方案", "实时性", "安全性", "扩展性", "易用性", "性能", "成本", "生态", "综合评分"],
        ["A2A协议", "<10ms", "零信任", "K8s原生", "极高", "极高", "中", "成长", "9.8"],
        ["传统API", ">100ms", "OAuth2", "中等", "极高", "低", "低", "成熟", "6.2"],
        ["消息队列", "<50ms", "TLS", "中等", "中等", "中", "中", "成熟", "7.5"],
        ["gRPC", "<30ms", "TLS", "高", "中等", "高", "中", "成熟", "7.8"],
        ["WebSocket", "<10ms", "TLS", "高", "高", "极高", "低", "成熟", "8.5"],
        ["GraphQL", ">100ms", "OAuth2", "高", "高", "低", "中", "成长", "7.2"],
        ["RESTful", ">100ms", "TLS", "高", "极高", "低", "低", "成熟", "6.8"],
        ["MQTT", "<50ms", "TLS", "高", "高", "高", "低", "成熟", "8.2"],
        ["NATS", "<10ms", "TLS", "高", "高", "极高", "低", "成长", "8.5"],
        ["对比结论", "A2A最优", "A2A最优", "A2A最优", "A2A最优", "A2A最优", "中等", "成长中", "A2A领先"],
    ]
    
    for i in range(len(data)):
        for j in range(len(data[i])):
            if i == 0:
                set_cell_style(table.cell(i, j), data[i][j], 11, True, RGBColor(0, 200, 255), RGBColor(20, 30, 50))
            elif j == 0:
                set_cell_style(table.cell(i, j), data[i][j], 10, True, RGBColor(200, 200, 255), RGBColor(15, 25, 45))
            elif "A2A最优" in data[i][j] or "A2A领先" in data[i][j]:
                set_cell_style(table.cell(i, j), data[i][j], 9, True, RGBColor(0, 255, 100), RGBColor(10, 30, 20))
            else:
                set_cell_style(table.cell(i, j), data[i][j], 9, False, RGBColor(220, 220, 240), RGBColor(10, 15, 30))
    
    # 项目9：市场预测
    slide = prs.slides.add_slide(slide_layout)
    add_tech_elements(slide)
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.8))
    title_text = title_box.text_frame
    title_para = title_text.add_paragraph()
    title_run = title_para.add_run()
    title_run.text = "项目9：市场预测（2026-2035年）"
    title_run.font.size = Pt(24)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(0, 200, 255)
    title_run.font.name = '微软雅黑'
    title_para.alignment = PP_ALIGN.CENTER
    
    table = slide.shapes.add_table(11, 8, left, top, width, height).table
    data = [
        ["年份", "市场规模", "增长率", "主要技术", "投资规模", "应用场景", "竞争格局", "技术突破"],
        ["2026", "2000亿", "85%", "GPT 5.0", "600亿", "企业服务", "格局初定", "AGI雏形"],
        ["2027", "3800亿", "90%", "DeepSeek V5.0", "1200亿", "具身智能", "巨头主导", "人形机器人"],
        ["2028", "6500亿", "71%", "GPT 6.0", "2000亿", "服务机器人", "生态竞争", "通用机器人"],
        ["2029", "10000亿", "54%", "DeepSeek V6.0", "3000亿", "家庭助手", "寡头垄断", "脑机接口"],
        ["2030", "15000亿", "50%", "GPT 7.0", "4500亿", "人机融合", "全球格局", "通用AGI"],
        ["2031", "22000亿", "47%", "量子AI", "6000亿", "智能经济", "生态竞争", "量子融合"],
        ["2032", "32000亿", "45%", "GPT 8.0", "8000亿", "万物互联", "平台竞争", "边缘大脑"],
        ["2033", "45000亿", "41%", "DeepSeek V7.0", "10000亿", "人机共生", "寡头格局", "神经接口"],
        ["2034", "62000亿", "38%", "量子AGI", "12000亿", "智能文明", "全球垄断", "量子AGI"],
        ["2035", "85000亿", "37%", "通用AGI", "15000亿", "全面智能", "全球一体", "超级智能"],
    ]
    
    for i in range(len(data)):
        for j in range(len(data[i])):
            if i == 0:
                set_cell_style(table.cell(i, j), data[i][j], 11, True, RGBColor(0, 200, 255), RGBColor(20, 30, 50))
            elif j == 0:
                set_cell_style(table.cell(i, j), data[i][j], 10, True, RGBColor(200, 200, 255), RGBColor(15, 25, 45))
            else:
                set_cell_style(table.cell(i, j), data[i][j], 9, False, RGBColor(220, 220, 240), RGBColor(10, 15, 30))
    
    # 项目10：总结与规划
    slide = prs.slides.add_slide(slide_layout)
    add_tech_elements(slide)
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.8))
    title_text = title_box.text_frame
    title_para = title_text.add_paragraph()
    title_run = title_para.add_run()
    title_run.text = "项目10：总结与未来规划"
    title_run.font.size = Pt(24)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(0, 200, 255)
    title_run.font.name = '微软雅黑'
    title_para.alignment = PP_ALIGN.CENTER
    
    table = slide.shapes.add_table(12, 4, left, top, width, height).table
    data = [
        ["项目", "核心内容", "当前状态", "未来规划"],
        ["项目1", "封面设计", "完成", "季度更新"],
        ["项目2", "价值维度", "完成", "持续优化"],
        ["项目3", "协议层级", "完成", "A2A 3.0升级"],
        ["项目4", "安全组件", "完成", "零信任深化"],
        ["项目5", "部署策略", "完成", "自动化升级"],
        ["项目6", "优化维度", "完成", "AI驱动优化"],
        ["项目7", "案例场景", "完成", "扩展至20领域"],
        ["项目8", "技术对比", "完成", "季度更新"],
        ["项目9", "市场预测", "完成", "滚动预测"],
        ["项目10", "总结展望", "完成", "年度更新"],
        ["合计", "完整AI方案", "全部完成", "持续演进"],
    ]
    
    for i in range(len(data)):
        for j in range(len(data[i])):
            if i == 0:
                set_cell_style(table.cell(i, j), data[i][j], 12, True, RGBColor(0, 200, 255), RGBColor(20, 30, 50))
            elif j == 0:
                set_cell_style(table.cell(i, j), data[i][j], 11, True, RGBColor(200, 200, 255), RGBColor(15, 25, 45))
            elif "完成" in data[i][j]:
                set_cell_style(table.cell(i, j), data[i][j], 10, True, RGBColor(0, 255, 100), RGBColor(10, 30, 20))
            else:
                set_cell_style(table.cell(i, j), data[i][j], 10, False, RGBColor(220, 220, 240), RGBColor(10, 15, 30))
    
    # 保存
    output_path = r"f:\个人作品\legal-rag-qa-system\A2A_PROTOCOL_AI_AGENT_2026_科技感完整版.pptx"
    prs.save(output_path)
    print(f"✅ 已创建科技感PPT！保存到: {output_path}")

if __name__ == "__main__":
    create_tech_pptx()