from pptx import Presentation
import os

def check_all_slides():
    ppt_path = r"f:\个人作品\legal-rag-qa-system\A2A_PROTOCOL_AI_AGENT_V14_终极完整版.pptx"
    
    print("=" * 120)
    print("🔍 终极完整版内容详细检查")
    print("=" * 120)
    
    prs = Presentation(ppt_path)
    
    for slide_idx, slide in enumerate(prs.slides):
        print(f"\n📄 幻灯片 {slide_idx + 1}:")
        
        # 查找表格
        table_found = False
        for shape_idx, shape in enumerate(slide.shapes):
            if shape.has_table:
                table_found = True
                table = shape.table
                print(f"   └─ 表格: {len(table.rows)}行 × {len(table.columns)}列")
                
                # 打印第一列内容（层级/标题）
                first_col = []
                for i in range(len(table.rows)):
                    cell_text = table.cell(i, 0).text.strip()
                    first_col.append(cell_text)
                
                print(f"   └─ 第一列内容: {', '.join(first_col)}")
                
                # 如果是项目3，检查L10
                if slide_idx == 2:  # 项目3
                    has_l10 = any("L10" in c for c in first_col)
                    print(f"   └─ 是否包含L10: {'✅ 是' if has_l10 else '❌ 否'}")
        
        if not table_found:
            print("   └─ 无表格")

def check_file_properties():
    ppt_path = r"f:\个人作品\legal-rag-qa-system\A2A_PROTOCOL_AI_AGENT_V14_终极完整版.pptx"
    
    print("\n" + "=" * 120)
    print("📁 文件属性")
    print("=" * 120)
    
    if os.path.exists(ppt_path):
        stat = os.stat(ppt_path)
        print(f"文件路径: {ppt_path}")
        print(f"文件大小: {stat.st_size} 字节")
        print(f"创建时间: {stat.st_ctime}")
        print(f"修改时间: {stat.st_mtime}")
        
        # 尝试读取内部结构
        prs = Presentation(ppt_path)
        print(f"\n内部结构:")
        print(f"  幻灯片数: {len(prs.slides)}")
        
        # 检查项目3表格
        if len(prs.slides) > 2:
            slide = prs.slides[2]
            for shape in slide.shapes:
                if shape.has_table:
                    table = shape.table
                    print(f"  项目3表格: {len(table.rows)}行 × {len(table.columns)}列")
                    
                    # 打印所有层级
                    print("  层级列表:")
                    for i in range(1, len(table.rows)):
                        layer = table.cell(i, 0).text.strip()
                        name = table.cell(i, 1).text.strip()
                        print(f"    {layer}: {name}")
                    break

if __name__ == "__main__":
    check_all_slides()
    check_file_properties()