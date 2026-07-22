from pptx import Presentation

ppt_path = r"f:\个人作品\legal-rag-qa-system\A2A_PROTOCOL_AI_AGENT_V14_完美清晰版.pptx"
prs = Presentation(ppt_path)

print("=" * 60)
print("PPT表格结构诊断")
print("=" * 60)

for slide_idx, slide in enumerate(prs.slides):
    print(f"\n📊 第 {slide_idx + 1} 张幻灯片:")
    for shape_idx, shape in enumerate(slide.shapes):
        if shape.has_table:
            table = shape.table
            rows = len(table.rows)
            cols = len(table.columns)
            print(f"  表格 {shape_idx + 1}: {rows} 行 x {cols} 列")
            print(f"  表头内容: ", end="")
            header = []
            for col in range(min(cols, 5)):
                cell = table.cell(0, col)
                text = cell.text.strip()[:10]
                header.append(text)
            print(header)
            print(f"  第二行内容: ", end="")
            if rows > 1:
                row2 = []
                for col in range(min(cols, 5)):
                    cell = table.cell(1, col)
                    text = cell.text.strip()[:10]
                    row2.append(f"'{text}'" if text else "'(空)'")
                print(row2)