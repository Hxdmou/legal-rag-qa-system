from pptx import Presentation

def check_current_state():
    print("=" * 120)
    print("🔍 紧急检查：当前PPT实际内容")
    print("=" * 120)
    
    # 检查最新创建的文件
    ppt_path = r"f:\个人作品\legal-rag-qa-system\AI_AGENT_PROTOCOL_FINAL_2026_最新完整版.pptx"
    
    prs = Presentation(ppt_path)
    print(f"幻灯片总数: {len(prs.slides)}")
    
    # 检查幻灯片3（项目3）
    if len(prs.slides) > 2:
        slide = prs.slides[2]
        print(f"\n📄 幻灯片3 (项目3):")
        
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                print(f"   表格: {len(table.rows)}行 × {len(table.columns)}列")
                
                # 打印列标题
                headers = []
                for j in range(min(10, len(table.columns))):
                    headers.append(table.cell(0, j).text.strip())
                print(f"   列标题: {', '.join(headers)}")
                
                # 打印所有行的第一列
                rows = []
                for i in range(len(table.rows)):
                    rows.append(table.cell(i, 0).text.strip())
                print(f"   行内容: {', '.join(rows)}")
                
                break

def find_all_ppts():
    import os
    print("\n" + "=" * 120)
    print("📁 查找所有PPT文件")
    print("=" * 120)
    
    for root, dirs, files in os.walk(r"f:\个人作品\legal-rag-qa-system"):
        for f in files:
            if f.endswith(".pptx"):
                full_path = os.path.join(root, f)
                size = os.path.getsize(full_path)
                print(f"✅ {f} ({size}字节)")

if __name__ == "__main__":
    find_all_ppts()
    check_current_state()