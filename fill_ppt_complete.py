from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

def set_cell_content_white(cell, text, font_size=10, bold=False):
    cell.text_frame.clear()
    para = cell.text_frame.add_paragraph()
    run = para.add_run()
    run.text = str(text) if text else ""
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(255, 255, 255)
    run.font.name = '微软雅黑'
    para.alignment = PP_ALIGN.CENTER
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

def fill_project2(slide):
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            rows, cols = len(table.rows), len(table.columns)
            data = [
                ["章节", "标题", "页码", "核心要点", "技术深度", "实用性", "创新性", "完整性", "可读性", "视觉设计", "交互体验", "数据支撑", "案例丰富度", "更新频率", "参考价值", "学习曲线"],
                ["1", "协议概述", "03", "A2A核心概念", "深入", "强", "强", "完整", "优", "优", "优", "丰富", "中", "每月", "高", "简单"],
                ["2", "技术架构", "06", "分层设计", "深入", "强", "强", "完整", "优", "优", "优", "丰富", "多", "每季度", "高", "中等"],
                ["3", "协议核心", "09", "通信机制", "深入", "强", "强", "完整", "优", "优", "优", "丰富", "多", "每月", "高", "困难"],
                ["4", "AI智能体", "12", "能力矩阵", "深入", "强", "强", "完整", "优", "优", "优", "丰富", "多", "每周", "高", "困难"],
                ["5", "最新趋势", "15", "2026热点", "深入", "强", "强", "完整", "优", "优", "优", "丰富", "多", "每周", "高", "中等"],
            ]
            for i in range(rows):
                for j in range(cols):
                    text = ""
                    if i < len(data) and j < len(data[i]):
                        text = data[i][j]
                    set_cell_content_white(table.cell(i, j), text, 12 if i == 0 else 9, i == 0)
            break

def fill_project3(slide):
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            rows, cols = len(table.rows), len(table.columns)
            data = [
                ["层级", "名称", "核心功能", "技术实现", "通信协议", "数据格式", "加密方式", "性能指标", "可靠性", "扩展性", "安全性"],
                ["L1", "传输层", "消息路由", "HTTP/2加gRPC", "HTTP/2", "Protobuf", "TLS1.3", "小于50毫秒", "99.999%", "线性扩展", "AES-256加密"],
                ["L2", "协议层", "通信协议", "A2A-Spec v2.0", "WebSocket", "JSON", "JWT", "小于20毫秒", "99.99%", "插件化", "OAuth2认证"],
                ["L3", "智能层", "AI决策", "LLM加RAG", "SSE", "JSON-RPC", "AES-256", "小于100毫秒", "99.9%", "模块化", "RBAC权限"],
                ["L4", "应用层", "业务逻辑", "多智能体", "gRPC", "Protobuf", "mTLS", "小于30毫秒", "99.99%", "微服务", "ABAC策略"],
                ["L5", "表现层", "用户交互", "Web前端", "HTTP/3", "JSON", "TLS1.3", "小于10毫秒", "99.99%", "CDN加速", "WAF防护"],
                ["L6", "数据层", "存储计算", "分布式存储", "TCP", "二进制", "SM4加密", "小于5毫秒", "99.999%", "数据分片", "数据脱敏"],
                ["L7", "安全层", "身份认证", "零信任架构", "mTLS", "X509", "国密算法", "小于15毫秒", "99.99%", "动态策略", "审计日志"],
                ["L8", "监控层", "运维监控", "可观测性", "HTTP", "OpenTelemetry", "TLS加密", "实时监控", "99.9%", "自动告警", "链路追踪"],
                ["L9", "网关层", "流量管理", "API网关", "HTTP/2", "JSON", "JWT加RSA", "小于25毫秒", "99.99%", "负载均衡", "限流熔断"],
            ]
            for i in range(rows):
                for j in range(cols):
                    text = ""
                    if i < len(data) and j < len(data[i]):
                        text = data[i][j]
                    set_cell_content_white(table.cell(i, j), text, 12 if i == 0 else 9, i == 0)
            break

def fill_project4(slide):
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            rows, cols = len(table.rows), len(table.columns)
            data = [
                ["层级", "能力描述", "技术支撑", "成熟度", "训练数据", "推理成本", "响应时间", "准确率", "可解释性", "可控性", "可扩展性", "安全性", "商业价值", "竞争优势", "未来潜力", "学习成本"],
                ["L1", "基础执行", "脚本引擎", "95%", "小规模", "低成本", "小于10毫秒", "99%", "高", "高", "高", "高", "中等", "低", "低", "低"],
                ["L2", "工具调用", "Function Calling", "90%", "中规模", "中成本", "小于50毫秒", "95%", "中", "中", "中", "中", "中高", "中", "中", "中"],
                ["L3", "规划推理", "Chain of Thought", "85%", "大规模", "中高成本", "小于200毫秒", "90%", "中", "中", "中", "中", "高", "高", "高", "高"],
                ["L4", "多模态理解", "Vision加Text", "80%", "大规模", "高成本", "小于500毫秒", "85%", "低", "中", "中", "中", "高", "高", "高", "高"],
                ["L5", "自我反思", "Self-Correction", "75%", "大规模", "高成本", "小于1秒", "80%", "低", "低", "低", "中", "极高", "极高", "极高", "极高"],
            ]
            for i in range(rows):
                for j in range(cols):
                    text = ""
                    if i < len(data) and j < len(data[i]):
                        text = data[i][j]
                    set_cell_content_white(table.cell(i, j), text, 12 if i == 0 else 9, i == 0)
            break

def fill_project5(slide):
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            rows, cols = len(table.rows), len(table.columns)
            data = [
                ["趋势", "成熟度", "市场影响", "代表技术", "商业价值", "技术难度", "投资热度", "竞争格局", "预期收益", "风险等级", "时间窗口", "人才需求", "基础设施", "政策影响"],
                ["LLM Agent", "高", "革命性", "自主智能体多任务规划", "极高", "高", "极热", "激烈", "10倍", "中", "1至3年", "紧缺", "高", "积极"],
                ["具身智能", "中", "深远", "机器人感知运动控制", "极高", "极高", "很热", "中等", "20倍", "高", "3至5年", "极缺", "高", "积极"],
                ["AI安全对齐", "中高", "关键", "RLHF加Constitutional AI", "极高", "高", "很热", "中等", "5倍", "低", "1至2年", "紧缺", "中", "积极"],
                ["多模态融合", "高", "广泛", "视觉理解语音识别", "高", "中高", "很热", "激烈", "3倍", "中", "1至2年", "紧缺", "高", "积极"],
                ["边缘AI", "中", "普及", "端侧推理隐私计算", "中高", "中", "较热", "中等", "2倍", "低", "2至3年", "中", "中", "积极"],
                ["AI编程", "中高", "效率革命", "CodeLlama加StarCoder", "高", "中", "很热", "激烈", "5倍", "低", "1至2年", "紧缺", "中", "积极"],
                ["RAG 2.0", "高", "落地加速", "多向量检索混合搜索", "高", "中", "较热", "激烈", "2倍", "低", "1年", "中", "高", "积极"],
            ]
            for i in range(rows):
                for j in range(cols):
                    text = ""
                    if i < len(data) and j < len(data[i]):
                        text = data[i][j]
                    set_cell_content_white(table.cell(i, j), text, 12 if i == 0 else 9, i == 0)
            break

def fill_project6(slide):
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            rows, cols = len(table.rows), len(table.columns)
            data = [
                ["领域", "典型场景", "智能体角色", "预期收益", "技术难度", "实施周期", "投资规模", "ROI周期", "风险等级", "市场成熟度", "竞争格局", "政策支持", "人才需求"],
                ["企业办公", "智能助理", "任务自动化", "效率提升50%", "中", "3个月", "100万", "18个月", "低", "高", "激烈", "强", "中"],
                ["金融服务", "风控合规", "风险评估", "准确率提升30%", "高", "6个月", "500万", "12个月", "中", "中高", "中等", "强", "高"],
                ["医疗健康", "辅助诊断", "医学分析", "误诊率降低40%", "极高", "12个月", "1000万", "24个月", "高", "中", "中等", "强", "极高"],
                ["智能制造", "质检维护", "设备监控", "故障预警提升60%", "中高", "9个月", "800万", "15个月", "中", "中高", "中等", "强", "高"],
                ["教育培训", "个性化学习", "智能推荐", "学习效果提升40%", "中", "6个月", "200万", "12个月", "低", "高", "激烈", "强", "中"],
                ["零售电商", "智能营销", "精准推荐", "转化率提升35%", "中", "4个月", "150万", "10个月", "低", "高", "激烈", "中", "中"],
                ["物流运输", "路径优化", "智能调度", "配送效率提升45%", "中高", "8个月", "300万", "14个月", "中", "中高", "中等", "中", "高"],
                ["能源电力", "智能调度", "负荷预测", "能耗降低25%", "高", "10个月", "600万", "18个月", "中", "中", "中等", "强", "高"],
            ]
            for i in range(rows):
                for j in range(cols):
                    text = ""
                    if i < len(data) and j < len(data[i]):
                        text = data[i][j]
                    set_cell_content_white(table.cell(i, j), text, 12 if i == 0 else 9, i == 0)
            break

def fill_project7(slide):
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            rows, cols = len(table.rows), len(table.columns)
            data = [
                ["案例", "行业", "应用场景", "核心成果", "技术方案", "投资规模", "ROI周期", "实施周期", "风险控制", "可复制性", "创新性", "商业价值"],
                ["某大型银行", "金融", "智能风控", "欺诈识别提升40%", "LLM加RAG", "200万", "2个月", "6个月", "高", "高", "高", "极高"],
                ["某电商平台", "零售", "智能客服", "成本降低60%", "多智能体", "150万", "3个月", "3个月", "中", "高", "中", "高"],
                ["某医疗机构", "医疗", "辅助诊断", "准确率提升25%", "Vision加LLM", "300万", "2.4个月", "12个月", "高", "中", "高", "极高"],
                ["某制造企业", "制造", "设备维护", "停机时间降低50%", "IoT加AI", "800万", "15个月", "9个月", "中", "高", "中", "高"],
                ["某保险公司", "金融", "智能理赔", "理赔效率提升70%", "LLM加知识图谱", "250万", "3个月", "5个月", "高", "高", "高", "高"],
            ]
            for i in range(rows):
                for j in range(cols):
                    text = ""
                    if i < len(data) and j < len(data[i]):
                        text = data[i][j]
                    set_cell_content_white(table.cell(i, j), text, 12 if i == 0 else 9, i == 0)
            break

def fill_project8(slide):
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            rows, cols = len(table.rows), len(table.columns)
            data = [
                ["特性", "A2A协议", "传统API", "消息队列", "gRPC", "微服务", "WebSocket", "GraphQL", "RESTful", "SOAP", "Webhook", "SSE"],
                ["实时性", "极强", "弱", "强", "强", "弱", "极强", "弱", "弱", "弱", "中等", "极强"],
                ["双向通信", "极强", "弱", "弱", "强", "弱", "极强", "弱", "弱", "弱", "弱", "极强"],
                ["智能路由", "极强", "无", "中等", "中等", "弱", "中等", "无", "无", "无", "中等", "弱"],
                ["状态管理", "极强", "弱", "中等", "中等", "中等", "强", "弱", "中等", "弱", "中等", "弱"],
                ["LLM集成", "极强", "无", "无", "中等", "中等", "中等", "无", "无", "无", "中等", "弱"],
            ]
            for i in range(rows):
                for j in range(cols):
                    text = ""
                    if i < len(data) and j < len(data[i]):
                        text = data[i][j]
                    set_cell_content_white(table.cell(i, j), text, 12 if i == 0 else 9, i == 0)
            break

def main():
    ppt_path = r"f:\个人作品\legal-rag-qa-system\A2A_PROTOCOL_AI_AGENT_V14_完美清晰版.pptx"

    if not os.path.exists(ppt_path):
        print("未找到PPT文件！")
        return

    print(f"📁 使用PPT文件: {ppt_path}")

    prs = Presentation(ppt_path)

    project_fillers = [
        None,
        fill_project2,
        fill_project3,
        fill_project4,
        fill_project5,
        fill_project6,
        fill_project7,
        fill_project8,
    ]

    for i, filler in enumerate(project_fillers):
        if filler and i < len(prs.slides):
            print(f"🔄 正在处理第 {i} 张幻灯片 (项目{i})...")
            filler(prs.slides[i])

    output_path = r"f:\个人作品\legal-rag-qa-system\A2A_PROTOCOL_AI_AGENT_V14_完整填充版.pptx"
    prs.save(output_path)
    print(f"\n✅ PPT表格完整填充完成！已保存到: {output_path}")

if __name__ == "__main__":
    main()