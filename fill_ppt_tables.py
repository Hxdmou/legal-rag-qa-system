from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

def set_cell_content_white(cell, text, font_size=10, bold=False):
    """设置单元格内容，强制使用白色字体确保在深色背景上清晰可见"""
    # 清除现有内容
    cell.text_frame.clear()
    
    # 添加新段落
    para = cell.text_frame.add_paragraph()
    run = para.add_run()
    run.text = str(text)
    
    # 设置格式 - 纯白色字体在深色背景上最清晰
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(255, 255, 255)  # 纯白色
    run.font.name = '微软雅黑'
    run.font.sans_serif = True
    
    # 对齐
    para.alignment = PP_ALIGN.CENTER
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

def fill_project1(slide):
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            data = [
                ["元素", "规格", "技术实现", "渲染引擎", "分辨率", "色彩空间", "动画效果", "交互方式", "文件格式", "压缩率", "加载时间", "兼容性", "性能优化", "缓存策略", "CDN加速", "响应式"],
                ["主标题", "48号字", "渐变渲染", "SVG", "1920×1080", "sRGB", "淡入", "静态", "SVG", "85%", "<100ms", "全平台", "GPU加速", "LRU", "支持", "是"],
                ["副标题", "22号字", "霓虹光效", "Canvas", "1920×1080", "P3", "闪烁", "静态", "PNG", "90%", "<50ms", "全平台", "WebGL", "LRU", "支持", "是"],
                ["背景", "全屏", "AI纹理", "SDXL", "3840×2160", "Rec.2020", "缩放", "静态", "WebP", "75%", "<200ms", "全平台", "WebGL", "LRU", "支持", "是"],
                ["视觉元素", "多层", "SVG动画", "GSAP", "1920×1080", "sRGB", "旋转", "悬停", "SVG", "80%", "<150ms", "全平台", "GPU", "LRU", "支持", "是"],
            ]
            
            for i, row in enumerate(data):
                for j, cell_value in enumerate(row):
                    if i < len(table.rows) and j < len(table.columns):
                        font_size = 12 if i == 0 else 10
                        bold = (i == 0)
                        set_cell_content_white(table.cell(i, j), cell_value, font_size, bold)
            break

def fill_project2(slide):
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            data = [
                ["章节", "标题", "页码", "核心要点", "技术深度", "实用性", "创新性", "完整性", "可读性", "视觉设计", "交互体验", "数据支撑", "案例丰富度", "更新频率", "参考价值", "学习曲线", "推荐指数"],
                ["1", "协议概述", "03", "A2A核心概念", "深入", "强", "强", "完整", "优", "优", "优", "丰富", "中", "每月", "高", "简单", "9.5"],
                ["2", "技术架构", "06", "分层设计", "深入", "强", "强", "完整", "优", "优", "优", "丰富", "多", "每季度", "高", "中等", "9.8"],
                ["3", "协议核心", "09", "通信机制", "深入", "强", "强", "完整", "优", "优", "优", "丰富", "多", "每月", "高", "困难", "9.7"],
                ["4", "AI智能体", "12", "能力矩阵", "深入", "强", "强", "完整", "优", "优", "优", "丰富", "多", "每周", "高", "困难", "9.9"],
                ["5", "最新趋势", "15", "2026热点", "深入", "强", "强", "完整", "优", "优", "优", "丰富", "多", "每周", "高", "中等", "10.0"],
                ["6", "应用场景", "18", "四大领域", "中等", "强", "中", "完整", "优", "优", "优", "丰富", "中", "每月", "高", "简单", "9.6"],
                ["7", "案例分析", "21", "真实验证", "深入", "强", "强", "完整", "优", "优", "优", "丰富", "多", "每季度", "高", "中等", "9.8"],
                ["8", "技术对比", "24", "竞争优势", "深入", "强", "强", "完整", "优", "优", "优", "丰富", "多", "每月", "高", "困难", "9.9"],
                ["9", "市场展望", "27", "5年规划", "深入", "强", "强", "完整", "优", "优", "优", "丰富", "多", "每季度", "高", "中等", "9.7"],
                ["10", "结语", "30", "行动建议", "中等", "强", "中", "完整", "优", "优", "优", "丰富", "少", "每年", "高", "简单", "9.5"],
            ]
            
            for i, row in enumerate(data):
                for j, cell_value in enumerate(row):
                    if i < len(table.rows) and j < len(table.columns):
                        font_size = 12 if i == 0 else 10
                        bold = (i == 0)
                        set_cell_content_white(table.cell(i, j), cell_value, font_size, bold)
            break

def fill_project3(slide):
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            data = [
                ["层级", "名称", "核心功能", "技术实现", "通信协议", "数据格式", "加密方式", "性能指标", "可靠性", "扩展性", "安全性", "监控能力", "调试工具", "文档完整度", "社区活跃度", "学习成本", "推荐指数"],
                ["L1", "传输层", "消息路由", "HTTP/2 + gRPC", "HTTP/2", "Protobuf", "TLS1.3", "<50ms", "99.999%", "线性", "AES-256", "Prometheus", "grpcurl", "95%", "高", "中", "10"],
                ["L2", "协议层", "通信协议", "A2A-Spec v2.0", "WebSocket", "JSON", "JWT", "<20ms", "99.99%", "插件", "OAuth2", "Grafana", "Postman", "98%", "高", "高", "10"],
                ["L3", "智能层", "AI决策", "LLM + RAG", "SSE", "JSON-RPC", "AES-256", "<100ms", "99.9%", "模块化", "RBAC", "ELK", "LangSmith", "90%", "中", "高", "9.5"],
                ["L4", "应用层", "业务逻辑", "多智能体", "gRPC", "Protobuf", "mTLS", "<30ms", "99.99%", "微服务", "ABAC", "Jaeger", "OpenTelemetry", "85%", "中", "中", "9"],
            ]
            
            for i, row in enumerate(data):
                for j, cell_value in enumerate(row):
                    if i < len(table.rows) and j < len(table.columns):
                        font_size = 12 if i == 0 else 10
                        bold = (i == 0)
                        set_cell_content_white(table.cell(i, j), cell_value, font_size, bold)
            break

def fill_project4(slide):
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            data = [
                ["层级", "能力描述", "技术支撑", "成熟度", "训练数据", "推理成本", "响应时间", "准确率", "可解释性", "可控性", "可扩展性", "安全性", "商业价值", "竞争优势", "未来潜力", "学习成本", "推荐度"],
                ["L1", "基础执行", "脚本引擎", "95%", "小", "低", "<10ms", "99%", "高", "高", "高", "高", "中", "低", "低", "低", "8"],
                ["L2", "工具调用", "Function Calling", "90%", "中", "中", "<50ms", "95%", "中", "中", "中", "中", "中高", "中", "中", "中", "9"],
                ["L3", "规划推理", "Chain of Thought", "85%", "大", "中高", "<200ms", "90%", "中", "中", "中", "中", "高", "高", "高", "高", "9.5"],
                ["L4", "多模态理解", "Vision+Text", "80%", "大", "高", "<500ms", "85%", "低", "中", "中", "中", "高", "高", "高", "高", "9"],
                ["L5", "自我反思", "Self-Correction", "75%", "大", "高", "<1s", "80%", "低", "低", "低", "中", "高", "极高", "极高", "极高", "9.8"],
                ["L6", "多智能体协作", "A2A协议", "70%", "大", "高", "<2s", "75%", "低", "低", "低", "中", "极高", "极高", "极高", "极高", "10"],
            ]
            
            for i, row in enumerate(data):
                for j, cell_value in enumerate(row):
                    if i < len(table.rows) and j < len(table.columns):
                        font_size = 12 if i == 0 else 10
                        bold = (i == 0)
                        set_cell_content_white(table.cell(i, j), cell_value, font_size, bold)
            break

def fill_project5(slide):
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            data = [
                ["趋势", "成熟度", "市场影响", "代表技术", "商业价值", "技术难度", "投资热度", "竞争格局", "预期收益", "风险等级", "时间窗口", "人才需求", "基础设施", "政策影响", "推荐度"],
                ["LLM Agent", "高", "革命性", "自主智能体+多任务规划", "极高", "高", "极热", "激烈", "10倍", "中", "1-3年", "紧缺", "高", "积极", "10"],
                ["具身智能", "中", "深远", "机器人感知+运动控制", "极高", "极高", "很热", "中等", "20倍", "高", "3-5年", "极缺", "高", "积极", "9.5"],
                ["AI安全对齐", "中高", "关键", "RLHF+Constitutional AI", "极高", "高", "很热", "中等", "5倍", "低", "1-2年", "紧缺", "中", "积极", "10"],
                ["多模态融合", "高", "广泛", "视觉理解+语音识别", "高", "中高", "很热", "激烈", "3倍", "中", "1-2年", "紧缺", "高", "积极", "9.5"],
                ["边缘AI", "中", "普及", "端侧推理+隐私计算", "中高", "中", "较热", "中等", "2倍", "低", "2-3年", "中", "中", "积极", "9"],
                ["AI编程", "中高", "效率革命", "CodeLlama+StarCoder", "高", "中", "很热", "激烈", "5倍", "低", "1-2年", "紧缺", "中", "积极", "9.5"],
                ["RAG 2.0", "高", "落地加速", "多向量检索+混合搜索", "高", "中", "较热", "激烈", "2倍", "低", "1年", "中", "高", "积极", "9"],
                ["长上下文", "高", "范式转变", "百万级Token处理", "极高", "中高", "很热", "激烈", "5倍", "中", "1-2年", "紧缺", "高", "积极", "10"],
            ]
            
            for i, row in enumerate(data):
                for j, cell_value in enumerate(row):
                    if i < len(table.rows) and j < len(table.columns):
                        font_size = 12 if i == 0 else 10
                        bold = (i == 0)
                        set_cell_content_white(table.cell(i, j), cell_value, font_size, bold)
            break

def fill_project6(slide):
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            data = [
                ["领域", "典型场景", "智能体角色", "预期收益", "技术难度", "实施周期", "投资规模", "ROI", "风险等级", "市场成熟度", "竞争格局", "政策支持", "人才需求", "推荐度"],
                ["企业办公", "智能助理", "任务自动化", "效率+50%", "中", "3个月", "100万", "18个月", "低", "高", "激烈", "强", "中", "9.5"],
                ["金融服务", "风控合规", "风险评估", "准确率+30%", "高", "6个月", "500万", "12个月", "中", "中高", "中等", "强", "高", "9"],
                ["医疗健康", "辅助诊断", "医学分析", "误诊率-40%", "极高", "12个月", "1000万", "24个月", "高", "中", "中等", "强", "极高", "8.5"],
                ["智能制造", "质检维护", "设备监控", "故障预警+60%", "中高", "9个月", "800万", "15个月", "中", "中高", "中等", "强", "高", "9"],
            ]
            
            for i, row in enumerate(data):
                for j, cell_value in enumerate(row):
                    if i < len(table.rows) and j < len(table.columns):
                        font_size = 12 if i == 0 else 10
                        bold = (i == 0)
                        set_cell_content_white(table.cell(i, j), cell_value, font_size, bold)
            break

def fill_project7(slide):
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            data = [
                ["案例", "行业", "应用场景", "核心成果", "技术方案", "投资规模", "ROI", "实施周期", "风险控制", "可复制性", "创新性", "商业价值", "推荐度"],
                ["某大型银行", "金融", "智能风控", "欺诈识别+40%", "LLM+RAG", "200万", "2个月", "6个月", "高", "高", "高", "极高", "10"],
                ["某电商平台", "零售", "智能客服", "成本-60%", "多智能体", "150万", "3个月", "3个月", "中", "高", "中", "高", "9.5"],
                ["某医疗机构", "医疗", "辅助诊断", "准确率+25%", "Vision+LLM", "300万", "2.4个月", "12个月", "高", "中", "高", "极高", "9"],
                ["某制造企业", "制造", "设备维护", "停机-50%", "IoT+AI", "800万", "15个月", "9个月", "中", "高", "中", "高", "9"],
            ]
            
            for i, row in enumerate(data):
                for j, cell_value in enumerate(row):
                    if i < len(table.rows) and j < len(table.columns):
                        font_size = 12 if i == 0 else 10
                        bold = (i == 0)
                        set_cell_content_white(table.cell(i, j), cell_value, font_size, bold)
            break

def fill_project8(slide):
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            data = [
                ["特性", "A2A协议", "传统API", "消息队列", "gRPC", "微服务", "WebSocket", "GraphQL", "RESTful", "SOAP", "Webhook", "SSE", "MQTT", "AMQP", "NATS", "推荐度"],
                ["实时性", "极强", "弱", "强", "强", "弱", "极强", "弱", "弱", "弱", "中等", "极强", "强", "强", "极强", "10"],
                ["双向通信", "极强", "弱", "弱", "强", "弱", "极强", "弱", "弱", "弱", "弱", "极强", "极强", "弱", "极强", "10"],
                ["智能路由", "极强", "无", "中等", "中等", "弱", "中等", "无", "无", "无", "中等", "弱", "弱", "中等", "弱", "10"],
                ["状态管理", "极强", "弱", "中等", "中等", "中等", "强", "弱", "中等", "弱", "中等", "弱", "弱", "中等", "弱", "10"],
                ["LLM集成", "极强", "无", "无", "中等", "中等", "中等", "无", "无", "无", "中等", "弱", "弱", "中等", "弱", "10"],
                ["安全性", "极强", "中等", "强", "强", "中等", "强", "中等", "中等", "中等", "中等", "强", "强", "强", "强", "10"],
                ["扩展性", "极强", "强", "中等", "强", "中等", "强", "强", "强", "中等", "中等", "强", "强", "强", "强", "10"],
                ["易用性", "极强", "强", "中等", "中等", "中等", "强", "强", "极强", "弱", "中等", "强", "强", "中等", "强", "10"],
            ]
            
            for i, row in enumerate(data):
                for j, cell_value in enumerate(row):
                    if i < len(table.rows) and j < len(table.columns):
                        font_size = 12 if i == 0 else 10
                        bold = (i == 0)
                        set_cell_content_white(table.cell(i, j), cell_value, font_size, bold)
            break

def main():
    # 查找可用的PPT文件
    ppt_candidates = [
        r"f:\个人作品\legal-rag-qa-system\A2A_PROTOCOL_AI_AGENT_V14.pptx",
        r"f:\个人作品\legal-rag-qa-system\A2A_PROTOCOL_AI_AGENT_V14_完美配色.pptx",
        r"f:\个人作品\legal-rag-qa-system\A2A_PROTOCOL_AI_AGENT_V14_蓝色字体.pptx",
        r"f:\个人作品\legal-rag-qa-system\A2A_PROTOCOL_AI_AGENT_V14_智能配色.pptx",
        r"f:\个人作品\legal-rag-qa-system\A2A_PROTOCOL_AI_AGENT_V14_完美清晰版.pptx",
    ]
    
    ppt_path = None
    for candidate in ppt_candidates:
        if os.path.exists(candidate):
            ppt_path = candidate
            break
    
    if ppt_path is None:
        print("未找到可用的PPT文件！")
        return
    
    print(f"📁 使用PPT文件: {ppt_path}")
    
    prs = Presentation(ppt_path)
    
    project_fillers = [
        fill_project1,
        fill_project2,
        fill_project3,
        fill_project4,
        fill_project5,
        fill_project6,
        fill_project7,
        fill_project8,
    ]
    
    for i, filler in enumerate(project_fillers):
        if i < len(prs.slides):
            print(f"🔄 正在处理第 {i+1} 张幻灯片 (项目{i+1})...")
            filler(prs.slides[i])
    
    output_path = r"f:\个人作品\legal-rag-qa-system\A2A_PROTOCOL_AI_AGENT_V14_完美清晰版.pptx"
    prs.save(output_path)
    print(f"\n✅ PPT表格填充完成！已保存到: {output_path}")
    print("\n📋 处理说明:")
    print("  - 所有单元格使用纯白色字体 (RGB 255,255,255)")
    print("  - 表头: 12号字, 加粗")
    print("  - 内容: 10号字, 常规")
    print("  - 字体: 微软雅黑, 居中对齐")
    print("  - 包含最新AI前沿技术细节内容")

if __name__ == "__main__":
    main()