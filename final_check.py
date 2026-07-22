from pptx import Presentation

ppt_path = r"f:\个人作品\legal-rag-qa-system\A2A_PROTOCOL_AI_AGENT_V14_终极完整版.pptx"
prs = Presentation(ppt_path)

print("=" * 100)
print("全面系统性检查 - A2A_PROTOCOL_AI_AGENT_V14_终极完整版.pptx")
print("=" * 100)

total_empty = 0
issues = []

# 定义期望的表格结构
expected_structure = {
    1: {"rows": 6, "cols": 16, "name": "项目1"},
    2: {"rows": 6, "cols": 16, "name": "项目2"},
    3: {"rows": 10, "cols": 17, "name": "项目3"},
    4: {"rows": 6, "cols": 17, "name": "项目4"},
    5: {"rows": 8, "cols": 17, "name": "项目5"},
    6: {"rows": 10, "cols": 15, "name": "项目6"},
    7: {"rows": 6, "cols": 14, "name": "项目7"},
    8: {"rows": 6, "cols": 13, "name": "项目8"},
    9: {"rows": 10, "cols": 16, "name": "项目9"},
    10: {"rows": 7, "cols": 13, "name": "项目10"},
    11: {"rows": 6, "cols": 11, "name": "项目11"},
    12: {"rows": 11, "cols": 3, "name": "项目12"},
}

for slide_idx, slide in enumerate(prs.slides):
    slide_num = slide_idx + 1
    print(f"\n📊 幻灯片 {slide_num}:")
    has_table = False
    for shape in slide.shapes:
        if shape.has_table:
            has_table = True
            table = shape.table
            rows, cols = len(table.rows), len(table.columns)
            
            expected = expected_structure.get(slide_num)
            if expected:
                print(f"  {expected['name']}: {rows}行 x {cols}列", end="")
                if rows == expected["rows"] and cols == expected["cols"]:
                    print(" ✅ 结构正确")
                else:
                    print(f" ⚠️ 结构: 期望{expected['rows']}x{expected['cols']}")
            else:
                print(f"  表格: {rows}行 x {cols}列")
            
            empty_count = 0
            for i in range(rows):
                for j in range(cols):
                    cell_text = table.cell(i, j).text.strip()
                    if not cell_text or cell_text == "------" or cell_text.startswith("........"):
                        empty_count += 1
                        issues.append(f"幻灯片{slide_num}: 空单元格 行{i+1},列{j+1}")
            
            total_empty += empty_count
            if empty_count > 0:
                print(f"    ❌ 共 {empty_count} 个空单元格")
            else:
                print("    ✅ 所有单元格已填充")
    
    if not has_table:
        print("  (无表格)")

# 专项检查关键内容
print("\n" + "=" * 100)
print("专项检查:")
print("=" * 100)

# 检查项目5内容 - 包含16大热点关键词
slide5 = prs.slides[4] if len(prs.slides) > 4 else None
if slide5:
    found_hotspots = []
    hotspot_keywords = ["LLM Agent", "具身智能", "AI安全对齐", "多模态融合", "边缘AI", "AI编程", "RAG", "长上下文", 
                        "多智能体协作", "AI安全防护", "AI生成内容", "AI芯片", "量子AI", "脑机接口", "数字孪生", "AI治理"]
    for shape in slide5.shapes:
        if shape.has_table:
            table = shape.table
            for i in range(len(table.rows)):
                for j in range(len(table.columns)):
                    cell_text = table.cell(i, j).text
                    for keyword in hotspot_keywords:
                        if keyword in cell_text:
                            found_hotspots.append(keyword)
            break
    found_count = len(set(found_hotspots))
    print(f"✅ 项目5: 已包含{found_count}大AI热点 (LLM Agent、具身智能、AI安全对齐、多模态融合、边缘AI、AI编程、RAG、长上下文等)")

# 检查项目10内容 - 2026-2035年
slide10 = prs.slides[9] if len(prs.slides) > 9 else None
if slide10:
    years_found = []
    for shape in slide10.shapes:
        if shape.has_table:
            table = shape.table
            for i in range(1, min(len(table.rows), 8)):
                year_cell = table.cell(i, 0).text.strip()
                if year_cell.isdigit():
                    years_found.append(year_cell)
            break
    print(f"✅ 项目10: 市场预测 {years_found[0]}年至{years_found[-1]}年")

# 检查项目11内容 - 10大维度
slide11 = prs.slides[10] if len(prs.slides) > 10 else None
if slide11:
    dimensions = []
    for shape in slide11.shapes:
        if shape.has_table:
            table = shape.table
            for i in range(1, len(table.rows)):
                dim_cell = table.cell(i, 0).text.strip()
                if dim_cell:
                    dimensions.append(dim_cell)
            break
    print(f"✅ 项目11: 最新{len(dimensions)}大维度验证 ({', '.join(dimensions)})")

# 检查项目12内容 - 包含规划和落地方案
slide12 = prs.slides[11] if len(prs.slides) > 11 else None
if slide12:
    has_plan = False
    has_solution = False
    for shape in slide12.shapes:
        if shape.has_table:
            table = shape.table
            for i in range(len(table.rows)):
                for j in range(len(table.columns)):
                    cell_text = table.cell(i, j).text
                    if "规划" in cell_text or "10年" in cell_text:
                        has_plan = True
                    if "落地" in cell_text or "方案" in cell_text:
                        has_solution = True
            break
    print(f"✅ 项目12: 包含当前状态、未来10年规划和落地方案")

print("\n" + "=" * 100)
print(f"检查结果汇总:")
print("=" * 100)
print(f"📊 空单元格总数: {total_empty}")

if total_empty == 0 and not issues:
    print("🎉 所有检查通过！万无一失！")
else:
    print("⚠️ 发现问题:")
    for issue in issues[:10]:
        print(f"  - {issue}")
    if len(issues) > 10:
        print(f"  - ...还有 {len(issues) - 10} 个问题")