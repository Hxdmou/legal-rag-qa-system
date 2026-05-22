from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml import parse_xml

def fill_cell(cell, text, font_size=9, bold=False):
    cell.text_frame.clear()
    para = cell.text_frame.add_paragraph()
    run = para.add_run()
    run.text = str(text) if text else "---"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(255, 255, 255)
    run.font.name = '微软雅黑'
    para.alignment = PP_ALIGN.CENTER
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

def update_project3(table):
    """更新项目3为完整10层架构"""
    # 完整的10层数据
    data = [
        ["层级", "名称", "核心功能", "技术实现", "通信协议", "数据格式", "加密方式", "性能指标", "可靠性", "扩展性", "安全性", "监控能力", "调试工具", "文档完整度", "社区活跃度", "学习成本", "推荐指数"],
        ["L1", "传输层", "消息路由+负载均衡+熔断", "HTTP/3+QUIC+Envoy", "HTTP/3", "Protobuf/JSON", "TLS1.3/AES-GCM", "<30ms", "99.999%", "线性/弹性", "AES-256/GCM", "Prometheus/Thanos", "grpcurl/tcpdump", "95%", "高", "中", "10"],
        ["L2", "协议层", "通信协议+认证+授权", "A2A-Spec v2.0+OAuth2", "WebSocket/SSE", "JSON/JSON-RPC", "JWT/OAuth2/mTLS", "<20ms", "99.99%", "插件化/热更新", "mTLS/OIDC", "Grafana/Tempo", "Postman/Charles", "98%", "高", "高", "10"],
        ["L3", "智能层", "AI决策+RAG+Agent调度", "GPT 5.0+DeepSeek V5.0+LangChain", "gRPC/SSE", "JSON/Protobuf", "AES-256", "<100ms", "99.9%", "模块化/动态扩展", "RBAC/ABAC", "ELK/LangSmith", "LangSmith/Debugger", "90%", "中", "高", "9.5"],
        ["L4", "应用层", "业务逻辑+工作流+编排", "多智能体协作+Flowise", "gRPC/HTTP", "Protobuf/JSON", "mTLS", "<30ms", "99.99%", "微服务/K8s", "ABAC/Security", "Jaeger/Zipkin", "OpenTelemetry", "85%", "中", "中", "9"],
        ["L5", "表现层", "用户交互+UI+实时渲染", "React/Vue+AI生成UI", "HTTP/3/WebSocket", "JSON/GraphQL", "TLS1.3", "<10ms", "99.99%", "CDN/边缘", "WAF/DDoS", "Datadog/NewRelic", "ChromeDev/Fiddler", "95%", "高", "低", "9.5"],
        ["L6", "数据层", "存储计算+缓存+搜索", "Redis+PostgreSQL+Elasticsearch", "TCP/gRPC", "二进制/JSON", "SM4/AES", "<5ms", "99.999%", "分片/读写分离", "数据脱敏/加密", "Prometheus/PGHero", "pgAdmin/redis-cli", "90%", "中", "中", "9"],
        ["L7", "安全层", "身份认证+审计+合规", "零信任架构+OPA", "mTLS", "X509/JWT", "国密SM2/SM4", "<15ms", "99.99%", "动态策略/热更新", "审计日志/SOC", "SIEM/Splunk", "Keycloak/Okta", "88%", "中", "高", "9.8"],
        ["L8", "监控层", "可观测性+告警+追踪", "APM+全链路追踪", "HTTP/gRPC", "OTLP/JSON", "TLS", "实时", "99.9%", "自动告警/扩展", "链路加密", "Jaeger/Grafana", "Grafana/Loki", "92%", "高", "中", "9"],
        ["L9", "网关层", "流量管理+限流+WAF", "API网关+Kong/APISIX", "HTTP/2", "JSON/Protobuf", "JWT+RSA", "<25ms", "99.99%", "负载均衡/灰度", "限流熔断/WAF", "Kong/APISIX", "curl/httpie", "93%", "高", "中", "9.5"],
        ["L10", "边缘层", "边缘计算+AI推理+缓存", "EdgeAI+FastAPI+Redis", "QUIC/HTTP", "Protobuf", "TLS1.3", "<10ms", "99.99%", "边缘部署/CDN", "本地化加密", "Prometheus", "CLI/debugpy", "85%", "低", "高", "8.5"],
    ]
    
    # 添加缺失的行
    while len(table.rows) < len(data):
        tbl = table._tbl
        tr = tbl.tr_lst[0]
        new_tr = parse_xml(tr.xml)
        tbl.append(new_tr)
    
    # 填充数据
    for i in range(len(data)):
        for j in range(len(data[i])):
            if i < len(table.rows) and j < len(table.columns):
                fill_cell(table.cell(i, j), data[i][j], 12 if i == 0 else 9, i == 0)

def update_project6(table):
    """更新项目6为15个应用场景"""
    data = [
        ["领域", "典型场景", "智能体角色", "大模型方案", "预期收益", "技术难度", "实施周期", "投资规模", "ROI周期", "风险等级", "市场成熟度", "竞争格局", "政策支持", "人才需求", "推荐度", "优先级"],
        ["企业办公", "智能助理+自动化+文档处理", "任务自动化", "GPT 5.0+DeepSeek V5.0", "效率+50%", "中", "3个月", "100万", "18个月", "低", "高", "激烈", "强", "中", "9.5", "P0"],
        ["金融服务", "风控合规+智能投顾+量化交易", "风险评估", "GPT 5.0+DeepSeek V5.0+千问3.6Plus", "准确率+30%", "高", "6个月", "500万", "12个月", "中", "中高", "中等", "强", "高", "9", "P0"],
        ["医疗健康", "辅助诊断+药物研发+医学影像", "医学分析", "Gemini 2.5+DeepSeek V5.0+千问3.6Plus", "误诊率-40%", "极高", "12个月", "1000万", "24个月", "高", "中", "中等", "强", "极高", "8.5", "P1"],
        ["智能制造", "质检维护+预测性维护+工艺优化", "设备监控", "GPT 5.0+DeepSeek V5.0+Claude 3.5", "故障预警+60%", "中高", "9个月", "800万", "15个月", "中", "中高", "中等", "强", "高", "9", "P0"],
        ["教育培训", "个性化学习+智能辅导+内容生成", "智能推荐", "GPT 5.0+千问3.6Plus+Gemini 2.5", "效果+40%", "中", "6个月", "200万", "12个月", "低", "高", "激烈", "强", "中", "9", "P1"],
        ["零售电商", "智能营销+供应链优化+用户画像", "精准推荐", "GPT 5.0+DeepSeek V5.0+千问3.6Plus", "转化+35%", "中", "4个月", "150万", "10个月", "低", "高", "激烈", "中", "中", "8.5", "P1"],
        ["物流运输", "路径优化+智能调度+仓储管理", "智能调度", "GPT 5.0+DeepSeek V5.0", "效率+45%", "中高", "8个月", "300万", "14个月", "中", "中高", "中等", "中", "高", "8.5", "P1"],
        ["能源电力", "智能调度+负荷预测+故障诊断", "负荷预测", "GPT 5.0+Claude 3.5", "能耗-25%", "高", "10个月", "600万", "18个月", "中", "中", "中等", "强", "高", "8.5", "P1"],
        ["智慧城市", "城市管理+数据融合+应急响应", "数据融合", "GPT 5.0+DeepSeek V5.0+千问3.6Plus", "效率+30%", "高", "12个月", "1500万", "24个月", "中", "中", "中等", "强", "极高", "8", "P2"],
        ["农业科技", "智慧农业+精准种植+病虫害检测", "精准种植", "GPT 5.0+Gemini 2.5", "产量+20%", "中高", "10个月", "400万", "20个月", "中", "中", "中等", "强", "中", "8", "P2"],
        ["法律服务", "智能合同+法律咨询+案件分析", "法律顾问", "GPT 5.0+DeepSeek V5.0+千问3.6Plus", "效率+60%", "高", "8个月", "350万", "16个月", "中", "中", "强", "高", "高", "8.5", "P1"],
        ["媒体娱乐", "内容创作+智能剪辑+个性化推荐", "内容生成", "GPT 5.0+Gemini 2.5+DeepSeek V5.0", "效率+70%", "中", "5个月", "250万", "10个月", "低", "高", "激烈", "中", "中", "8.5", "P1"],
        ["交通出行", "智能调度+路径规划+自动驾驶", "交通管理", "GPT 5.0+Claude 3.5+千问3.6Plus", "效率+40%", "极高", "14个月", "1200万", "30个月", "高", "中", "强", "极高", "极高", "8", "P2"],
        ["游戏娱乐", "NPC智能+关卡生成+剧情创作", "游戏AI", "GPT 5.0+DeepSeek V5.0+Gemini 2.5", "体验+50%", "中高", "7个月", "450万", "12个月", "中", "高", "中等", "高", "高", "8.5", "P1"],
        ["科研创新", "文献分析+实验设计+数据挖掘", "科研助手", "GPT 5.0+千问3.6Plus+DeepSeek V5.0", "效率+80%", "高", "10个月", "800万", "24个月", "中", "中", "强", "极高", "极高", "8", "P2"],
    ]
    
    while len(table.rows) < len(data):
        tbl = table._tbl
        tr = tbl.tr_lst[0]
        new_tr = parse_xml(tr.xml)
        tbl.append(new_tr)
    
    for i in range(len(data)):
        for j in range(len(data[i])):
            if i < len(table.rows) and j < len(table.columns):
                fill_cell(table.cell(i, j), data[i][j], 12 if i == 0 else 9, i == 0)

def update_project7(table):
    """更新项目7为10个案例"""
    data = [
        ["案例", "行业", "应用场景", "核心成果", "技术方案", "大模型选型", "投资规模", "ROI周期", "实施周期", "风险控制", "可复制性", "创新性", "商业价值", "推荐度", "后续计划"],
        ["某大型银行", "金融", "智能风控+反欺诈+智能客服", "欺诈识别+40%,成本-60%", "LLM+RAG+知识图谱", "GPT 5.0+DeepSeek V5.0", "200万", "2个月", "6个月", "高", "高", "高", "极高", "10", "全国推广+跨产品线"],
        ["某电商平台", "零售", "智能客服+个性化推荐+供应链优化", "成本-60%,转化+35%", "多智能体协作+RAG", "DeepSeek V5.0+千问3.6Plus", "150万", "3个月", "3个月", "中", "高", "中", "高", "9.5", "持续优化+国际化"],
        ["某医疗机构", "医疗", "辅助诊断+影像分析+药物研发", "准确率+25%,效率+50%", "Vision+LLM+RAG", "Gemini 2.5+DeepSeek V5.0", "300万", "2.4个月", "12个月", "高", "中", "高", "极高", "9", "二期扩展+AI制药"],
        ["某制造企业", "制造", "设备维护+预测性维护+质量检测", "停机-50%,质检效率+80%", "IoT+AI+RAG", "GPT 5.0+Claude 3.5", "800万", "15个月", "9个月", "中", "高", "中", "高", "9", "全产线覆盖+智能工厂"],
        ["某保险公司", "金融", "智能理赔+核保+客户服务", "效率+70%,成本-50%", "LLM+知识图谱+RAG", "GPT 5.0+千问3.6Plus", "250万", "3个月", "5个月", "高", "高", "高", "高", "9.5", "跨产品线+全球服务"],
        ["某教育机构", "教育", "智能辅导+自适应学习+内容生成", "成绩+20%,效率+40%", "LLM+个性化+RAG", "GPT 5.0+Gemini 2.5", "180万", "4个月", "6个月", "低", "高", "中", "高", "9", "课程扩展+AI教师"],
        ["某物流公司", "物流", "智能调度+路径优化+仓储管理", "效率+45%,成本-30%", "LLM+IoT+优化算法", "GPT 5.0+DeepSeek V5.0", "300万", "14个月", "8个月", "中", "高", "中", "高", "8.5", "全国网络+无人配送"],
        ["某能源集团", "能源", "智能调度+负荷预测+故障诊断", "能耗-25%,效率+40%", "LLM+IoT+预测模型", "GPT 5.0+Claude 3.5", "600万", "18个月", "10个月", "中", "中", "高", "高", "8.5", "全国推广+新能源"],
        ["某媒体平台", "媒体", "内容创作+智能剪辑+个性化推荐", "效率+70%,用户+50%", "LLM+CV+推荐系统", "GPT 5.0+Gemini 2.5", "250万", "10个月", "5个月", "低", "高", "高", "高", "9", "多平台+全球化"],
        ["某律师事务所", "法律", "智能合同+法律咨询+案件分析", "效率+60%,准确率+40%", "LLM+知识图谱+NLP", "GPT 5.0+千问3.6Plus", "350万", "16个月", "8个月", "中", "高", "高", "高", "8.5", "全国服务+AI法官"],
    ]
    
    while len(table.rows) < len(data):
        tbl = table._tbl
        tr = tbl.tr_lst[0]
        new_tr = parse_xml(tr.xml)
        tbl.append(new_tr)
    
    for i in range(len(data)):
        for j in range(len(data[i])):
            if i < len(table.rows) and j < len(table.columns):
                fill_cell(table.cell(i, j), data[i][j], 12 if i == 0 else 9, i == 0)

def update_project8(table):
    """更新项目8为15个对比维度"""
    data = [
        ["特性", "A2A协议", "传统API", "消息队列", "gRPC", "微服务", "WebSocket", "GraphQL", "RESTful", "SOAP", "Webhook", "SSE", "推荐度"],
        ["实时性", "极强(<10ms)", "弱(>100ms)", "强(<50ms)", "强(<30ms)", "弱(>100ms)", "极强(<10ms)", "弱(>100ms)", "弱(>100ms)", "弱(>200ms)", "中等(<80ms)", "极强(<10ms)", "10"],
        ["双向通信", "极强", "弱", "弱", "强", "弱", "极强", "弱", "弱", "弱", "弱", "极强", "10"],
        ["智能路由", "极强(AutoML)", "无", "中等", "中等", "弱", "中等", "无", "无", "无", "中等", "弱", "10"],
        ["状态管理", "极强(Redis集成)", "弱", "中等", "中等", "中等", "强", "弱", "中等", "弱", "中等", "弱", "10"],
        ["LLM集成", "极强(GPT5/DeepSeek)", "无", "无", "中等", "中等", "中等", "无", "无", "无", "中等", "弱", "10"],
        ["多智能体", "极强(A2A协议)", "无", "弱", "弱", "中等", "中等", "无", "无", "无", "弱", "弱", "10"],
        ["安全性", "极高(mTLS/零信任)", "中等", "强", "强", "中等", "强", "中等", "中等", "中等", "中等", "强", "9.5"],
        ["扩展性", "极强(K8s原生)", "强", "中等", "强", "中等", "强", "强", "强", "中等", "中等", "强", "9"],
        ["易用性", "极强", "强", "中等", "中等", "中等", "强", "强", "极强", "弱", "中等", "强", "9.5"],
        ["性能", "极高(<10ms)", "低(>100ms)", "中(<50ms)", "高(<30ms)", "低(>100ms)", "极高(<10ms)", "低(>100ms)", "低(>100ms)", "低(>200ms)", "中(<80ms)", "极高(<10ms)", "10"],
        ["可靠性", "极高(99.999%)", "中(99.9%)", "高(99.99%)", "高(99.99%)", "中(99.9%)", "高(99.99%)", "中(99.9%)", "中(99.9%)", "中(99.9%)", "中(99.9%)", "高(99.99%)", "9.5"],
        ["可观测性", "极强(全链路)", "弱", "中等", "强", "中等", "中等", "弱", "中等", "弱", "弱", "中等", "9.5"],
        ["成本", "低", "低", "中", "中", "高", "低", "中", "低", "高", "低", "低", "9"],
        ["生态", "快速成长", "成熟", "成熟", "成熟", "成熟", "成熟", "成长", "成熟", "衰退", "成长", "成长", "9"],
        ["综合评分", "9.8", "6.2", "7.5", "7.8", "7.0", "8.5", "7.2", "6.8", "5.5", "6.0", "8.0", "A2A最优"],
    ]
    
    while len(table.rows) < len(data):
        tbl = table._tbl
        tr = tbl.tr_lst[0]
        new_tr = parse_xml(tr.xml)
        tbl.append(new_tr)
    
    for i in range(len(data)):
        for j in range(len(data[i])):
            if i < len(table.rows) and j < len(table.columns):
                fill_cell(table.cell(i, j), data[i][j], 12 if i == 0 else 9, i == 0)

def update_project9(table):
    """更新项目9为完整协议对比"""
    data = [
        ["特性", "A2A协议", "传统API", "消息队列", "gRPC", "微服务", "WebSocket", "GraphQL", "RESTful", "SOAP", "Webhook", "SSE", "MQTT", "AMQP", "NATS", "对比结论"],
        ["实时性", "<10ms", ">100ms", "<50ms", "<30ms", ">100ms", "<10ms", ">100ms", ">100ms", ">200ms", "<80ms", "<10ms", "<50ms", "<30ms", "<10ms", "A2A最优"],
        ["双向通信", "支持", "不支持", "不支持", "支持", "不支持", "支持", "不支持", "不支持", "不支持", "不支持", "支持", "支持", "支持", "支持", "A2A最优"],
        ["智能路由", "AutoML", "无", "简单", "简单", "手动", "简单", "无", "无", "无", "简单", "无", "简单", "简单", "简单", "A2A最优"],
        ["状态管理", "Redis集成", "无", "有限", "有限", "有限", "有限", "无", "有限", "无", "有限", "无", "有限", "有限", "有限", "A2A最优"],
        ["LLM集成", "原生支持", "无", "无", "有限", "有限", "有限", "无", "无", "无", "有限", "无", "无", "无", "无", "A2A最优"],
        ["多智能体", "原生支持", "无", "无", "有限", "有限", "有限", "无", "无", "无", "有限", "无", "无", "无", "无", "A2A最优"],
        ["安全性", "零信任+mTLS", "OAuth2", "TLS", "TLS", "OAuth2", "TLS", "OAuth2", "TLS", "TLS", "TLS", "TLS", "TLS", "TLS", "TLS", "A2A最优"],
        ["扩展性", "K8s原生", "中等", "中等", "高", "高", "高", "中等", "高", "低", "低", "高", "高", "高", "高", "A2A最优"],
        ["易用性", "高", "极高", "中", "中", "中", "高", "高", "极高", "低", "中等", "高", "高", "中", "高", "A2A最优"],
        ["综合评价", "9.8", "6.2", "7.5", "7.8", "7.0", "8.5", "7.2", "6.8", "5.5", "6.0", "8.0", "8.2", "7.5", "8.5", "A2A领先"],
    ]
    
    while len(table.rows) < len(data):
        tbl = table._tbl
        tr = tbl.tr_lst[0]
        new_tr = parse_xml(tr.xml)
        tbl.append(new_tr)
    
    for i in range(len(data)):
        for j in range(len(data[i])):
            if i < len(table.rows) and j < len(table.columns):
                fill_cell(table.cell(i, j), data[i][j], 12 if i == 0 else 9, i == 0)

def update_project10(table):
    """更新项目10为完整市场预测"""
    data = [
        ["年份", "市场规模", "增长率", "主要驱动力", "投资规模", "技术突破", "应用场景", "竞争格局", "政策环境", "人才供给", "技术成熟度", "商业化程度", "投资建议"],
        ["2026", "2000亿", "85%", "GPT 5.0发布", "600亿", "GPT 5.0/DeepSeek V5.0", "企业服务爆发", "格局初定", "政策完善", "紧缺", "高", "高", "重仓"],
        ["2027", "3800亿", "90%", "具身智能崛起", "1200亿", "GPT 5.5/人形机器人", "智能硬件+工业质检", "巨头主导", "监管落地", "紧缺", "高", "极高", "长期"],
        ["2028", "6500亿", "71%", "AI机器人爆发", "2000亿", "GPT 6.0/Optimus 2.0", "服务机器人+物流配送", "生态竞争", "全球标准", "平衡", "极高", "极高", "稳健"],
        ["2029", "10000亿", "54%", "脑机接口突破", "3000亿", "DeepSeek V6.0/BCI", "家庭助手+医疗护理", "寡头垄断", "伦理规范", "充足", "极高", "极高", "价值"],
        ["2030", "15000亿", "50%", "AGI原型落地", "4500亿", "GPT 7.0/通用AGI", "人机融合+教育陪伴", "全球格局", "国际治理", "充足", "极高", "极高", "战略"],
        ["2031", "22000亿", "47%", "量子AI商用", "6000亿", "量子AI融合/IBM Q1000", "智能经济+农业机器人", "生态竞争", "标准统一", "充足", "极高", "极高", "价值"],
        ["2032", "32000亿", "45%", "边缘智能爆发", "8000亿", "GPT 8.0/边缘大脑", "万物互联+城市服务", "平台竞争", "全球协作", "充足", "极高", "极高", "稳健"],
        ["2033", "45000亿", "41%", "脑机融合深化", "10000亿", "DeepSeek V7.0/神经接口", "人机共生+养老服务", "寡头格局", "伦理框架", "充足", "极高", "极高", "长期"],
        ["2034", "62000亿", "38%", "量子AGI突破", "12000亿", "量子AGI融合", "智能文明+太空探索", "全球垄断", "国际治理", "充足", "极高", "极高", "战略"],
        ["2035", "85000亿", "37%", "AGI全面普及", "15000亿", "通用AGI/超级智能", "智能时代+全面智能化", "全球一体", "AI治理", "充足", "极高", "极高", "长期价值"],
    ]
    
    while len(table.rows) < len(data):
        tbl = table._tbl
        tr = tbl.tr_lst[0]
        new_tr = parse_xml(tr.xml)
        tbl.append(new_tr)
    
    for i in range(len(data)):
        for j in range(len(data[i])):
            if i < len(table.rows) and j < len(table.columns):
                fill_cell(table.cell(i, j), data[i][j], 12 if i == 0 else 9, i == 0)

def update_project11(table):
    """更新项目11为完整技术评估"""
    data = [
        ["维度", "价值描述", "量化指标", "测试环境", "数据规模", "准确率", "响应时间", "资源消耗", "可扩展性", "安全性", "商业价值"],
        ["推理能力", "复杂任务+逻辑推理+数学", "96%", "GPT 5.0", "500万token", "93%", "<150ms", "中", "高", "高", "极高"],
        ["多模态", "图文音视频+3D理解", "92%", "Gemini 2.5", "50万样本", "90%", "<400ms", "高", "中", "中", "高"],
        ["长上下文", "百万Token+记忆+检索", "88%", "Claude 3.5", "1M+ token", "87%", "<800ms", "高", "中", "中", "极高"],
        ["工具调用", "API+代码执行+调试", "99%", "GPT 5.0", "10000次", "99.5%", "<80ms", "低", "高", "高", "高"],
        ["自我反思", "错误纠正+优化+迭代", "85%", "DeepSeek V5.0", "10万样本", "82%", "<1.5s", "高", "中", "中", "极高"],
        ["综合评分", "全面能力评估", "91%", "GPT 5.0/DeepSeek", "综合测试", "90%", "<300ms", "中", "高", "高", "极高"],
    ]
    
    for i in range(len(data)):
        for j in range(len(data[i])):
            if i < len(table.rows) and j < len(table.columns):
                fill_cell(table.cell(i, j), data[i][j], 12 if i == 0 else 9, i == 0)

def update_project12(table):
    """更新项目12为完整总结"""
    data = [
        ["项目", "核心内容", "未来10年规划"],
        ["1", "封面升级", "季度更新迭代"],
        ["2", "目录优化", "AI智能导航"],
        ["3", "协议架构", "A2A 3.0升级(2027)"],
        ["4", "能力矩阵", "扩展至10层(2028)"],
        ["5", "16大热点", "年度更新跟踪"],
        ["6", "应用场景", "扩展至20领域(2029)"],
        ["7", "AI机器人案例", "2028年服务机器人普及"],
        ["8", "技术对比", "季度对比报告"],
        ["9", "市场预测", "2026-2035年滚动预测"],
        ["10", "10大维度", "扩展至15维度(2030)"],
        ["合计", "完整AI方案", "分阶段落地"],
    ]
    
    for i in range(len(data)):
        for j in range(len(data[i])):
            if i < len(table.rows) and j < len(table.columns):
                fill_cell(table.cell(i, j), data[i][j], 12 if i == 0 else 9, i == 0)

def main():
    input_path = r"f:\个人作品\legal-rag-qa-system\A2A_PROTOCOL_AI_AGENT_V14_完整填充版.pptx"
    output_path = r"f:\个人作品\legal-rag-qa-system\A2A_PROTOCOL_AI_AGENT_V14_终极完整版.pptx"
    
    print("🔄 正在从旧版本更新到新版本...")
    prs = Presentation(input_path)
    
    # 更新所有项目
    if len(prs.slides) > 2:
        for shape in prs.slides[2].shapes:
            if shape.has_table:
                update_project3(shape.table)
                print("✅ 项目3已更新")
                break
    
    if len(prs.slides) > 5:
        for shape in prs.slides[5].shapes:
            if shape.has_table:
                update_project6(shape.table)
                print("✅ 项目6已更新")
                break
    
    if len(prs.slides) > 6:
        for shape in prs.slides[6].shapes:
            if shape.has_table:
                update_project7(shape.table)
                print("✅ 项目7已更新")
                break
    
    if len(prs.slides) > 7:
        for shape in prs.slides[7].shapes:
            if shape.has_table:
                update_project8(shape.table)
                print("✅ 项目8已更新")
                break
    
    if len(prs.slides) > 8:
        for shape in prs.slides[8].shapes:
            if shape.has_table:
                update_project9(shape.table)
                print("✅ 项目9已更新")
                break
    
    if len(prs.slides) > 9:
        for shape in prs.slides[9].shapes:
            if shape.has_table:
                update_project10(shape.table)
                print("✅ 项目10已更新")
                break
    
    if len(prs.slides) > 10:
        for shape in prs.slides[10].shapes:
            if shape.has_table:
                update_project11(shape.table)
                print("✅ 项目11已更新")
                break
    
    if len(prs.slides) > 11:
        for shape in prs.slides[11].shapes:
            if shape.has_table:
                update_project12(shape.table)
                print("✅ 项目12已更新")
                break
    
    prs.save(output_path)
    print(f"\n✅ 已完成！保存到: {output_path}")
    print("\n💡 请确保您打开的是这个文件！")

if __name__ == "__main__":
    main()