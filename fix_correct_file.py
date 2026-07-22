from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml import parse_xml

def fill_cell(cell, text, font_size=9, bold=False):
    cell.text_frame.clear()
    para = cell.text_frame.add_paragraph()
    run = para.add_run()
    run.text = str(text) if text else "---"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(255, 255, 255)
    run.font.name = '微软雅黑'
    para.alignment = PP_ALIGN.CENTER
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

def update_project3(table):
    """更新项目3为完整内容"""
    data = [
        ["协议层级", "业务价值", "技术亮点", "成功指标", "适用场景", "性能指标", "投资回报", "技术架构", "关键挑战"],
        ["应用层", "高", "多模态融合", "95%", "企业服务", "小于10毫秒", "中规模", "Function Calling", "高"],
        ["消息层", "高", "异步通信", "99%", "微服务", "小于50毫秒", "中规模", "Event-Driven", "中高"],
        ["传输层", "高", "HTTP/3+QUIC", "99.99%", "大规模", "小于5毫秒", "大规模", "gRPC+Envoy", "中"],
        ["安全层", "极高", "零信任架构", "99.999%", "金融医疗", "小于15毫秒", "高成本", "mTLS+OPA", "中高"],
        ["发现层", "中高", "服务发现", "99%", "云原生", "小于30毫秒", "中规模", "Consul+DNS", "低"],
        ["事件层", "高", "实时事件流", "99.9%", "IoT", "小于20毫秒", "大规模", "Kafka+Pulsar", "中"],
        ["数据层", "极高", "分布式存储", "99.999%", "大数据", "小于10毫秒", "大规模", "Redis+PostgreSQL", "中"],
        ["控制层", "中", "流量控制", "99%", "网关", "小于25毫秒", "中规模", "Kong+APISIX", "低"],
        ["编排层", "高", "工作流编排", "99%", "业务流程", "小于100毫秒", "大规模", "Flowise+Airflow", "中高"],
        ["语义层", "高", "AI理解", "95%", "智能助理", "小于150毫秒", "大规模", "GPT 5.0+DeepSeek", "高"],
        ["Agent通信", "极高", "A2A协议", "99.9%", "多智能体", "小于50毫秒", "大规模", "LangChain+gRPC", "中高"],
        ["Agent发现", "中", "智能发现", "98%", "智能生态", "小于100毫秒", "中规模", "AI-Powered", "中"],
        ["边缘层(L10)", "高", "边缘AI推理", "99.9%", "边缘计算", "小于10毫秒", "大规模", "EdgeAI+FastAPI", "中"],
    ]
    
    # 添加缺失的行
    while len(table.rows) < len(data):
        tbl = table._tbl
        tr = tbl.tr_lst[0]
        new_tr = parse_xml(tr.xml)
        tbl.append(new_tr)
    
    # 填充数据
    for i in range(len(data)):
        for j in range(len(data[i])):
            if i < len(table.rows) and j < len(table.columns):
                fill_cell(table.cell(i, j), data[i][j], 12 if i == 0 else 9, i == 0)

def main():
    # 正确的文件路径！
    input_path = r"f:\个人作品\legal-rag-qa-system\docs\presentations\A2A_ENTERPRISE_PPT_V12.pptx"
    output_path = r"f:\个人作品\legal-rag-qa-system\docs\presentations\A2A_ENTERPRISE_PPT_V12_完整版.pptx"
    
    print("🔄 正在修复正确的文件...")
    prs = Presentation(input_path)
    
    # 更新项目3（幻灯片3）
    if len(prs.slides) > 2:
        for shape in prs.slides[2].shapes:
            if shape.has_table:
                update_project3(shape.table)
                print("✅ 项目3已更新")
                break
    
    prs.save(output_path)
    print(f"\n✅ 已完成！保存到: {output_path}")
    print("\n💡 请打开这个文件！")

if __name__ == "__main__":
    main()