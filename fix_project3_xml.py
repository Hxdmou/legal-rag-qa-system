from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml import parse_xml

def fill_cell(cell, text, font_size=9, bold=False):
    cell.text_frame.clear()
    para = cell.text_frame.add_paragraph()
    run = para.add_run()
    run.text = str(text) if text else "---"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(255, 255, 255)
    run.font.name = '微软雅黑'
    para.alignment = PP_ALIGN.CENTER
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

def add_table_row(table):
    """通过XML添加新行"""
    # 获取表格的XML
    tbl = table._tbl
    
    # 获取第一行的XML作为模板
    tr = tbl.tr_lst[0]
    
    # 复制第一行
    new_tr = parse_xml(tr.xml)
    
    # 添加到表格
    tbl.append(new_tr)
    
    print(f"  已添加新行，现在表格有 {len(table.rows)} 行")

def fix_project3_xml(prs):
    """通过XML修复项目3"""
    if len(prs.slides) > 2:
        slide = prs.slides[2]
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                print(f"📌 原始表格: {len(table.rows)}行 × {len(table.columns)}列")
                
                # 检查是否需要添加行
                if len(table.rows) < 11:
                    print(f"🔄 需要添加 {11 - len(table.rows)} 行")
                    
                    # 添加新行
                    for _ in range(11 - len(table.rows)):
                        try:
                            add_table_row(table)
                        except Exception as e:
                            print(f"  添加行失败: {e}")
                            # 尝试直接操作XML
                            tbl = table._tbl
                            tr = tbl.tr_lst[0]
                            new_tr = parse_xml(tr.xml)
                            tbl.append(new_tr)
                            print(f"  通过XML添加行成功")
                    
                    print(f"✅ 现在表格: {len(table.rows)}行 × {len(table.columns)}列")
                
                # 完整的10层数据
                data = [
                    ["层级", "名称", "核心功能", "技术实现", "通信协议", "数据格式", "加密方式", "性能指标", "可靠性", "扩展性", "安全性", "监控能力", "调试工具", "文档完整度", "社区活跃度", "学习成本", "推荐指数"],
                    ["L1", "传输层", "消息路由+负载均衡+熔断", "HTTP/3+QUIC+Envoy", "HTTP/3", "Protobuf/JSON", "TLS1.3/AES-GCM", "<30ms", "99.999%", "线性/弹性", "AES-256/GCM", "Prometheus/Thanos", "grpcurl/tcpdump", "95%", "高", "中", "10"],
                    ["L2", "协议层", "通信协议+认证+授权", "A2A-Spec v2.0+OAuth2", "WebSocket/SSE", "JSON/JSON-RPC", "JWT/OAuth2/mTLS", "<20ms", "99.99%", "插件化/热更新", "mTLS/OIDC", "Grafana/Tempo", "Postman/Charles", "98%", "高", "高", "10"],
                    ["L3", "智能层", "AI决策+RAG+Agent调度", "GPT 5.0+DeepSeek V5.0+LangChain", "gRPC/SSE", "JSON/Protobuf", "AES-256", "<100ms", "99.9%", "模块化/动态扩展", "RBAC/ABAC", "ELK/LangSmith", "LangSmith/Debugger", "90%", "中", "高", "9.5"],
                    ["L4", "应用层", "业务逻辑+工作流+编排", "多智能体协作+Flowise", "gRPC/HTTP", "Protobuf/JSON", "mTLS", "<30ms", "99.99%", "微服务/K8s", "ABAC/Security", "Jaeger/Zipkin", "OpenTelemetry", "85%", "中", "中", "9"],
                    ["L5", "表现层", "用户交互+UI+实时渲染", "React/Vue+AI生成UI", "HTTP/3/WebSocket", "JSON/GraphQL", "TLS1.3", "<10ms", "99.99%", "CDN/边缘", "WAF/DDoS", "Datadog/NewRelic", "ChromeDev/Fiddler", "95%", "高", "低", "9.5"],
                    ["L6", "数据层", "存储计算+缓存+搜索", "Redis+PostgreSQL+Elasticsearch", "TCP/gRPC", "二进制/JSON", "SM4/AES", "<5ms", "99.999%", "分片/读写分离", "数据脱敏/加密", "Prometheus/PGHero", "pgAdmin/redis-cli", "90%", "中", "中", "9"],
                    ["L7", "安全层", "身份认证+审计+合规", "零信任架构+OPA", "mTLS", "X509/JWT", "国密SM2/SM4", "<15ms", "99.99%", "动态策略/热更新", "审计日志/SOC", "SIEM/Splunk", "Keycloak/Okta", "88%", "中", "高", "9.8"],
                    ["L8", "监控层", "可观测性+告警+追踪", "APM+全链路追踪", "HTTP/gRPC", "OTLP/JSON", "TLS", "实时", "99.9%", "自动告警/扩展", "链路加密", "Jaeger/Grafana", "Grafana/Loki", "92%", "高", "中", "9"],
                    ["L9", "网关层", "流量管理+限流+WAF", "API网关+Kong/APISIX", "HTTP/2", "JSON/Protobuf", "JWT+RSA", "<25ms", "99.99%", "负载均衡/灰度", "限流熔断/WAF", "Kong/APISIX", "curl/httpie", "93%", "高", "中", "9.5"],
                    ["L10", "边缘层", "边缘计算+AI推理+缓存", "EdgeAI+FastAPI+Redis", "QUIC/HTTP", "Protobuf", "TLS1.3", "<10ms", "99.99%", "边缘部署/CDN", "本地化加密", "Prometheus", "CLI/debugpy", "85%", "低", "高", "8.5"],
                ]
                
                # 填充所有数据
                for i in range(len(data)):
                    for j in range(len(data[i])):
                        if i < len(table.rows) and j < len(table.columns):
                            fill_cell(table.cell(i, j), data[i][j], 12 if i == 0 else 9, i == 0)
                
                print(f"✅ 已填充{len(data)}行数据")
                print(f"✅ 包含L1-L10完整10层架构")
                
                break

def verify_project3(prs):
    """验证项目3内容"""
    print("\n" + "=" * 120)
    print("🔍 验证项目3内容")
    print("=" * 120)
    
    if len(prs.slides) > 2:
        slide = prs.slides[2]
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                print(f"\n表格: {len(table.rows)}行 × {len(table.columns)}列")
                
                # 检查每一层
                layers = []
                for i in range(1, len(table.rows)):
                    layer_id = table.cell(i, 0).text.strip()
                    layer_name = table.cell(i, 1).text.strip()
                    layers.append(f"{layer_id}: {layer_name}")
                
                print(f"\n包含的层级:")
                for layer in layers:
                    print(f"  {layer}")
                
                # 检查是否有L10
                if "L10" in [table.cell(i, 0).text for i in range(len(table.rows))]:
                    print("\n✅ 包含L10边缘层！")
                    return True
                else:
                    print("\n❌ 缺少L10边缘层！")
                    return False
                
                break

def comprehensive_check(prs):
    """全面检查"""
    print("\n" + "=" * 120)
    print("🔍 全面检查")
    print("=" * 120)
    
    total_cells = 0
    empty_cells = 0
    ellipsis_cells = 0
    
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                for i in range(len(table.rows)):
                    for j in range(len(table.columns)):
                        total_cells += 1
                        cell_text = table.cell(i, j).text.strip()
                        if not cell_text:
                            empty_cells += 1
                        if "......" in cell_text or "........" in cell_text:
                            ellipsis_cells += 1
    
    print(f"📊 总单元格数: {total_cells}")
    print(f"❌ 空单元格: {empty_cells}")
    print(f"❌ 省略号单元格: {ellipsis_cells}")
    
    if empty_cells == 0 and ellipsis_cells == 0:
        print("🎉🎉🎉 所有检查通过！")
        return True
    else:
        print("⚠️ 发现问题")
        return False

def main():
    ppt_path = r"f:\个人作品\legal-rag-qa-system\A2A_PROTOCOL_AI_AGENT_V14_终极完整版.pptx"
    prs = Presentation(ppt_path)
    
    print("🔄 正在修复项目3（通过XML添加行）...")
    fix_project3_xml(prs)
    
    output_path = r"f:\个人作品\legal-rag-qa-system\A2A_PROTOCOL_AI_AGENT_V14_终极完整版.pptx"
    prs.save(output_path)
    
    # 重新打开进行验证
    prs = Presentation(output_path)
    has_l10 = verify_project3(prs)
    comprehensive_check(prs)
    
    print(f"\n✅ 已完成！保存到: {output_path}")
    
    if has_l10:
        print("\n🎉 项目3现在包含完整的10层架构（L1-L10）！")
    else:
        print("\n⚠️ 项目3仍然缺少L10边缘层")

if __name__ == "__main__":
    main()