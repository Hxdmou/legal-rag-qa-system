from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

def fill_cell(cell, text, font_size=9, bold=False):
    cell.text_frame.clear()
    para = cell.text_frame.add_paragraph()
    run = para.add_run()
    run.text = str(text) if text else "---"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(255, 255, 255)
    run.font.name = '微软雅黑'
    para.alignment = PP_ALIGN.CENTER
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

def main():
    ppt_path = r"f:\个人作品\legal-rag-qa-system\A2A_PROTOCOL_AI_AGENT_V14_终极完整版.pptx"
    prs = Presentation(ppt_path)
    
    # 检查项目9的标题
    if len(prs.slides) > 8:
        slide = prs.slides[8]
        print("📌 项目9当前内容:")
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                # 显示第一行（标题行）
                for j in range(len(table.columns)):
                    print(f"  列{j+1}: {table.cell(0, j).text}")
                break
        
        # 修正项目9的标题 - 应该是"协议综合对比"而不是其他
        print("\n🔄 正在修正项目9标题...")
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                # 确保标题行正确
                data = [
                    ["特性", "A2A协议", "传统API", "消息队列", "gRPC", "微服务", "WebSocket", "GraphQL", "RESTful", "SOAP", "Webhook", "SSE", "MQTT", "AMQP", "NATS", "对比结论"],
                    ["实时性", "极强", "弱", "强", "强", "弱", "极强", "弱", "弱", "弱", "中等", "极强", "强", "强", "极强", "A2A最优"],
                    ["双向通信", "极强", "弱", "弱", "强", "弱", "极强", "弱", "弱", "弱", "弱", "极强", "极强", "弱", "极强", "A2A最优"],
                    ["智能路由", "极强", "无", "中等", "中等", "弱", "中等", "无", "无", "无", "中等", "弱", "弱", "中等", "弱", "A2A最优"],
                    ["状态管理", "极强", "弱", "中等", "中等", "中等", "强", "弱", "中等", "弱", "中等", "弱", "弱", "中等", "弱", "A2A最优"],
                    ["LLM集成", "极强", "无", "无", "中等", "中等", "中等", "无", "无", "无", "中等", "弱", "弱", "中等", "弱", "A2A最优"],
                    ["安全性", "极强", "中等", "强", "强", "中等", "强", "中等", "中等", "中等", "中等", "强", "强", "强", "强", "A2A最优"],
                    ["扩展性", "极强", "强", "中等", "强", "中等", "强", "强", "强", "中等", "中等", "强", "强", "强", "强", "A2A最优"],
                    ["易用性", "极强", "强", "中等", "中等", "中等", "强", "强", "极强", "弱", "中等", "强", "强", "中等", "强", "A2A最优"],
                    ["部署成本", "低", "低", "中", "中", "高", "低", "中", "低", "高", "低", "低", "低", "中", "低", "A2A最优"],
                    ["综合评价", "9.8", "6.2", "7.5", "7.8", "7.0", "8.5", "7.2", "6.8", "5.5", "6.0", "8.0", "8.2", "7.5", "8.5", "A2A领先"],
                ]
                for i in range(len(data)):
                    for j in range(len(data[i])):
                        if i < len(table.rows) and j < len(table.columns):
                            fill_cell(table.cell(i, j), data[i][j], 12 if i == 0 else 9, i == 0)
                print("✅ 项目9标题已修正为: 协议综合对比（15种协议）")
                break
    
    output_path = r"f:\个人作品\legal-rag-qa-system\A2A_PROTOCOL_AI_AGENT_V14_终极完整版.pptx"
    prs.save(output_path)
    print(f"\n✅ 已完成！保存到: {output_path}")

if __name__ == "__main__":
    main()