from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

def fill_cell(cell, text, font_size=9, bold=False):
    cell.text_frame.clear()
    para = cell.text_frame.add_paragraph()
    run = para.add_run()
    run.text = str(text) if text else "---"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(255, 255, 255)
    run.font.name = '微软雅黑'
    para.alignment = PP_ALIGN.CENTER
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

def fill_all_slides(prs):
    # 幻灯片5 - 项目5: 16大AI热点 (适配8行17列表格)
    if len(prs.slides) > 4:
        slide = prs.slides[4]
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                rows, cols = len(table.rows), len(table.columns)
                data = [
                    ["趋势", "成熟度", "市场影响", "代表技术", "商业价值", "技术难度", "投资热度", "竞争格局", "预期收益", "风险等级", "时间窗口", "人才需求", "落地建议", "投资回报", "10年规划", "优先级", "推荐度"],
                    ["LLM Agent", "高", "革命性", "自主智能体+多任务规划", "极高", "高", "极热", "激烈", "10倍", "中", "1-3年", "紧缺", "优先布局", "18个月回本", "持续研发", "P0", "10"],
                    ["具身智能", "中", "深远", "机器人感知+运动控制", "极高", "极高", "很热", "中等", "20倍", "高", "3-5年", "极缺", "战略投入", "战略价值", "生态闭环", "P1", "9.5"],
                    ["AI安全对齐", "中高", "关键", "RLHF+Constitutional AI", "极高", "高", "很热", "中等", "5倍", "低", "1-2年", "紧缺", "必投合规", "合规保障", "安全体系", "P0", "10"],
                    ["多模态融合", "高", "广泛", "视觉+语音+视频", "高", "中高", "很热", "激烈", "3倍", "中", "1-2年", "紧缺", "垂直落地", "年回报30%", "行业深耕", "P0", "9.5"],
                    ["边缘AI", "中", "普及", "端侧推理+隐私计算", "中高", "中", "较热", "中等", "2倍", "低", "2-3年", "中", "逐步推进", "稳健增长", "边缘协同", "P1", "9"],
                    ["AI编程", "中高", "效率革命", "CodeLlama+GPT4o", "高", "中", "很热", "激烈", "5倍", "低", "1-2年", "紧缺", "效率优先", "6个月见效", "工具整合", "P0", "9.5"],
                    ["RAG 2.0", "高", "落地加速", "多向量+Rerank", "高", "中", "较热", "激烈", "2倍", "低", "1年", "中", "技术升级", "平稳增长", "知识库", "P1", "9"],
                    ["长上下文", "高", "范式转变", "百万Token+记忆", "极高", "中高", "很热", "激烈", "5倍", "中", "1-2年", "紧缺", "核心能力", "战略布局", "持续迭代", "P0", "10"],
                ]
                for i in range(rows):
                    for j in range(cols):
                        text = data[i][j] if i < len(data) and j < len(data[i]) else "---"
                        fill_cell(table.cell(i, j), text, 12 if i == 0 else 9, i == 0)

    # 幻灯片10 - 项目10: 市场预测2026-2035 (适配7行13列表格)
    if len(prs.slides) > 9:
        slide = prs.slides[9]
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                rows, cols = len(table.rows), len(table.columns)
                data = [
                    ["年份", "市场规模", "增长率", "主要驱动力", "投资规模", "技术突破", "应用场景", "竞争格局", "政策环境", "人才供给", "技术成熟度", "商业化程度", "投资建议"],
                    ["2026", "2000亿", "85%", "多模态+Agent", "600亿", "GPT-4o", "全场景", "格局初定", "政策完善", "紧缺", "高", "高", "重仓"],
                    ["2027", "3800亿", "90%", "具身智能", "1200亿", "AGI雏形", "智能硬件", "巨头主导", "监管落地", "紧缺", "高", "极高", "长期"],
                    ["2028", "6500亿", "71%", "边缘AI", "2000亿", "通用Agent", "万物互联", "生态竞争", "全球标准", "平衡", "极高", "极高", "稳健"],
                    ["2029", "10000亿", "54%", "脑机接口", "3000亿", "神经连接", "人机融合", "寡头垄断", "伦理规范", "充足", "极高", "极高", "价值"],
                    ["2030", "15000亿", "50%", "AGI落地", "4500亿", "通用AI", "全面渗透", "全球格局", "国际治理", "充足", "极高", "极高", "战略"],
                    ["2031", "22000亿", "47%", "AI生态", "6000亿", "Agent网络", "智能经济", "生态竞争", "标准统一", "充足", "极高", "极高", "价值"],
                    ["2035", "85000亿", "37%", "AGI普及", "15000亿", "通用AGI", "智能时代", "全球一体", "AI治理", "充足", "极高", "极高", "长期价值"],
                ]
                for i in range(rows):
                    for j in range(cols):
                        text = data[i][j] if i < len(data) and j < len(data[i]) else "---"
                        fill_cell(table.cell(i, j), text, 12 if i == 0 else 9, i == 0)

    # 幻灯片11 - 项目11: 最新10大维度 (适配6行11列表格)
    if len(prs.slides) > 10:
        slide = prs.slides[10]
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                rows, cols = len(table.rows), len(table.columns)
                data = [
                    ["维度", "价值描述", "量化指标", "测试环境", "数据规模", "准确率", "响应时间", "资源消耗", "可扩展性", "安全性", "商业价值"],
                    ["推理能力", "复杂任务+逻辑推理", "96%", "GPT-4o", "500万token", "93%", "<150ms", "中", "高", "高", "极高"],
                    ["多模态理解", "图文音视频+3D", "92%", "Gemini Ultra", "50万样本", "90%", "<400ms", "高", "中", "中", "高"],
                    ["长上下文", "百万Token+记忆", "88%", "Claude 3", "1M+ token", "87%", "<800ms", "高", "中", "中", "极高"],
                    ["工具调用", "API+代码执行", "99%", "GPT-4", "10000次", "99.5%", "<80ms", "低", "高", "高", "高"],
                    ["自我反思", "错误纠正+优化", "85%", "GPT-4o", "10万样本", "82%", "<1.5s", "高", "中", "中", "极高"],
                    ["综合评分", "全面能力评估", "91%", "All Models", "综合测试", "90%", "<300ms", "中", "高", "高", "极高"],
                ]
                for i in range(rows):
                    for j in range(cols):
                        text = data[i][j] if i < len(data) and j < len(data[i]) else "---"
                        fill_cell(table.cell(i, j), text, 12 if i == 0 else 9, i == 0)

    # 幻灯片12 - 项目12: 总结+10年规划 (适配11行3列表格)
    if len(prs.slides) > 11:
        slide = prs.slides[11]
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                rows, cols = len(table.rows), len(table.columns)
                data = [
                    ["项目", "核心内容", "未来10年规划"],
                    ["1", "封面升级", "季度更新迭代"],
                    ["2", "目录优化", "AI智能导航"],
                    ["3", "协议架构", "A2A 3.0升级"],
                    ["4", "能力矩阵", "扩展至10层"],
                    ["5", "16大热点", "年度更新跟踪"],
                    ["6", "应用场景", "扩展至20领域"],
                    ["7", "案例分析", "扩展至20案例"],
                    ["8", "技术对比", "季度对比报告"],
                    ["9", "市场预测", "年度滚动预测"],
                    ["10", "10大维度", "扩展至15维度"],
                    ["合计", "完整AI方案", "分阶段落地"],
                ]
                for i in range(rows):
                    for j in range(cols):
                        text = data[i][j] if i < len(data) and j < len(data[i]) else "---"
                        fill_cell(table.cell(i, j), text, 12 if i == 0 else 9, i == 0)

def main():
    ppt_path = r"f:\个人作品\legal-rag-qa-system\A2A_PROTOCOL_AI_AGENT_V14_终极完整版.pptx"
    prs = Presentation(ppt_path)
    
    print("🔄 正在修复项目5、10、11、12...")
    fill_all_slides(prs)
    
    output_path = r"f:\个人作品\legal-rag-qa-system\A2A_PROTOCOL_AI_AGENT_V14_终极完整版.pptx"
    prs.save(output_path)
    print(f"✅ 已完成！保存到: {output_path}")

if __name__ == "__main__":
    main()