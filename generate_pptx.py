from pptx import Presentation
from pptx.util import Pt, Inches, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def set_cell_style(cell, text, font_size=8, bold=False, text_color=RGBColor(255,255,255), fill_color=RGBColor(30,30,60)):
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill_color
    cell.text_frame.clear()
    cell.text_frame.word_wrap = True
    
    text = str(text) if text else "-"
    
    if "\n" in text:
        lines = text.split("\n")
        for i, line in enumerate(lines):
            para = cell.text_frame.add_paragraph()
            run = para.add_run()
            run.text = line.strip()
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = text_color
            run.font.name = 'Microsoft YaHei'
            run.font.latin_name = 'Arial'
            para.alignment = PP_ALIGN.CENTER
            para.space_after = Pt(0)
            para.space_before = Pt(0)
    else:
        para = cell.text_frame.add_paragraph()
        run = para.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = text_color
        run.font.name = 'Microsoft YaHei'
        run.font.latin_name = 'Arial'
        para.alignment = PP_ALIGN.CENTER
        para.space_after = Pt(0)
        para.space_before = Pt(0)
    
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

def add_watermark(slide):
    watermark_text = "© 2026 A2A PROTOCOL - 仅供学习参考，禁止商用"
    watermark_box = slide.shapes.add_textbox(Inches(1.5), Inches(4), Inches(9), Inches(2))
    watermark_frame = watermark_box.text_frame
    para = watermark_frame.add_paragraph()
    run = para.add_run()
    run.text = watermark_text
    run.font.size = Pt(24)
    run.font.color.rgb = RGBColor(150, 150, 150)
    run.font.name = 'Microsoft YaHei'
    run.font.bold = True
    watermark_box.rotation = 30
    watermark_box.zorder = 1

def create_full_pptx(add_watermark_flag=True, output_filename="output.pptx"):
    prs = Presentation()
    prs.slide_width = Inches(11.69)
    prs.slide_height = Inches(6.56)
    
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(5, 10, 25)
    
    title = slide.shapes.title
    title.text = "A2A协议2.0：AI智能体演进路线图（2026-2035）"
    for para in title.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(22)
            run.font.bold = True
            run.font.color.rgb = RGBColor(100, 180, 255)
            run.font.name = 'Microsoft YaHei'
    
    subtitle = slide.placeholders[1]
    subtitle.text = "基于行业趋势推演的合理预测 | 模拟数据，仅供参考\n免责声明：所有预测内容均为模拟推演，不代表真实未来；技术术语引用自公开文献或合理推测"
    for para in subtitle.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(12)
            run.font.color.rgb = RGBColor(150, 200, 255)
            run.font.name = 'Microsoft YaHei'
    
    if add_watermark_flag:
        add_watermark(slide)

    projects = [
        {"title": "项目1：A2A协议核心架构演进（2026-2035）",
            "bg_color": RGBColor(15, 25, 45),
            "header_color": RGBColor(46, 90, 157),
            "accent_color": RGBColor(100, 180, 255),
            "data": [
                ["年份", "主题类型", "负责人/团队", "协议层级", "AI大模型", "核心技术栈", "性能指标", "技术难点", "应对措施", "年度总结/关键亮点"],
                ["2026上半年", "技术突破", "架构组", "应用层", "GPT-5.0/DeepSeek-V5", "LangChain/RAG/LlamaIndex", "<10ms, QPS>50万", "长上下文理解", "动态管理", "突破: RAG v3.0.0升级; 多模态API发布 | 多模态融合"],
                ["2026下半年", "市场影响", "产品组", "应用层", "Gemini-2.5/GPT-5.0", "Multi-Agent/AutoGPT/CrewAI", "效率+200%, >1万 task/s", "Agent协调", "智能调度", "突破: 企业级Agent平台; A轮融资 | 企业级部署"],
                ["2027上半年", "技术突破", "架构组", "消息层", "GPT-5.5/DeepSeek-V6", "Kafka/NATS/eBPF", ">100万 msg/s, 99.99%可用", "消息稳定性", "幂等性设计", "突破: 事件驱动架构; 技术白皮书发布 | 事件驱动"],
                ["2027下半年", "投资需求", "财务组", "消息层", "Claude-3.5/GPT-5.5", "Milvus/Redis/TiDB", "<10ms, 索引10亿+", "分布式一致", "HNSW优化", "突破: 向量数据库融合; B轮融资 | 向量数据库"],
                ["2028上半年", "技术突破", "架构组", "网络层", "GPT-6.0/Claude-3.5", "HTTP/3/QUIC/gRPC", "延迟降低40%, >20万 req/s", "多路复用", "拥塞控制", "突破: HTTP/3升级; 性能白皮书发布 | 多路复用"],
                ["2028下半年", "技术突破", "安全组", "网络层", "Gemini-2.5/DeepSeek-V6", "OPA/Vault/WAF", "威胁检测99.9%, 防护100%", "AI对抗", "统一策略", "突破: 安全架构升级; 三级认证 | 安全升级"],
                ["2029上半年", "技术突破", "数据组", "数据层", "GPT-6.5/DeepSeek-V7", "Milvus/Redis/TiDB", "精度99.9%, >50万 qps", "向量检索", "增量索引", "突破: 10亿级向量; 数据平台上线 | 向量升级"],
                ["2029下半年", "技术突破", "安全组", "数据层", "Gemini-3.0/Qianwen-4.0", "CRYSTALS-Kyber/TLS1.3", "性能+50%, NIST认证", "后量子性能", "硬件加速", "突破: 后量子试点; 量子白皮书发布 | 后量子"],
                ["2030上半年", "技术突破", "架构组", "安全层", "GPT-7.0/量子AI", "QKD/量子加密/Qiskit", "距离>1000km, 无条件安全", "量子分发", "量子中继", "突破: 量子密码部署; 安全联盟成立 | 量子安全"],
                ["2030下半年", "技术突破", "研发组", "安全层", "Claude-4.0/DeepSeek-V7", "量子-AI融合", "加速比>100x, 效率+500%", "量子融合", "混合计算", "突破: 量子-AI融合; 技术峰会 | 量子融合"],
                ["2035上半年", "技术突破", "前沿组", "量子层", "Quantum-AI/GPT-8.0", "QKD/量子加密/Qiskit", "覆盖率95%, >3000km", "量子纠错", "中继器部署", "突破: 量子通信网络试点; 量子中继器部署 | 量子通信试点"],
                ["2035下半年", "技术突破", "前沿组", "量子层", "Quantum-AI/GPT-8.0", "QKD/量子加密/Qiskit", "覆盖率100%, >5000km", "量子纠错", "中继器", "突破: 量子通信网络; 量子互联网试验 | 量子通信"]
            ]
        },
        {"title": "项目2：价值维度评估体系（2026-2035）",
            "bg_color": RGBColor(10, 45, 35),
            "header_color": RGBColor(46, 125, 50),
            "accent_color": RGBColor(50, 255, 200),
            "data": [
                ["年份", "主题类型", "负责人/团队", "评估维度", "AI大模型", "核心技术栈", "性能指标", "技术难点", "应对措施", "年度总结/关键亮点"],
                ["2026上半年", "技术突破", "架构组", "互操作性", "GPT-5.0/DeepSeek-V5", "Kong/LangChain", "集成率95%, 兼容100%", "系统集成", "标准化接口", "突破: API网关升级; 标准接口发布 | API标准化"],
                ["2026下半年", "市场影响", "产品组", "弹性扩展", "Gemini-2.5/GPT-5.0", "K8s/KEDA/HPA", "扩缩容-80%, 准确率99%", "资源调度", "智能调度", "突破: 弹性伸缩; 云原生白皮书发布 | 弹性伸缩"],
                ["2027上半年", "技术突破", "安全组", "安全合规", "GPT-5.5/Qianwen-3.6", "ZTNA/mTLS/OPA", "合规100%, <5ms认证", "策略管理", "统一策略", "突破: 零信任部署; 安全白皮书发布 | 零信任"],
                ["2027下半年", "技术突破", "运维组", "可观测性", "Claude-3.5/GPT-5.5", "Prometheus/Loki/Tempo", "MTTD<5min, 检测99%", "异常检测", "AI预测", "突破: 智能监控; 运维升级 | 智能监控"],
                ["2028上半年", "技术突破", "研发组", "智能协作", "GPT-6.0/Claude-3.5", "Multi-Agent/AutoGPT/CrewAI", "效率+200%, 完成率98%", "Agent协调", "智能调度", "突破: 协作平台; 效率提升 | Agent协作"],
                ["2028下半年", "技术突破", "架构组", "跨云部署", "Gemini-2.5/DeepSeek-V6", "Crossplane/Terraform/Karmada", "<300s, 统一管理100%", "多云管理", "统一编排", "突破: 多云统一; 多云白皮书发布 | 多云编排"],
                ["2029上半年", "技术突破", "前沿组", "量子-AI评估", "GPT-6.5/DeepSeek-V7", "量子评估框架", "准确率99.99%, 加速>50x", "量子精度", "量子算法", "突破: 量子评估; 评估框架 | 量子评估"],
                ["2029下半年", "技术突破", "安全组", "AI安全评估", "Gemini-3.0/Qianwen-4.0", "AI安全评估框架", "检测99.9%, 合规100%", "伦理评估", "检测引擎", "突破: AI安全体系; 评估标准 | AI安全评估"],
                ["2030上半年", "技术突破", "前沿组", "量子安全评估", "GPT-7.0/量子AI", "量子安全评估", "精度99.99%, 量子安全", "量子安全", "安全框架", "突破: 量子安全评估; 量子白皮书发布 | 量子安全评估"],
                ["2030下半年", "技术突破", "研发组", "综合评估", "Claude-4.0/DeepSeek-V7", "综合AI评估", "覆盖100%, 准确率99.9%", "多维度评估", "统一框架", "突破: 综合评估体系; 评估框架 | 综合评估"],
                ["2035上半年", "技术突破", "前沿组", "量子评估体系", "Quantum-AI/GPT-8.0", "量子评估/AI验证", "精度99.99%, 覆盖95%", "算法优化", "量子纠错试点", "突破: 量子评估体系试点; 评估标准制定 | 量子评估试点"],
                ["2035下半年", "技术突破", "前沿组", "量子评估体系", "Quantum-AI/GPT-8.0", "量子评估/AI验证", "精度99.999%, 覆盖100%", "算法优化", "量子纠错", "突破: 量子评估体系; 评估标准 | 量子评估"]
            ]
        },
        {"title": "项目3：安全组件体系（2026-2035）",
            "bg_color": RGBColor(15, 35, 25),
            "header_color": RGBColor(46, 125, 50),
            "accent_color": RGBColor(0, 200, 150),
            "data": [
                ["年份", "主题类型", "负责人/团队", "安全组件", "AI大模型", "核心技术栈", "安全能力", "技术难点", "应对措施", "年度总结/关键亮点"],
                ["2026上半年", "技术突破", "安全组", "API网关", "GPT-5.0/DeepSeek-V5", "Kong/APISIX/Envoy", "防护率100%", "API安全", "AI分析", "突破: API网关安全升级; 安全架构发布 | AI流量调度"],
                ["2026下半年", "技术突破", "安全组", "零信任架构", "Gemini-2.5/GPT-5.0", "BeyondCorp/ZTNA/SASE", "防护等级领先", "网络隔离", "微分段", "突破: 网络隔离; 零信任白皮书发布 | 零信任"],
                ["2027上半年", "技术突破", "安全组", "数据加密", "GPT-5.5/DeepSeek-V6", "AES-256/TLS1.3/Kyber", "NIST认证", "后量子性能", "硬件加速", "突破: 后量子密码; 加密升级 | 后量子加密"],
                ["2027下半年", "伦理挑战", "合规组", "隐私计算", "Claude-3.5/GPT-5.5", "MPC/FHE/联邦学习", "PII保护", "计算效率", "MPC优化", "突破: 隐私计算平台; 隐私白皮书发布 | 隐私保护"],
                ["2028上半年", "技术突破", "安全组", "AI安全", "GPT-6.0/Claude-3.5", "AI安全检测/防护", "检测率99.99%", "后门检测", "水印溯源", "突破: 模型防护; AI安全标准 | AI模型安全"],
                ["2028下半年", "技术突破", "安全组", "云原生安全", "Gemini-2.5/DeepSeek-V6", "Trivy/Falco/eBPF", "漏洞检测100%", "容器安全", "漏洞扫描", "突破: 云原生安全; 云安全白皮书发布 | 云原生安全"],
                ["2029上半年", "技术突破", "前沿组", "量子-AI安全", "GPT-6.5/DeepSeek-V7", "量子认证", "无条件安全", "量子协议", "安全协议", "突破: 量子认证; 量子白皮书发布 | 量子AI安全"],
                ["2029下半年", "技术突破", "安全组", "AI安全架构", "Gemini-3.0/Qianwen-4.0", "AI防护框架", "检测99.99%, 覆盖100%", "对抗攻击", "Guardrails", "突破: 安全架构; AI白皮书发布 | AI安全架构"],
                ["2030上半年", "技术突破", "前沿组", "量子安全", "GPT-7.0/量子AI", "量子协议/QKD", ">1000km, 覆盖100%", "量子分发", "量子中继", "突破: 量子部署; 量子框架 | 量子安全"],
                ["2030下半年", "技术突破", "安全组", "综合安全", "Claude-4.0/DeepSeek-V7", "综合防护/零信任", "覆盖100%, 等级领先", "威胁检测", "统一策略", "突破: 综合安全体系; 安全白皮书发布 | 综合安全"],
                ["2035上半年", "技术突破", "前沿组", "量子安全体系", "Quantum-Security/GPT-8.0", "量子协议/QKD", ">3000km, 覆盖95%", "密钥管理", "密钥轮换试点", "突破: 量子安全体系试点; 量子标准制定 | 量子安全试点"],
                ["2035下半年", "技术突破", "前沿组", "量子安全体系", "Quantum-Security/GPT-8.0", "量子协议/QKD", ">5000km, 覆盖100%", "密钥管理", "密钥轮换", "突破: 量子安全体系; 量子标准 | 量子安全"]
            ]
        },
        {"title": "项目4：部署策略规划（2026-2035）",
            "bg_color": RGBColor(45, 30, 15),
            "header_color": RGBColor(46, 90, 157),
            "accent_color": RGBColor(255, 150, 50),
            "data": [
                ["年份", "主题类型", "负责人/团队", "部署策略", "AI大模型", "核心技术栈", "性能指标", "技术难点", "应对措施", "年度总结/关键亮点"],
                ["2026上半年", "技术突破", "运维组", "云原生部署", "GPT-5.0/DeepSeek-V5", "K8s/Istio/ArgoCD", "<30s部署, <10s扩缩容", "GitOps", "自动化部署", "突破: 云原生升级; GitOps落地 | GitOps"],
                ["2026下半年", "市场影响", "产品组", "Serverless", "Gemini-2.5/GPT-5.0", "Lambda/Knative/OpenFaaS", "<60s启动, >100万/s", "冷启动", "预热策略", "突破: 自动扩缩容; Serverless平台 | 无服务器"],
                ["2027上半年", "技术突破", "研发组", "AI推理", "GPT-5.5/DeepSeek-V6", "vLLM/TensorRT-LLM/TGI", "<5ms推理, >500 tokens/s", "推理延迟", "动态批处理", "突破: 吞吐量+300%; 推理升级 | 高效推理"],
                ["2027下半年", "伦理挑战", "合规组", "AGI部署", "Claude-3.5/GPT-5.5", "AGI框架/分布式推理", "<50ms响应, 99.5%准确", "AGI部署", "容错机制", "突破: AGI试点; AGI部署 | AGI落地"],
                ["2028上半年", "技术突破", "运维组", "边缘AI", "GPT-6.0/Claude-3.5", "EdgeAI/TensorRT/ONNX", "<5ms延迟, >100万 req/s", "资源受限", "模型压缩", "突破: 边缘优化; 边缘AI平台 | 边缘AI加速"],
                ["2028下半年", "技术突破", "架构组", "多云优化", "Gemini-2.5/DeepSeek-V6", "Crossplane/Terraform/Karmada", "<5min跨云, 统一100%", "跨云延迟", "状态同步", "突破: 多云统一; 多云白皮书发布 | 多云优化"],
                ["2029上半年", "技术突破", "安全组", "量子加密", "GPT-6.5/DeepSeek-V7", "QKD/量子加密/BB84", ">500km, 无条件安全", "退相干", "量子中继", "突破: 量子试点; 量子启动 | 量子加密"],
                ["2029下半年", "技术突破", "安全组", "后量子密码", "Gemini-3.0/Qianwen-4.0", "CRYSTALS-Kyber/TLS1.3", "性能+50%, NIST认证", "部署复杂", "硬件加速", "突破: 后量子试点; 后量子白皮书发布 | 后量子"],
                ["2030上半年", "技术突破", "架构组", "量子融合", "GPT-7.0/量子AI", "量子-AI融合框架", "加速>100x, 效率+500%", "量子融合", "混合计算", "突破: 量子融合; 技术峰会 | 量子融合"],
                ["2030下半年", "技术突破", "前沿组", "量子通信", "Claude-4.0/DeepSeek-V7", "QKD/IBM Qiskit", ">1000km, 无条件安全", "量子网络", "卫星链路", "突破: 量子通信; 量子网络 | 量子通信"],
                ["2035上半年", "技术突破", "前沿组", "量子网络", "Quantum-AI/GPT-8.0", "量子通信/量子中继", "覆盖95%, >3000km", "量子纠错", "纠错码试点", "突破: 量子互联网试点; 量子中继部署 | 量子网络试点"],
                ["2035下半年", "技术突破", "前沿组", "量子网络", "Quantum-AI/GPT-8.0", "量子通信/量子中继", "覆盖100%, >5000km", "量子纠错", "纠错码", "突破: 量子互联网; 量子发布 | 量子网络"]
            ]
        },
        {"title": "项目5：AI能力演进（2026-2035）",
            "bg_color": RGBColor(35, 15, 45),
            "header_color": RGBColor(90, 46, 157),
            "accent_color": RGBColor(200, 100, 255),
            "data": [
                ["年份", "主题类型", "负责人/团队", "AI能力", "AI大模型", "核心技术栈", "性能指标", "技术难点", "应对措施", "年度总结/关键亮点"],
                ["2026上半年", "技术突破", "研发组", "自然语言理解", "GPT-5.0/DeepSeek-V5", "Transformer/LoRA/RLHF", "准确率95%+, >10万 req/s", "长上下文理解", "动态管理", "突破: NLU升级; NLU平台 | NLU优化"],
                ["2026下半年", "市场影响", "产品组", "代码生成", "Gemini-2.5/GPT-5.0", "CodeLlama/Copilot", "效率+200%, 85%准确", "代码质量", "安全扫描", "突破: 代码效率; 代码助手 | 代码生成"],
                ["2027上半年", "技术突破", "研发组", "推理能力", "GPT-5.5/DeepSeek-V6", "思维链/工具使用", "准确率+30%, 80%通过", "复杂推理", "CoT", "突破: 推理升级; 推理引擎 | 推理升级"],
                ["2027下半年", "人才需求", "HR组", "AI工程化", "Claude-3.5/GPT-5.5", "MLOps/模型监控", "效率+100%, 覆盖100%", "模型部署", "自动化", "突破: AI工程体系; MLOps平台 | AI工程"],
                ["2028上半年", "技术突破", "研发组", "自主学习", "GPT-6.0/Claude-3.5", "AutoML/自监督", "效率+200%", "自主发现", "元学习", "突破: 自主学习; 自主框架 | 自主学习"],
                ["2028下半年", "技术突破", "研发组", "多Agent协作", "Gemini-2.5/DeepSeek-V6", "Multi-Agent/Swarm", "效率+200%, 完成98%", "Agent协调", "智能调度", "突破: 协作平台; 多Agent平台 | 多Agent"],
                ["2029上半年", "技术突破", "前沿组", "量子AI融合", "GPT-6.5/DeepSeek-V7", "量子神经网络", "加速>100x, 效率+500%", "量子纠错", "错误纠正", "突破: 量子加速; 量子AI白皮书发布 | 量子AI"],
                ["2029下半年", "技术突破", "研发组", "高级推理", "Gemini-3.0/Qianwen-4.0", "逻辑/数学/符号推理", "通过90%+, 数学85%+", "复杂推理", "程序合成", "突破: 推理90%; 高级引擎 | 高级推理"],
                ["2030上半年", "技术突破", "研发组", "多模态融合", "GPT-7.0/量子AI", "多模态Transformer", "理解95%+", "跨模态推理", "融合训练", "突破: 深度融合; 多模态白皮书发布 | 多模态融合"],
                ["2030下半年", "技术突破", "安全组", "AI安全", "Claude-4.0/DeepSeek-V7", "AI安全检测/防护", "检测99.9%", "对抗攻击", "Guardrails", "突破: AI安全体系; 安全框架 | AI安全"],
                ["2035上半年", "技术突破", "前沿组", "量子安全AI", "Quantum-AI/GPT-8.0", "量子安全/AI融合", "推理效率提升50%", "量子融合试点", "安全框架试点", "突破: 量子安全AI试点; 量子框架发布 | 量子安全AI试点"],
                ["2035下半年", "技术突破", "前沿组", "量子安全AI", "Quantum-AI/GPT-8.0", "量子安全/AI融合", "推理效率持续提升", "量子融合", "安全框架", "突破: 量子安全AI; 量子发布 | 量子安全AI"]
            ]
        },
        {"title": "项目6：性能指标体系（2026-2035）",
            "bg_color": RGBColor(25, 45, 15),
            "header_color": RGBColor(125, 157, 46),
            "accent_color": RGBColor(180, 255, 100),
            "data": [
                ["年份", "主题类型", "负责人/团队", "指标维度", "AI大模型", "核心技术栈", "性能指标", "技术难点", "应对措施", "年度总结/关键亮点"],
                ["2026上半年", "技术突破", "研发组", "延迟优化", "GPT-5.0/DeepSeek-V5", "vLLM/TensorRT-LLM", "<100ms, >1000 tokens/s", "推理延迟", "KV缓存", "突破: 延迟-50%; 低延迟平台 | 低延迟推理"],
                ["2026下半年", "市场影响", "运维组", "可用性", "Gemini-2.5/GPT-5.0", "K8s/ArgoCD/Istio", "99.99%可用, MTTR<1800s", "故障恢复", "故障转移", "突破: 99.99%可用; 高可用架构 | 高可用"],
                ["2027上半年", "技术突破", "研发组", "准确率", "GPT-5.5/DeepSeek-V6", "AI评估/监控", "98%+准确, 95%+召回", "模型漂移", "自动微调", "突破: 准确率98%; 评估升级 | 高精度"],
                ["2027下半年", "伦理挑战", "合规组", "可解释性", "Claude-3.5/GPT-5.5", "可解释AI/可视化", "95%+可解释, 透明达标", "决策解释", "LIME/SHAP", "突破: 可解释体系; 可解释框架 | 可解释AI"],
                ["2028上半年", "技术突破", "研发组", "量子性能", "GPT-6.0/Claude-3.5", "量子计算/量子神经网络", "加速>100x, 效率+500%", "量子纠错", "错误纠正", "突破: 量子加速; 量子白皮书发布 | 量子性能"],
                ["2028下半年", "技术突破", "架构组", "扩展能力", "Gemini-2.5/DeepSeek-V6", "多云/边缘/容器", "扩展+200%, >100万 req/s", "分布式一致", "自动扩缩容", "突破: 弹性扩展; 扩展架构 | 弹性扩展"],
                ["2029上半年", "技术突破", "安全组", "后量子性能", "GPT-6.5/DeepSeek-V7", "后量子密码/量子加密", "性能+50%, NIST认证", "部署复杂", "硬件加速", "突破: 后量子部署; 后量子白皮书发布 | 后量子性能"],
                ["2029下半年", "技术突破", "前沿组", "量子-AI性能", "Gemini-3.0/Qianwen-4.0", "量子-AI融合", "加速>100x, 效率持续提升", "量子融合", "错误纠正", "突破: 量子-AI融合; 量子性能 | 量子AI性能"],
                ["2030上半年", "技术突破", "前沿组", "量子通信", "GPT-7.0/量子AI", "量子通信/量子加密", ">1000km, 无条件安全", "量子分发", "量子中继", "突破: 量子通信; 量子发布 | 量子通信"],
                ["2030下半年", "技术突破", "研发组", "综合性能", "Claude-4.0/DeepSeek-V7", "综合优化/AI调优", "性能+300%, 能效+200%", "能效平衡", "能效管理", "突破: 综合性能; 综合白皮书发布 | 综合优化"],
                ["2035上半年", "技术突破", "前沿组", "量子安全性能", "Quantum-AI/GPT-8.0", "量子安全/AI融合", "扩展能力提升50%", "性能平衡试点", "优化框架试点", "突破: 量子安全性能试点; 量子优化发布 | 量子安全试点"],
                ["2035下半年", "技术突破", "前沿组", "量子安全性能", "Quantum-AI/GPT-8.0", "量子安全/AI融合", "无限扩展能力", "性能平衡", "优化框架", "突破: 量子安全性能; 量子发布 | 量子安全"]
            ]
        },
        {"title": "项目7：生态系统建设（2026-2035）",
            "bg_color": RGBColor(15, 25, 35),
            "header_color": RGBColor(46, 125, 157),
            "accent_color": RGBColor(100, 200, 255),
            "data": [
                ["年份", "主题类型", "负责人/团队", "生态组件", "AI大模型", "合作伙伴", "生态规模", "技术难点", "应对措施", "年度总结/关键亮点"],
                ["2026上半年", "市场影响", "市场组", "开发者社区", "GPT-5.0/DeepSeek-V5", "GitHub/GitLab", "开发者10万+, 贡献者1万+", "社区运营", "开源治理", "突破: 开发者社区; 开源发布 | 开发者生态"],
                ["2026下半年", "投资需求", "商务组", "合作伙伴", "Gemini-2.5/GPT-5.0", "云厂商/ISV/SI", "伙伴100+, 覆盖20+行业", "伙伴管理", "生态合作", "突破: 合作伙伴体系; 合作伙伴大会 | 合作伙伴"],
                ["2027上半年", "市场影响", "HR组", "培训教育", "GPT-5.5/DeepSeek-V6", "在线教育/认证", "培训10万+, 认证5万+", "教育内容", "课程开发", "突破: 培训体系; 培训平台 | 教育培训"],
                ["2027下半年", "人才需求", "HR组", "人才培养", "Claude-3.5/GPT-5.5", "高校合作/人才平台", "培养1万+, 引进100+", "人才培养", "校企合作", "突破: 人才体系; 校企签约 | 人才生态"],
                ["2028上半年", "市场影响", "市场组", "全球扩张", "GPT-6.0/Claude-3.5", "全球化/本地化", "覆盖50+国家", "全球化运营", "本地化团队", "突破: 全球布局; 全球发布会 | 全球生态"],
                ["2028下半年", "市场影响", "商务组", "产业联盟", "Gemini-2.5/DeepSeek-V6", "行业联盟/标准组织", "成员1000+, 标准50+", "联盟治理", "标准推动", "突破: 产业联盟; 联盟大会 | 产业联盟"],
                ["2029上半年", "市场影响", "前沿组", "量子生态", "GPT-6.5/DeepSeek-V7", "量子计算/量子通信", "量子节点100+", "量子网络", "节点部署", "突破: 量子生态; 量子启动 | 量子生态"],
                ["2029下半年", "市场影响", "安全组", "安全生态", "Gemini-3.0/Qianwen-4.0", "安全平台/威胁检测", "覆盖100%, 检测99.9%", "威胁检测", "AI检测", "突破: 安全生态; 安全发布 | 安全生态"],
                ["2030上半年", "市场影响", "研发组", "AI安全生态", "GPT-7.0/量子AI", "AI安全检测/防护", "检测99.9%, 防护领先", "对抗攻击", "Guardrails", "突破: AI安全生态; AI安全白皮书发布 | AI安全生态"],
                ["2030下半年", "市场影响", "前沿组", "量子安全生态", "Claude-4.0/DeepSeek-V7", "量子安全/量子加密", "覆盖100%, >1000km", "量子通信", "安全框架", "突破: 量子安全生态; 量子大会 | 量子安全生态"],
                ["2035上半年", "市场影响", "前沿组", "全球量子生态", "Quantum-AI/GPT-8.0", "全球量子网络", "节点800+, 全球覆盖95%", "全球部署试点", "全球架构试点", "突破: 全球量子生态试点; 量子互联网试点 | 全球量子试点"],
                ["2035下半年", "市场影响", "前沿组", "全球量子生态", "Quantum-AI/GPT-8.0", "全球量子网络", "节点1000+, 全球覆盖", "全球部署", "全球架构", "突破: 全球量子生态; 量子互联网 | 全球量子生态"]
            ]
        },
        {"title": "项目8：技术方案对比（2026-2035）",
            "bg_color": RGBColor(30, 30, 60),
            "header_color": RGBColor(60, 60, 120),
            "accent_color": RGBColor(150, 150, 255),
            "data": [
                ["年份", "主题类型", "负责人/团队", "技术方案", "AI大模型", "核心技术栈", "性能对比", "技术难点", "应对措施", "年度总结/关键亮点"],
                ["2026上半年", "技术突破", "架构组", "A2A协议", "GPT-5.0/DeepSeek-V5", "A2A协议/gRPC", "<10ms, >100万 req/s", "协议设计", "标准化", "突破: A2A协议; 协议白皮书发布 | A2A协议"],
                ["2026下半年", "技术突破", "数据组", "向量数据库", "Gemini-2.5/GPT-5.0", "Milvus/FAISS/Pinecone", "<10ms, 索引10亿+", "向量检索", "HNSW优化", "突破: 向量DB; 向量平台 | 向量DB"],
                ["2027上半年", "技术突破", "运维组", "Serverless", "GPT-5.5/DeepSeek-V6", "Lambda/Knative", "<60s启动, 自动扩缩容", "冷启动", "预热策略", "突破: Serverless; Serverless平台 | Serverless"],
                ["2027下半年", "技术突破", "前沿组", "量子加密", "Claude-3.5/GPT-5.5", "QKD/量子加密/BB84", "无条件安全", "量子分发", "量子中继", "突破: 量子加密; 量子白皮书发布 | 量子加密"],
                ["2028上半年", "技术突破", "架构组", "HTTP/3", "GPT-6.0/Claude-3.5", "HTTP/3/QUIC/Envoy", "<15ms, 多路复用", "连接迁移", "拥塞控制", "突破: HTTP/3; 性能白皮书发布 | HTTP/3"],
                ["2028下半年", "技术突破", "运维组", "跨云部署", "Gemini-2.5/DeepSeek-V6", "Crossplane/Terraform/Karmada", "<5min, 多云统一", "多云管理", "统一编排", "突破: 跨云部署; 多云白皮书发布 | 跨云"],
                ["2029上半年", "技术突破", "前沿组", "量子-AI融合", "GPT-6.5/DeepSeek-V7", "量子-AI融合", "加速>100x, 效率+500%", "量子纠错", "错误纠正", "突破: 量子融合; 量子白皮书发布 | 量子融合"],
                ["2029下半年", "技术突破", "安全组", "后量子密码", "Gemini-3.0/Qianwen-4.0", "CRYSTALS-Kyber/TLS1.3", "性能+50%, NIST认证", "部署复杂", "硬件加速", "突破: 后量子; 后量子白皮书发布 | 后量子"],
                ["2030上半年", "技术突破", "前沿组", "量子通信", "GPT-7.0/量子AI", "QKD/IBM Qiskit", ">1000km, 无条件安全", "量子分发", "量子中继", "突破: 量子通信; 量子白皮书发布 | 量子通信"],
                ["2030下半年", "技术突破", "架构组", "量子网络", "Claude-4.0/DeepSeek-V7", "量子网络/量子中继", "覆盖50%, >3000km", "量子中继", "中继部署", "突破: 量子网络; 量子发布 | 量子网络"],
                ["2035上半年", "技术突破", "前沿组", "全球量子网络", "Quantum-AI/GPT-8.0", "全球量子网络", "覆盖95%, >4000km", "全球部署试点", "全球架构试点", "突破: 全球量子试点; 量子互联网试点 | 全球量子试点"],
                ["2035下半年", "技术突破", "前沿组", "全球量子网络", "Quantum-AI/GPT-8.0", "全球量子网络", "覆盖100%, >5000km", "全球部署", "全球架构", "突破: 全球量子; 量子互联网试验 | 全球量子"]
            ]
        }
    ]

    for project in projects:
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = project["bg_color"]
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(10.7), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.clear()
        para = title_frame.add_paragraph()
        run = para.add_run()
        run.text = project["title"]
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.color.rgb = project["accent_color"]
        run.font.name = 'Microsoft YaHei'
        
        note_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(10.7), Inches(0.4))
        note_frame = note_box.text_frame
        note_frame.clear()
        p = note_frame.paragraphs[0]
        p.text = "注：2026-2027年半年总结基于原季度规划整合；2028-2030年为直接预测；2031-2034年技术平稳积累，合并展示。"
        p.runs[0].font.size = Pt(8)
        p.runs[0].font.color.rgb = RGBColor(200, 200, 220)
        p.runs[0].font.name = 'Microsoft YaHei'
        p.alignment = PP_ALIGN.LEFT
        
        rows = len(project["data"])
        cols = len(project["data"][0])
        
        table = slide.shapes.add_table(rows, cols, Inches(0.3), Inches(1.25), Inches(11.1), Inches(4.5)).table
        
        col_widths = [
            Cm(1.6),   # 年份 6%
            Cm(1.8),   # 主题类型 7%
            Cm(1.8),   # 负责人/团队 7%
            Cm(1.8),   # 协议层级/部署策略等 7%
            Cm(2.8),   # AI大模型 11% - 扩展
            Cm(3.5),   # 核心技术栈 14%
            Cm(2.8),   # 性能指标/生态规模等 11%
            Cm(2.4),   # 技术难点 9%
            Cm(2.4),   # 应对措施 9%
            Cm(3.8)    # 年度总结/关键亮点 15% - 合并列扩展
        ]
        
        for j in range(min(cols, len(col_widths))):
            table.columns[j].width = col_widths[j]
        
        for i, row in enumerate(project["data"]):
            for j, cell in enumerate(row):
                if i == 0:
                    set_cell_style(table.cell(i, j), cell, font_size=7, bold=True, text_color=RGBColor(255,255,255), fill_color=project["header_color"])
                else:
                    set_cell_style(table.cell(i, j), cell, font_size=6, bold=False, text_color=RGBColor(200,200,220), fill_color=RGBColor(20,30,50) if i % 2 == 0 else RGBColor(25,35,55))
        
        if add_watermark_flag:
            add_watermark(slide)

    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(5, 10, 25)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(10.7), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.clear()
    para = title_frame.add_paragraph()
    run = para.add_run()
    run.text = "免责声明与数据来源说明"
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = RGBColor(100, 180, 255)
    run.font.name = 'Microsoft YaHei'
    
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(10.2), Inches(4.5))
    content_frame = content_box.text_frame
    content_frame.clear()
    
    disclaimer_text = """【免责声明】
1. 本PPT所有内容均为基于行业趋势推演的合理预测，不代表真实未来；
2. 技术术语引用自公开文献或合理推测，实际发展可能有所不同；
3. 本文件仅供学习参考，不构成任何投资建议或商业决策依据；
4. 所有预测数据均为模拟推演，实际结果可能与预测存在差异；
5. 本文件内容受著作权保护，未经授权禁止商用传播；
6. 技术预测截至2035年，更长远的未来（2040-2050）因存在技术奇点不确定性，暂未展开；
7. 2031-2034年属于技术平稳积累期，无重大突破，故合并展示。

【数据来源说明】
- AI模型信息参考自OpenAI、Google、Anthropic、DeepSeek、百度等厂商公开信息；
- 技术趋势参考Gartner技术趋势报告、IDC全球AI支出指南、工信部人工智能发展规划；
- 量子计算信息参考IBM、Google、中科院量子信息与量子科技创新研究院公开进展；
- 行业标准参考ISO/IEC、W3C、IEEE等国际标准化组织公开文档；
- 市场预测参考麦肯锡全球AI报告、德勤科技趋势报告、信通院人工智能白皮书；
- 所有时间节点为合理预测，实际落地时间可能因技术突破或阻碍而变化；
- 所有模拟数据均不构成对实际技术发展路径的承诺。

【风险提示】
- 技术发展存在不确定性，实际进度可能因技术突破或阻碍而变化；
- 监管政策可能对技术发展方向产生影响；
- 市场竞争格局可能发生变化；
- 人才供应可能影响技术落地进度。

© 2026 A2A PROTOCOL - 仅供学习参考，禁止商用"""
    
    para = content_frame.add_paragraph()
    run = para.add_run()
    run.text = disclaimer_text
    run.font.size = Pt(10)
    run.font.color.rgb = RGBColor(180, 200, 230)
    run.font.name = 'Microsoft YaHei'
    para.alignment = PP_ALIGN.LEFT
    
    if add_watermark_flag:
        add_watermark(slide)
    
    prs.save(output_filename)
    print(f"PPT文件已生成: {output_filename}")

if __name__ == "__main__":
    create_full_pptx(add_watermark_flag=True, output_filename="A2A协议2.0_AI智能体演进路线图_2026_V17_带水印.pptx")
    create_full_pptx(add_watermark_flag=False, output_filename="A2A协议2.0_AI智能体演进路线图_2026_V17_无水印.pptx")