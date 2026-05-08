from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import re

CHAPTER_COLORS = {
    1: {'bg': RGBColor(10, 15, 28), 'accent': RGBColor(0, 180, 255), 'text': RGBColor(230, 240, 250), 'highlight': RGBColor(255, 80, 80), 'secondary': RGBColor(50, 120, 180)},
    2: {'bg': RGBColor(12, 25, 22), 'accent': RGBColor(0, 220, 180), 'text': RGBColor(220, 240, 230), 'highlight': RGBColor(255, 100, 100), 'secondary': RGBColor(30, 150, 130)},
    3: {'bg': RGBColor(25, 15, 28), 'accent': RGBColor(220, 100, 200), 'text': RGBColor(240, 220, 245), 'highlight': RGBColor(255, 100, 100), 'secondary': RGBColor(150, 60, 150)},
    4: {'bg': RGBColor(25, 22, 12), 'accent': RGBColor(255, 200, 60), 'text': RGBColor(250, 245, 220), 'highlight': RGBColor(255, 100, 100), 'secondary': RGBColor(180, 140, 40)},
    5: {'bg': RGBColor(10, 25, 32), 'accent': RGBColor(160, 120, 255), 'text': RGBColor(220, 235, 250), 'highlight': RGBColor(255, 100, 100), 'secondary': RGBColor(100, 80, 200)},
    6: {'bg': RGBColor(28, 20, 15), 'accent': RGBColor(255, 120, 160), 'text': RGBColor(250, 235, 220), 'highlight': RGBColor(255, 100, 100), 'secondary': RGBColor(180, 80, 110)},
    7: {'bg': RGBColor(15, 30, 25), 'accent': RGBColor(0, 200, 150), 'text': RGBColor(220, 245, 230), 'highlight': RGBColor(255, 100, 100), 'secondary': RGBColor(50, 150, 100)},
    8: {'bg': RGBColor(30, 18, 35), 'accent': RGBColor(200, 150, 255), 'text': RGBColor(240, 230, 250), 'highlight': RGBColor(255, 100, 100), 'secondary': RGBColor(120, 80, 180)},
    9: {'bg': RGBColor(20, 25, 15), 'accent': RGBColor(180, 220, 80), 'text': RGBColor(235, 245, 220), 'highlight': RGBColor(255, 100, 100), 'secondary': RGBColor(120, 150, 60)},
    10: {'bg': RGBColor(35, 20, 40), 'accent': RGBColor(255, 150, 200), 'text': RGBColor(245, 230, 245), 'highlight': RGBColor(255, 100, 100), 'secondary': RGBColor(150, 100, 180)},
    11: {'bg': RGBColor(25, 20, 15), 'accent': RGBColor(139, 90, 43), 'text': RGBColor(245, 235, 220), 'highlight': RGBColor(139, 69, 19), 'secondary': RGBColor(180, 120, 60)}
}

CHAPTER_TITLES = [
    "第一章：A2A协议架构与核心价值",
    "第二章：协议规范与技术细节",
    "第三章：安全与合规体系",
    "第四章：部署与运维实践",
    "第五章：性能优化策略",
    "第六章：实战案例与最佳实践",
    "第七章：AI智能体发展趋势展望",
    "第八章：蚌埠市人才招聘规划",
    "第九章：创业指南与政策支持",
    "第十章：AI大模型发展趋势与软件生态",
    "第十一章：核心内容总结"
]

WATERMARK_TEXT = "© 何晓冬 | A2A协议深度解析"
COPYRIGHT_INFO = "Copyright © 2026 何晓冬. All Rights Reserved."

DEFAULT_TABLES = {
    1: [
        ['价值维度', '业务收益', '技术实现', '成功指标', '适用场景', '成熟度', '投资回报', '技术架构', '关键挑战'],
        ['互操作性', '降低集成成本🔴60%+', '标准化消息格式+Agent Cards', '集成时间缩短80%', '跨平台协作、生态整合', '高', '18个月', 'A2A协议栈', '协议版本兼容'],
        ['弹性扩展', '支持🔴百万级并发', '去中心化架构+K8s调度', '吞吐量提升10x', '大规模部署、弹性业务', '高', '12个月', 'K8s+Istio', '资源调度优化'],
        ['安全合规', '通过🔴SOC2/ISO27001认证', '端到端加密+零信任架构', '合规成本降低50%', '金融级应用、医疗数据', '高', '24个月', '零信任架构', '证书管理'],
        ['可观测性', '全链路追踪+实时监控', 'OpenTelemetry+Prometheus', 'MTTR降低60%', '运维监控、故障排查', '中', '18个月', 'O11y Stack', '监控告警风暴'],
        ['智能协作', 'Agent自主协同+技能发现', 'A2A技能发现协议', '任务完成率提升40%', '多Agent系统、任务分工', '中', '24个月', '多Agent架构', '协调一致性'],
        ['跨云部署', '多云无缝迁移', 'K8s原生支持+CNCF认证', '迁移时间缩短70%', '混合云架构、多云战略', '中', '36个月', '混合云架构', '数据同步'],
        ['成本优化', '降低运维成本🔴35%', '自动化运维+Serverless', 'TCO降低35%', '中大型企业、成本敏感', '中', '24个月', 'Serverless架构', '冷启动延迟'],
        ['快速迭代', '缩短开发周期🔴45%', '模块化设计+API优先', '上市时间缩短45%', '敏捷团队、快速上市', '高', '12个月', 'API优先架构', '版本兼容'],
        ['生态扩展', '第三方Agent接入', '标准化API+开发者SDK', '生态伙伴增长3x', '平台生态、合作伙伴', '中', '30个月', '开放平台架构', '生态治理'],
        ['数据价值', '数据资产变现', '隐私计算+联邦学习', '数据价值提升2x', '数据合作、合规共享', '低', '48个月', '隐私计算架构', '合规边界'],
        ['边缘计算', '边缘延迟🔴降低70%', 'K3s+边缘节点', '响应时间<20ms', '边缘场景、低延迟', '中', '24个月', '边缘架构', '边缘管理'],
        ['AI协同', 'AI任务效率🔴提升3x', 'Agent技能组合+编排', '复杂任务完成率+50%', 'AI应用、智能系统', '中', '30个月', 'AI Agent架构', '技能组合'],
        ['Agent自主', '自主决策能力🔴提升5x', '强化学习+决策引擎', '人工干预减少80%', '自主系统、智能决策', '中', '24个月', '自主Agent架构', '决策准确性'],
        ['多模态交互', '交互体验🔴提升4x', '文本+语音+视觉融合', '用户满意度提升60%', '智能客服、虚拟助手', '高', '18个月', '多模态Agent', '模态融合']
    ],
    2: [
        ['协议层级', '业务价值', '技术亮点', '成功指标', '适用场景', '性能指标', '投资回报', '技术架构', '关键挑战'],
        ['应用层', '技能发现效率🔴提升5x', 'Agent Cards/ADR规范', '技能发现<10ms', 'Agent协作、技能共享', '<10ms响应', '12个月', 'REST+GraphQL', '技能描述标准化'],
        ['消息层', '消息吞吐量🔴提升10x', 'JSON-RPC 2.0/JSON-LD', '支持10万TPS', '异步通信、事件驱动', '10万TPS', '18个月', 'Kafka/Pulsar', '消息顺序保证'],
        ['传输层', '延迟🔴降低60%', 'HTTP/2、gRPC、QUIC', '端到端<50ms', '低延迟通信、实时交互', '<50ms延迟', '12个月', 'gRPC+QUIC', '协议协商'],
        ['安全层', '安全合规成本🔴降低40%', 'OAuth2、mTLS、JWT', '零信任架构', '身份认证、数据加密', '零信任', '24个月', 'PKI+OAuth2', '密钥轮换'],
        ['发现层', '服务发现效率🔴提升8x', '服务注册中心/ETCD', '实时更新<1s', '动态服务、微服务架构', '实时更新', '18个月', 'ETCD+Consul', '一致性保证'],
        ['事件层', '实时推送延迟🔴降低80%', 'WebSocket、SSE', '<100ms延迟', '实时通知、状态同步', '<100ms延迟', '12个月', 'WebSocket集群', '连接管理'],
        ['数据层', '序列化效率🔴提升2x', 'Protobuf/JSON/BSON', '压缩率60%', '数据交换、存储', '60%压缩', '18个月', 'Protobuf+Schema', '版本演进'],
        ['控制层', '路由灵活性🔴提升3x', 'gRPC/REST/gRPC-Web', '弹性路由', 'API网关、流量控制', '弹性路由', '12个月', 'Envoy/Kong', '配置复杂度'],
        ['编排层', '工作流效率🔴提升4x', 'Dapr/Knative', '复杂流程编排', '业务流程自动化', '复杂流程', '24个月', 'Dapr+Airflow', '状态管理'],
        ['语义层', '知识理解能力🔴提升3x', 'RDF/OWL', '语义推理', '知识图谱、AI理解', '语义理解', '36个月', '知识图谱引擎', '推理性能'],
        ['Agent通信', 'Agent间通信效率🔴提升3x', 'A2A消息格式+协议栈', '通信延迟<5ms', '多Agent协作、消息传递', '<5ms延迟', '12个月', 'A2A协议栈', '消息可靠性'],
        ['Agent发现', 'Agent发现效率🔴提升6x', '服务发现+Agent注册', '发现时间<1s', '动态Agent网络、自动发现', '<1s发现', '18个月', 'Agent注册中心', '注册一致性']
    ],
    3: [
        ['安全组件', '业务价值', '技术亮点', '成功指标', '适用场景', '防护级别', '投资回报', '技术架构', '关键挑战'],
        ['API网关', 'API安全成本🔴降低50%', 'Kong/APISIX/Envoy', '攻击拦截率99%', 'API安全、流量控制', '高', '18个月', 'API网关集群', '性能损耗'],
        ['WAF防火墙', 'Web攻击防护🔴提升95%', 'ModSecurity/Akamai', 'OWASP Top10防护', 'Web应用、API保护', '极高', '24个月', 'WAF+CDN', '规则误报'],
        ['OAuth 2.0', '认证效率🔴提升3x', '授权码/PKCE/JWT', '认证延迟<3ms', '身份管理、SSO', '高', '12个月', 'OAuth2+OIDC', '令牌管理'],
        ['mTLS', '通道安全🔴提升100%', '双向证书+证书轮换', '零中间人攻击', '敏感数据传输', '极高', '24个月', 'PKI+SPIFFE', '证书分发'],
        ['RBAC/ABAC', '权限管理效率🔴提升4x', '角色/属性权限', '细粒度控制', '企业级访问控制', '高', '18个月', 'IAM系统', '权限爆炸'],
        ['DDoS防护', 'DDoS攻击防护🔴提升99%', 'Cloudflare/Akamai', '流量清洗率99.9%', '大规模流量攻击', '极高', '24个月', '多层防护', '成本投入'],
        ['数据加密', '数据安全合规🔴100%', 'AES-256/GPG', '加密覆盖率100%', '数据存储、传输', '极高', '18个月', '加密服务', '性能影响'],
        ['审计日志', '合规审计效率🔴提升5x', 'ELK/Splunk/Datadog', '日志不可篡改', '合规审计、行为追踪', '高', '24个月', 'SIEM系统', '存储成本'],
        ['威胁检测', '威胁发现🔴快80%', 'SIEM/EDR', '攻击发现<5min', '高级威胁检测', '高', '30个月', 'AI检测引擎', '误报率'],
        ['密钥管理', '密钥安全🔴100%', 'HSM/Cloud KMS', '密钥零泄露', '密钥存储、轮换', '极高', '24个月', 'HSM集群', '可用性'],
        ['数据脱敏', '数据隐私保护🔴100%', '动态脱敏+静态脱敏', '敏感数据零泄露', '数据共享、合规', '高', '24个月', '脱敏引擎', '精度平衡'],
        ['访问审计', '合规审计能力🔴提升4x', '行为分析+操作记录', '审计覆盖率100%', '合规检查、安全审计', '高', '18个月', '审计系统', '日志存储'],
        ['Agent安全', 'Agent间通信安全🔴提升100%', 'A2A安全协议+加密', '通信零泄露', '多Agent协作、数据传输', '极高', '18个月', 'A2A安全栈', '协议兼容'],
        ['AI伦理', 'AI行为合规性🔴100%', '伦理框架+审计机制', '违规行为零发生', 'AI应用、智能决策', '高', '24个月', '伦理治理系统', '伦理标准']
    ],
    4: [
        ['部署策略', '业务价值', '技术亮点', '成功指标', '适用场景', '可用性', '投资回报', '技术架构', '关键挑战'],
        ['多活部署', '业务可用性🔴99.99%', 'K8s StatefulSet+DNS轮询', '故障转移<30s', '核心业务、高可用', '99.99%', '24个月', '多活K8s集群', '数据一致性'],
        ['自动扩缩容', '资源利用率🔴提升60%', 'HPA/VPA+KEDA', '弹性响应<2min', '波动流量、弹性业务', '弹性', '12个月', 'K8s自动伸缩', '冷启动'],
        ['健康检查', '故障自愈能力🔴提升80%', 'Liveness/Readiness/Startup', '故障恢复<1min', '微服务、高可用', '高', '6个月', 'K8s健康检查', '探针配置'],
        ['蓝绿发布', '发布风险🔴降低90%', 'Istio/Argo Rollouts', '零停机升级', '生产环境、频繁发布', '零停机', '18个月', 'Istio+Argo', '流量切分'],
        ['异地容灾', '灾难恢复能力🔴99.999%', '跨区域K8s集群+DRBD', 'RPO<5min,RTO<15min', '金融系统、核心数据', '99.999%', '36个月', '跨区集群', '数据同步'],
        ['边缘部署', '边缘延迟🔴降低70%', 'K3s/KubeEdge', '<20ms延迟', '边缘计算、就近访问', '低延迟', '24个月', '边缘K8s', '边缘管理'],
        ['金丝雀发布', '发布风险🔴降低80%', 'Argo Rollouts/Flagger', '渐进式发布', '核心服务、风险控制', '渐进式', '18个月', 'Argo+Prometheus', '指标评估'],
        ['Serverless', '运维成本🔴降低50%', 'KEDA+Knative', '按需扩缩', '非核心服务、突发流量', '按需', '12个月', 'KEDA+Knative', '冷启动延迟'],
        ['混合云部署', '云中立🔴100%', 'Anthos/Azure Arc', '多云统一管理', '企业级、多云战略', '多云', '30个月', '混合云管理', '跨云网络'],
        ['GitOps', '部署效率🔴提升3x', 'Argo CD/Flux', '声明式配置', 'CI/CD、版本控制', '高', '18个月', 'GitOps工具链', '配置漂移'],
        ['Agent部署', 'Agent部署效率🔴提升4x', 'Agent镜像+自动编排', '部署时间缩短70%', '多Agent系统、自动化部署', '高', '12个月', 'Agent编排平台', '版本兼容'],
        ['Agent监控', 'Agent监控能力🔴提升5x', '分布式追踪+指标采集', '故障发现<2min', 'Agent集群、智能监控', '极高', '18个月', 'O11y+Agent监控', '监控告警']
    ],
    5: [
        ['优化维度', '业务价值', '技术亮点', '成功指标', '适用场景', '预期收益', '投资回报', '技术架构', '关键挑战'],
        ['网络优化', '延迟🔴降低40%', 'HTTP/2、gRPC、QUIC', 'RTT降低40%', '低延迟通信', '延迟-40%', '6个月', 'gRPC+QUIC', '客户端兼容'],
        ['消息批量', '吞吐量🔴提升2x', '批量聚合+背压控制', 'QPS提升2x', '高吞吐量场景', 'QPS+2x', '12个月', 'Kafka+batching', '消息顺序'],
        ['缓存策略', '响应速度🔴快3x', 'Redis集群+多级缓存', 'P95延迟降低60%', '读多写少场景', '延迟-60%', '6个月', 'Redis+LocalCache', '数据一致性'],
        ['并发处理', '容量🔴提升5x', '异步队列+线程池', '并发能力+5x', '高并发场景', '并发+5x', '18个月', 'Vert.x/Netty', '线程安全'],
        ['序列化', '数据体积🔴减少60%', 'Protobuf/FlatBuffers', '序列化体积-60%', '数据传输、存储', '体积-60%', '12个月', 'Protobuf', '版本兼容'],
        ['CDN加速', '边缘延迟🔴降低70%', '边缘缓存+预热', 'P95延迟-70%', '静态资源、全球用户', '延迟-70%', '6个月', 'CDN+预热', '缓存失效'],
        ['数据库优化', '查询速度🔴快2x', '读写分离+分库分表', '查询延迟-50%', '大数据量场景', '延迟-50%', '12个月', 'Sharding+读写分离', '事务一致性'],
        ['代码优化', '性能🔴提升30%', '算法改进+JIT', 'CPU使用率-30%', '计算密集场景', '性能+30%', '24个月', 'Profiling+优化', '回归风险'],
        ['资源调度', '资源利用率🔴提升40%', 'K8s QoS+资源预留', '资源利用率+40%', '容器集群', '利用率+40%', '8个月', 'K8s调度优化', 'Pod驱逐'],
        ['连接池', '连接效率🔴提升30%', 'HikariCP/OkHttp', '连接复用率+30%', '高连接场景', '复用率+30%', '3个月', '连接池优化', '连接泄漏'],
        ['内存优化', '内存使用🔴降低30%', '对象池+GC调优', '内存占用-30%', '内存敏感场景', '内存-30%', '12个月', '内存管理', 'GC停顿'],
        ['异步IO', 'IO效率🔴提升2x', 'Reactor/CompletableFuture', '吞吐量+100%', 'IO密集场景', '吞吐量+100%', '18个月', '异步框架', '复杂度'],
        ['Agent优化', 'Agent性能🔴提升3x', 'Agent缓存+预加载', '响应时间-70%', '多Agent系统、智能协作', '高', '12个月', 'Agent性能优化', '缓存一致性'],
        ['AI推理优化', '推理速度🔴提升4x', '模型量化+TensorRT', '推理延迟-60%', 'AI应用、边缘部署', '高', '18个月', 'AI推理引擎', '精度损失']
    ],
    6: [
        ['案例场景', '业务价值', '技术亮点', '成功指标', '行业领域', '规模', '投资回报', '技术架构', '关键挑战'],
        ['供应链管理', '效率🔴提升40%', '多Agent协作+智能调度', '交付周期缩短30%', '制造业', '大型', '18个月', 'A2A+IoT', '数据孤岛'],
        ['金融风控', '风险🔴降低35%', '实时决策引擎+知识图谱', '准确率提升25%', '金融', '中型', '24个月', 'A2A+ML', '合规要求'],
        ['智能办公', '生产力🔴提升50%', '自动化流程+RPA集成', '成本节约40%', '企业服务', '中型', '12个月', 'A2A+低代码', '系统集成'],
        ['医疗诊断', '准确率🔴提升25%', '专业知识整合+推理引擎', '误诊率下降20%', '医疗', '大型', '36个月', 'A2A+NLP', '数据隐私'],
        ['智能客服', '满意度🔴提升30%', '多轮对话+情感分析', '响应时间缩短50%', '电商', '大型', '18个月', 'A2A+LLM', '上下文理解'],
        ['智能运维', '故障发现🔴快80%', '预测性维护+异常检测', 'MTTR降低60%', '运维', '中型', '24个月', 'A2A+监控', '告警风暴'],
        ['智能推荐', '转化率🔴提升35%', '个性化算法+实时特征', '用户留存提升20%', '互联网', '大型', '18个月', 'A2A+推荐系统', '冷启动'],
        ['智能教育', '学习效率🔴提升45%', '自适应学习+知识追踪', '成绩提升25%', '教育', '中型', '36个月', 'A2A+教育AI', '学习路径'],
        ['智能制造', '良率🔴提升15%', '数字孪生+质量检测', '成本降低20%', '制造业', '大型', '30个月', 'A2A+IoT+CV', '设备互联'],
        ['智慧城市', '能耗🔴降低25%', '多系统协同+预测优化', '效率提升30%', '政府', '大型', '48个月', 'A2A+城市平台', '数据整合'],
        ['Agent协作', '协作效率🔴提升5x', '多Agent协同+任务分配', '任务完成率+60%', '多Agent系统、智能协作', '大型', '18个月', 'A2A协作平台', '协调一致性'],
        ['Agent生态', '生态价值🔴提升4x', 'Agent平台+开发者社区', '生态规模+3x', '平台生态、合作伙伴', '大型', '30个月', '开放Agent平台', '生态治理'],
        ['Agent创新', '创新效率🔴提升3x', 'AI Agent+新技术融合', '创新周期缩短50%', 'AI创新、前沿探索', '中型', '24个月', 'Agent创新实验室', '技术风险']
    ],
    7: [
        ['时间阶段', '发展特征', '关键技术', '行业影响', '蚌埠机遇', '投资价值', '风险评估', '战略建议', '预期收益'],
        ['2026-2028', '标准化加速', 'Agent Cards、ADR规范', '行业标准形成', '承接长三角产业转移', '高', '中', '布局技术团队', '抢占先机'],
        ['2029-2033', '智能升级', '自主学习、情感理解', 'Agent能力飞跃', '制造业AI升级', '极高', '中', '研发投入', '技术领先'],
        ['2034-2036', '生态成熟', '去中心化网络', '全球Agent经济', '构建产业生态', '高', '低', '平台建设', '生态收益'],
        ['2037-2041', '通用智能', '通用Agent出现', '生产力革命', '全面智能化', '极高', '中', '战略转型', '指数增长'],
        ['2042-2046', '意识涌现', '类人智能可能', '社会变革', '新产业形态', '极高', '高', '前瞻布局', '颠覆创新'],
        ['当前趋势', '多模态融合', '文本+语音+视觉', '交互体验升级', '本地应用开发', '高', '低', '产品落地', '市场份额'],
        ['当前趋势', '边缘部署', 'K3s/KubeEdge', '低延迟服务', '边缘计算节点', '中', '低', '基础设施', '技术储备'],
        ['当前趋势', '安全增强', '零信任+隐私计算', '合规保障', '金融级安全', '高', '低', '安全体系', '信任背书'],
        ['产业趋势', '平台化', '云厂商Agent平台', '生态竞争', '对接云服务商', '中', '中', '合作共赢', '资源整合'],
        ['产业趋势', '垂直深耕', '行业专用Agent', '细分市场', '制造业AI应用', '高', '低', '行业聚焦', '差异化优势']
    ],
    8: [
        ['人群分类', '岗位定位', '招聘策略', '推荐软件', '薪资范围', '蚌埠特色', '政策支持', '预期成效', '关键挑战'],
        ['青年人才(22-28岁)', '初级开发、测试', '校招+实习转正', 'BOSS直聘、智联', '8-15K/月', '本地高校合作', '实习补贴', '人才储备', '经验不足'],
        ['中坚力量(29-38岁)', '技术骨干、管理', '社招+返乡人才', 'BOSS直聘、猎聘', '20-35K/月', '返乡计划', '安家补贴', '核心团队', '竞争激烈'],
        ['资深专家(39-45岁)', '架构师、总监', '猎头引进', '猎聘、脉脉', '60-150万/年', '高层次人才计划', '住房优惠', '技术突破', '成本较高'],
        ['高中学历', '数据标注、客服', '职业院校+线上', 'BOSS直聘、58同城', '4-10K/月', '本地职高合作', '技能补贴', '基础岗位', '技能提升'],
        ['精神残疾人', '数据处理、审核', '残联推荐', '官方平台、微信小程序', '同工同酬', '就业保障', '残保金优惠', '社会责任', '环境适配'],
        ['大专学历', '运维、技术支持', '本地招聘', '前程无忧、BOSS直聘', '6-12K/月', '职业学院合作', '就业补贴', '运维保障', '发展空间'],
        ['本科学历', '全栈开发、产品', '重点高校', '智联、拉勾网', '12-25K/月', '高校招聘', '人才政策', '技术核心', '留才挑战'],
        ['硕士学历', 'AI算法、架构', '顶尖高校', '猎聘、拉勾网', '20-35K/月', '人才引进', '科研补贴', '技术突破', '竞争激烈'],
        ['博士学历', '首席科学家', '高端引进', '猎聘、学术平台', '80-200万/年', '院士工作站', '科研经费', '行业引领', '稀缺性'],
        ['AI工程师', 'AI Agent开发', '高校+社招', 'BOSS直聘、拉勾', '15-30K/月', 'AI人才引进', '科研补贴', '技术核心', '竞争激烈'],
        ['Agent训练师', 'Agent训练优化', '专业培训+实战', '训练效率+50%', 'AI应用、智能系统', '中', '12个月', 'AI训练平台', '数据质量'],
        ['顾问导师(46岁+)', '技术顾问', '兼职返聘', '脉脉、领英', '顾问费', '退休专家', '荣誉职位', '知识传递', '时间投入']
    ],
    9: [
        ['办理环节', '所需材料', '办理流程', '注意事项', '蚌埠政策', '办理地点', '办理时间', '费用说明', '常见问题'],
        ['名称核准', '身份证、备选名称', '线上申请→审核', '准备3-5个名称', '全程电子化', '政务服务中心', '1-2天', '免费', '名称重复'],
        ['场地证明', '租赁合同/无偿使用证明', '提供材料→审核', '住宅需社区同意', '住宅改商用支持', '所在区县', '1天', '免费', '地址不符'],
        ['材料提交', '身份证、场地证明', '线上/窗口提交', '经营范围规范', '一窗通办', '政务服务中心', '1-2天', '免费', '材料不全'],
        ['领取执照', '受理通知书', '窗口领取/邮寄', '核对信息', '当场办结', '政务服务中心', '当场', '免费', '信息错误'],
        ['刻章备案', '营业执照', '指定机构刻章', '公章+财务章+法人章', '免费刻章', '公安局指定点', '1-2天', '免费', '章型选择'],
        ['税务登记', '营业执照、法人身份证', '电子税务局', '30日内完成', '线上办理', '税务局', '1天', '免费', '逾期罚款'],
        ['银行开户', '营业执照、公章', '选择银行办理', '基本户只能一个', '多家银行可选', '任意银行', '3-5天', '年费', '开户费用'],
        ['无房产证', '租赁合同+房东证明', '提交替代材料', '社区出具证明', '支持替代材料', '政务服务中心', '1-2天', '免费', '材料真实性'],
        ['未过户房产', '网签合同/交易协议', '提供合同+证明', '原房主同意', '支持网签合同', '政务服务中心', '1-2天', '免费', '产权纠纷'],
        ['原房主逝世', '继承证明+死亡证明', '继承公证→过户→注册', '需先办理继承公证', '支持继承证明', '政务服务中心', '7-15天', '公证费', '继承纠纷'],
        ['创业补贴', '营业执照、就业证明', '线上申请→审核', '毕业5年内/就业困难', '最高5000元', '人社局', '30天', '免费', '条件不符']
    ],
    10: [
        ['模型类型', '代表产品', '部署方式', 'GPU需求', '算力费用/月', 'API调用成本', '优点', '缺点', '适用场景'],
        ['多模态大模型', 'GPT-4V、Gemini Ultra', '🔴API调用', '4xA100/80GB', '¥8-15万', 'GPT-4V: $0.01/image', '能力最强、多模态融合', '成本高、数据安全风险', '创意设计、内容创作'],
        ['文本大模型', 'GPT-4、Claude 3', 'API调用', '8xA100', '¥12-25万', 'GPT-4: $0.03/1K tokens', '推理强、上下文深', '成本高、无本地化', '代码生成、数据分析'],
        ['开源大模型', 'Llama 3-70B、Qwen-72B', '私有化部署', '4xA100/40GB', '¥5-12万', '免费', '数据安全、定制灵活', '技术门槛高', '企业内部系统'],
        ['垂直领域模型', 'Med-PaLM、CodeLlama', '本地/私有化', '2xA100', '¥3-8万', '按需', '专业知识强、精准输出', '通用性弱', '医疗、法律、编程'],
        ['国产大模型', '文心一言、豆包', 'API+私有化', '8xH100', '¥15-30万', '豆包: ¥99/月起', '本土化支持、合规性好', '能力略逊', '政企服务'],
        ['轻量模型', 'Phi-3、Llama 3-8B', '本地部署', 'RTX 4090', '¥0.5-2万', '免费', '速度快、资源需求低', '能力有限', '移动端、IoT'],
        ['DeepSeek系列', 'DeepSeek-R1', 'API+本地', '4xA100', '¥6-12万', 'API: ¥0.002/1K tokens', '长上下文、数学推理强', '知名度低', '科研教育'],
        ['Qwen系列', 'Qwen 2-72B', '私有化部署', '4xA100', '¥5-10万', '免费商用', '多语言、社区活跃', '生态相对小', '跨境服务'],
        ['Doubao系列', '豆包企业版', 'API调用', '云端', '按需', '企业版: ¥1999/月起', '娱乐性强、生活化', '专业能力有限', '消费级应用'],
        ['云端租赁', 'AWS/GCP/Azure', '云部署', '按需配置', '¥2-20万', '按需付费', '弹性伸缩、免维护', '成本波动', '弹性业务'],
        ['边缘部署', 'NVIDIA Jetson', '本地边缘', 'Jetson AGX', '¥0.3-1万', '免费', '低延迟、隐私保护', '能力有限', 'IoT、边缘计算']
    ],
    11: [
        ['章节', '核心内容', '关键技术', '业务价值', '实战要点', '投资回报', '蚌埠机遇', '核心价值'],
        ['第1章', 'A2A协议架构', 'Agent Cards、标准化消息', '🔴互操作性提升60%', '跨平台协作', '18个月', '生态整合', '打破平台壁垒'],
        ['第2章', '协议规范', 'JSON-RPC 2.0、gRPC', '🔴通信效率提升10x', '低延迟通信', '12个月', '技术基础', '标准化接口'],
        ['第3章', '安全合规', 'OAuth2、mTLS、零信任', '🔴合规成本降低40%', '金融级安全', '24个月', '信任背书', '数据安全'],
        ['第4章', '部署运维', 'K8s、自动扩缩容', '🔴可用性99.99%', '高可用架构', '18个月', '基础设施', '业务连续性'],
        ['第5章', '性能优化', '缓存、异步IO、AI推理', '🔴性能提升3x', '大规模承载', '12个月', '效率提升', '用户体验'],
        ['第6章', '实战案例', '供应链、风控、医疗', '🔴效率提升40%', '行业落地', '18-36个月', '案例借鉴', '价值验证'],
        ['第7章', 'AI趋势', '2026-2046演进', '🔴抢占先机', '战略布局', '长期', '产业升级', '前瞻布局'],
        ['第8章', '人才规划', '全人群招聘策略', '🔴人才储备', '本地就业', '持续', '人力资源', '团队建设'],
        ['第9章', '创业指南', '执照办理、政策支持', '🔴创业便利化', '零成本注册', '短期', '创业环境', '营商环境'],
        ['第10章', '大模型生态', '多模态、开源、费用', '🔴选型决策', '成本优化', '持续', '技术选型', 'AI能力'],
        ['总结', 'A2A+AI智能体体系', '端到端解决方案', '🔴企业智能化升级', '技术+人才+生态', '长期价值', '蚌埠AI产业', '全面发展']
    ]
}

CHAPTER_DESCRIPTIONS = {
    1: "A2A协议是下一代Agent-to-Agent通信标准，实现异构AI Agent之间的无缝协作，打破平台壁垒，构建开放的AI生态系统。其核心价值包括互操作性、弹性扩展、安全合规和智能协作能力。",
    2: "协议规范定义了Agent间通信的标准化消息格式和交互模式，基于JSON-RPC 2.0实现跨平台、跨语言的互操作性，支持多层级的协议栈设计。",
    3: "企业级安全架构包括API网关、WAF防火墙、OAuth 2.0认证、mTLS加密等多层防护，确保数据安全和合规要求，满足金融、医疗等行业的严格监管标准。",
    4: "基于Kubernetes的高可用部署方案，支持自动扩缩容、健康检查、蓝绿发布等运维能力，保障业务连续性和系统稳定性。",
    5: "通过网络优化、消息批量、缓存策略、并发处理等手段，实现系统性能的全方位提升，支撑大规模业务场景下的高效运行。",
    6: "通过供应链管理、金融风控、智能办公、医疗诊断等真实案例，验证A2A协议的业务价值和技术可行性，为企业数字化转型提供参考。",
    7: "AI智能体正处于快速发展期，未来10-20年将经历标准化、智能化到生态化的演进。蚌埠市作为制造业基地，可抓住长三角一体化机遇，布局AI Agent产业，推动传统产业升级。",
    8: "针对蚌埠市产业特点，制定差异化人才策略：高中学历推荐使用BOSS直聘、58同城找基础岗位；精神残疾人可通过中国残疾人就业网络服务平台、微信小程序等官方渠道找工作；各年龄段和学历层次都有相应的招聘渠道和培养方案。",
    9: "蚌埠市为创业者提供完善的政策支持，包括免费注册、创业补贴、税收优惠等。针对无房产证、未过户房产及原房主逝世等特殊情况，提供灵活的场地证明解决方案：原房主逝世需先办理继承公证和房产过户，凭继承证明可正常注册。",
    10: "AI大模型正处于爆发式发展阶段，算力成本是关键考量因素。API调用方式：GPT-4V每图$0.01，GPT-4每千token$0.03，豆包个人版¥99/月起。私有化部署：Llama 3-70B需4xA100，月费用¥5-12万；轻量模型如Phi-3仅需RTX 4090，月费用¥0.5-2万。云端租赁（AWS/GCP/Azure）按需付费，月费用¥2-20万。企业需根据数据敏感性和预算选择合适的部署方案。",
    11: "本PPT系统阐述了A2A协议与AI智能体的完整知识体系：从协议架构到安全合规，从部署运维到性能优化，从实战案例到未来趋势，从人才规划到创业指南，全面覆盖AI智能体技术栈。**核心价值**在于构建端到端的企业智能化解决方案，助力蚌埠市AI产业升级与人才发展。"
}

CHAPTER_RESOURCES = {
    1: [
        '官方文档: https://a2a-protocol.org',
        'Linux Foundation: https://linuxfoundation.org/press-release/a2a-protocol/',
        'DeepLearning.AI课程: https://www.deeplearning.ai/courses/intro-to-a2a-protocol/',
        'GitHub仓库: https://github.com/a2aproject'
    ],
    2: [
        'JSON-RPC 2.0规范: https://www.jsonrpc.org/specification',
        'HTTP/2协议: https://http2.github.io/',
        'gRPC官方文档: https://grpc.io/docs/',
        'QUIC协议: https://datatracker.ietf.org/wg/quic/about/'
    ],
    3: [
        'OAuth 2.0规范: https://oauth.net/2/',
        'mTLS指南: https://www.envoyproxy.io/docs/envoy/latest/security/mozilla',
        'WAF配置: https://www.modsecurity.org/CRS/Documentation/',
        'Kong网关: https://docs.konghq.com/'
    ],
    4: [
        'Kubernetes文档: https://kubernetes.io/zh-cn/docs/',
        'Prometheus监控: https://prometheus.io/docs/introduction/overview/',
        'Grafana可视化: https://grafana.com/docs/grafana/latest/',
        'Istio服务网格: https://istio.io/latest/zh/docs/'
    ],
    5: [
        'Redis缓存: https://redis.io/docs/',
        'gRPC性能优化: https://grpc.io/docs/guides/performance/',
        '批量处理: https://docs.python.org/3/library/concurrent.futures.html',
        'CDN加速: https://www.cloudflare.com/learning/cdn/what-is-a-cdn/'
    ],
    6: [
        '供应链案例: https://www.supplychaindive.com/',
        '金融风控案例: https://www.fintechnexus.com/',
        '多Agent系统: https://arxiv.org/abs/2401.07580',
        '最佳实践: https://www.agentprotocol.io/'
    ],
    7: [
        'AI趋势报告: https://www.gartner.com/en/articles/what-s-new-in-artificial-intelligence-from-the-2024-gartner-hype-cycle',
        'Agent未来展望: https://arxiv.org/abs/2403.07691',
        '中国AI发展: https://www.caict.ac.cn/',
        '长三角一体化: https://www.yangtze-river-delta.gov.cn/'
    ],
    8: [
        '蚌埠人才网: http://www.bbrc.com.cn/',
        '蚌埠人社局: http://rsj.bengbu.gov.cn/',
        '残疾人就业平台: https://www.cdpee.cn/',
        'BOSS直聘: https://www.zhipin.com/'
    ],
    9: [
        '安徽政务服务网: http://www.ahzwfw.gov.cn/',
        '蚌埠市场监管局: http://amr.bengbu.gov.cn/',
        '创业补贴政策: http://rsj.bengbu.gov.cn/xxgk/gsgg/202403/t20240306_3567238.html',
        '皖事通APP: https://www.ahzwfw.gov.cn/col/col13206/index.html'
    ],
    10: [
        'GPT-4V: https://openai.com/research/gpt-4v-system-card',
        'Gemini: https://deepmind.google/technologies/gemini/',
        'Llama 3: https://ai.meta.com/llama/',
        'Qwen: https://qwenlm.github.io/',
        'DeepSeek: https://www.deepseek.com/',
        '豆包: https://www.doubao.com/'
    ]
}

def add_watermark(slide):
    """在幻灯片添加水印"""
    shapes = slide.shapes
    
    for i in range(3):
        for j in range(4):
            left = Inches(1.5 + i * 5)
            top = Inches(1.8 + j * 2)
            
            watermark_box = slide.shapes.add_textbox(left, top, Inches(4), Inches(1))
            watermark_frame = watermark_box.text_frame
            p = watermark_frame.add_paragraph()
            p.text = WATERMARK_TEXT
            p.font.size = Pt(10)
            p.font.color.rgb = RGBColor(100, 120, 150)
            p.font.bold = False
            p.alignment = PP_ALIGN.CENTER
            
            watermark_box.rotation = -15
            watermark_box.fill.transparency = 0.85

def add_copyright(slide):
    """在幻灯片底部添加版权信息"""
    copyright_box = slide.shapes.add_textbox(Inches(0.5), Inches(8.4), Inches(15), Inches(0.5))
    copyright_frame = copyright_box.text_frame
    p = copyright_frame.add_paragraph()
    p.text = COPYRIGHT_INFO
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(120, 140, 160)
    p.alignment = PP_ALIGN.LEFT

def create_robot_face(slide, left, top, width, colors):
    """创建机器人脸部图形"""
    shapes = slide.shapes
    
    head = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, width * 0.9)
    head.fill.solid()
    head.fill.fore_color.rgb = colors['accent']
    head.fill.transparency = 0.3
    head.line.fill.background()
    
    eye_width = width * 0.15
    eye_height = width * 0.15
    eye_left = left + width * 0.25
    eye_top = top + width * 0.25
    
    left_eye = shapes.add_shape(MSO_SHAPE.OVAL, eye_left, eye_top, eye_width, eye_height)
    left_eye.fill.solid()
    left_eye.fill.fore_color.rgb = RGBColor(255, 255, 255)
    left_eye.line.fill.background()
    
    right_eye = shapes.add_shape(MSO_SHAPE.OVAL, eye_left + width * 0.35, eye_top, eye_width, eye_height)
    right_eye.fill.solid()
    right_eye.fill.fore_color.rgb = RGBColor(255, 255, 255)
    right_eye.line.fill.background()
    
    pupil = shapes.add_shape(MSO_SHAPE.OVAL, eye_left + eye_width * 0.25, eye_top + eye_height * 0.25, eye_width * 0.5, eye_height * 0.5)
    pupil.fill.solid()
    pupil.fill.fore_color.rgb = RGBColor(0, 0, 0)
    pupil.line.fill.background()
    
    pupil2 = shapes.add_shape(MSO_SHAPE.OVAL, eye_left + width * 0.35 + eye_width * 0.25, eye_top + eye_height * 0.25, eye_width * 0.5, eye_height * 0.5)
    pupil2.fill.solid()
    pupil2.fill.fore_color.rgb = RGBColor(0, 0, 0)
    pupil2.line.fill.background()
    
    mouth = shapes.add_shape(MSO_SHAPE.RECTANGLE, left + width * 0.35, top + width * 0.6, width * 0.3, width * 0.08)
    mouth.fill.solid()
    mouth.fill.fore_color.rgb = RGBColor(255, 255, 255)
    mouth.fill.transparency = 0.5
    mouth.line.fill.background()

def add_ai_robot_element(slide, chapter_num, colors):
    """根据章节主题添加符合内容的AI机器人元素"""
    shapes = slide.shapes
    
    if chapter_num == 1:
        create_robot_face(slide, Inches(13.5), Inches(1), Inches(2), colors)
        
        circle = shapes.add_shape(MSO_SHAPE.OVAL, Inches(14), Inches(3.2), Inches(1.2), Inches(1.2))
        circle.fill.solid()
        circle.fill.fore_color.rgb = colors['accent']
        circle.fill.transparency = 0.4
        circle.line.fill.background()
        
        hexagon = shapes.add_shape(MSO_SHAPE.HEXAGON, Inches(14.2), Inches(4.5), Inches(0.8), Inches(0.8))
        hexagon.fill.solid()
        hexagon.fill.fore_color.rgb = colors['secondary']
        hexagon.fill.transparency = 0.5
        hexagon.line.fill.background()
        
    elif chapter_num == 2:
        left = Inches(13.8)
        top = Inches(1)
        
        for i in range(4):
            rect = shapes.add_shape(MSO_SHAPE.RECTANGLE, left + i * Inches(0.3), top + i * Inches(0.5), Inches(1), Inches(0.4))
            rect.fill.solid()
            rect.fill.fore_color.rgb = colors['accent']
            rect.fill.transparency = 0.5 - i * 0.1
            rect.line.fill.background()
        
        arrow = shapes.add_shape(MSO_SHAPE.UP_ARROW, Inches(14.3), Inches(3.5), Inches(0.8), Inches(0.8))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = colors['secondary']
        arrow.fill.transparency = 0.6
        arrow.line.fill.background()
        
    elif chapter_num == 3:
        outer_circle = shapes.add_shape(MSO_SHAPE.OVAL, Inches(13.5), Inches(1), Inches(2), Inches(2))
        outer_circle.fill.solid()
        outer_circle.fill.fore_color.rgb = colors['accent']
        outer_circle.fill.transparency = 0.2
        outer_circle.line.fill.background()
        
        middle_circle = shapes.add_shape(MSO_SHAPE.OVAL, Inches(13.8), Inches(1.3), Inches(1.4), Inches(1.4))
        middle_circle.fill.solid()
        middle_circle.fill.fore_color.rgb = colors['secondary']
        middle_circle.fill.transparency = 0.3
        middle_circle.line.fill.background()
        
        inner_circle = shapes.add_shape(MSO_SHAPE.OVAL, Inches(14.1), Inches(1.6), Inches(0.8), Inches(0.8))
        inner_circle.fill.solid()
        inner_circle.fill.fore_color.rgb = RGBColor(255, 255, 255)
        inner_circle.fill.transparency = 0.5
        inner_circle.line.fill.background()
        
        pentagon = shapes.add_shape(MSO_SHAPE.PENTAGON, Inches(14), Inches(3), Inches(1), Inches(1.2))
        pentagon.fill.solid()
        pentagon.fill.fore_color.rgb = colors['accent']
        pentagon.fill.transparency = 0.4
        pentagon.line.fill.background()
        
    elif chapter_num == 4:
        for i in range(3):
            layer = shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(13.5 + i * 0.2), Inches(1 + i * 0.3), Inches(2 - i * 0.4), Inches(0.6))
            layer.fill.solid()
            layer.fill.fore_color.rgb = colors['accent']
            layer.fill.transparency = 0.5 - i * 0.15
            layer.line.fill.background()
        
        server = shapes.add_shape(MSO_SHAPE.TRAPEZOID, Inches(13.8), Inches(2.5), Inches(1.5), Inches(1))
        server.fill.solid()
        server.fill.fore_color.rgb = colors['secondary']
        server.fill.transparency = 0.4
        server.line.fill.background()
        
    elif chapter_num == 5:
        bars = [0.8, 1.2, 0.9, 1.4, 1.1]
        for i, height in enumerate(bars):
            bar = shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(13.6 + i * 0.3), Inches(2.5 - height * 0.5), Inches(0.2), Inches(height))
            bar.fill.solid()
            bar.fill.fore_color.rgb = colors['accent']
            bar.fill.transparency = 0.5
            bar.line.fill.background()
        
        speed_lines = shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(14), Inches(1), Inches(1.2), Inches(1))
        speed_lines.fill.solid()
        speed_lines.fill.fore_color.rgb = colors['secondary']
        speed_lines.fill.transparency = 0.4
        speed_lines.line.fill.background()
        
    elif chapter_num == 6:
        create_robot_face(slide, Inches(13.5), Inches(1), Inches(1.8), colors)
        
        hexagon1 = shapes.add_shape(MSO_SHAPE.HEXAGON, Inches(13.8), Inches(3), Inches(1), Inches(1))
        hexagon1.fill.solid()
        hexagon1.fill.fore_color.rgb = colors['accent']
        hexagon1.fill.transparency = 0.4
        hexagon1.line.fill.background()
        
        hexagon2 = shapes.add_shape(MSO_SHAPE.HEXAGON, Inches(14.5), Inches(3.8), Inches(0.8), Inches(0.8))
        hexagon2.fill.solid()
        hexagon2.fill.fore_color.rgb = colors['secondary']
        hexagon2.fill.transparency = 0.5
        hexagon2.line.fill.background()

def add_chapter_background(slide, chapter_num):
    """添加章节特定风格背景"""
    colors = CHAPTER_COLORS.get(chapter_num, CHAPTER_COLORS[1])
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors['bg']
    
    shapes = slide.shapes
    
    gradient_oval = shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(-2), Inches(12), Inches(10))
    gradient_oval.fill.solid()
    gradient_oval.fill.fore_color.rgb = colors['accent']
    gradient_oval.fill.transparency = 0.12
    gradient_oval.line.fill.background()
    
    bottom_bar = shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(8.5), Inches(16), Inches(0.5))
    bottom_bar.fill.solid()
    bottom_bar.fill.fore_color.rgb = colors['accent']
    bottom_bar.fill.transparency = 0.2
    bottom_bar.line.fill.background()
    
    return colors

def create_title_slide(prs):
    """创建封面页"""
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(8, 12, 22)
    
    shapes = slide.shapes
    
    large_gradient = shapes.add_shape(MSO_SHAPE.OVAL, Inches(6), Inches(-3), Inches(14), Inches(12))
    large_gradient.fill.solid()
    large_gradient.fill.fore_color.rgb = RGBColor(0, 120, 200)
    large_gradient.fill.transparency = 0.18
    large_gradient.line.fill.background()
    
    hexagon_decor = shapes.add_shape(MSO_SHAPE.HEXAGON, Inches(13), Inches(0.5), Inches(2.5), Inches(2.5))
    hexagon_decor.fill.solid()
    hexagon_decor.fill.fore_color.rgb = RGBColor(0, 180, 255)
    hexagon_decor.fill.transparency = 0.4
    hexagon_decor.line.fill.background()
    
    circle_decor = shapes.add_shape(MSO_SHAPE.OVAL, Inches(14), Inches(2.5), Inches(1), Inches(1))
    circle_decor.fill.solid()
    circle_decor.fill.fore_color.rgb = RGBColor(100, 220, 255)
    circle_decor.fill.transparency = 0.6
    circle_decor.line.fill.background()
    
    create_robot_face(slide, Inches(0.5), Inches(5.5), Inches(1.5), {'accent': RGBColor(0, 180, 255), 'secondary': RGBColor(50, 120, 180)})
    
    title_box = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(12), Inches(2.5))
    title_frame = title_box.text_frame
    p = title_frame.add_paragraph()
    p.text = "A2A协议与AI智能体：技术、趋势、人才与大模型生态"
    p.font.size = Pt(34)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide.shapes.add_textbox(Inches(2.5), Inches(4.2), Inches(11), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    p = subtitle_frame.add_paragraph()
    p.text = "协议架构 · 安全合规 · 部署运维 · 大模型趋势 · 蚌埠实践"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0, 200, 255)
    p.alignment = PP_ALIGN.CENTER
    
    presenter_box = slide.shapes.add_textbox(Inches(3), Inches(5.5), Inches(10), Inches(1))
    presenter_frame = presenter_box.text_frame
    p = presenter_frame.add_paragraph()
    p.text = "【主讲人：何晓冬】"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(255, 200, 100)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    info_box = slide.shapes.add_textbox(Inches(3), Inches(7), Inches(10), Inches(1))
    info_frame = info_box.text_frame
    p = info_frame.add_paragraph()
    p.text = "版本: v3.1 | 日期: 2026年5月 | © 2026 何晓冬"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(150, 170, 190)
    p.alignment = PP_ALIGN.CENTER
    
    add_copyright(slide)
    
    return slide

def create_chapter_slide(prs, chapter):
    """创建章节幻灯片（含表格、AI元素、红色重点和拓展资源）"""
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    colors = add_chapter_background(slide, chapter['num'])
    
    add_watermark(slide)
    
    add_ai_robot_element(slide, chapter['num'], colors)
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(12), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.add_paragraph()
    p.text = chapter['title']
    p.font.size = Pt(22)
    p.font.color.rgb = colors['accent']
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    desc = CHAPTER_DESCRIPTIONS.get(chapter['num'], "")
    if desc:
        desc_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(12), Inches(0.6))
        desc_frame = desc_box.text_frame
        desc_frame.word_wrap = True
        p = desc_frame.add_paragraph()
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = colors['text']
        p.alignment = PP_ALIGN.LEFT
    
    if chapter['tables']:
        table_data = chapter['tables'][0]
        create_table(slide, table_data, colors)
    
    resources = CHAPTER_RESOURCES.get(chapter['num'], [])
    if resources:
        create_resources_box(slide, resources, colors)
    
    add_copyright(slide)
    
    return slide

def create_table(slide, table_data, colors):
    """在幻灯片中创建表格，红色标记重点内容"""
    num_rows = len(table_data)
    num_cols = len(table_data[0]) if num_rows > 0 else 3
    
    left = Inches(0.4)
    top = Inches(2.0)
    width = Inches(14.8)
    height = Inches(4.2)
    
    table = slide.shapes.add_table(num_rows, num_cols, left, top, width, height).table
    
    highlight_keywords = ['核心', '关键', '重要', '必须', '最佳', '提升', '降低', '99', 'SOC2', '百万', '🔴']
    
    for i, row in enumerate(table_data):
        for j, cell in enumerate(row):
            cell_text = cell.replace('🔴', '')
            
            table.cell(i, j).text = cell_text
            
            paragraph = table.cell(i, j).text_frame.paragraphs[0]
            if num_cols >= 9:
                paragraph.font.size = Pt(6.8)
            elif num_cols == 8:
                paragraph.font.size = Pt(7.2)
            elif num_cols >= 6:
                paragraph.font.size = Pt(7.8)
            elif num_cols == 5:
                paragraph.font.size = Pt(8.5)
            else:
                paragraph.font.size = Pt(9)
            paragraph.alignment = PP_ALIGN.CENTER
            
            if i == 0:
                paragraph.font.color.rgb = RGBColor(255, 255, 255)
                paragraph.font.bold = True
                table.cell(i, j).fill.solid()
                table.cell(i, j).fill.fore_color.rgb = colors['accent']
            else:
                is_highlight = any(keyword in cell for keyword in highlight_keywords)
                if is_highlight:
                    paragraph.font.color.rgb = colors['highlight']
                    paragraph.font.bold = True
                else:
                    paragraph.font.color.rgb = colors['text']
                table.cell(i, j).fill.solid()
                table.cell(i, j).fill.fore_color.rgb = RGBColor(25, 30, 40)
                table.cell(i, j).fill.transparency = 0.3

def create_resources_box(slide, resources, colors):
    """在幻灯片底部创建拓展资源区域"""
    res_box = slide.shapes.add_textbox(Inches(0.7), Inches(6.5), Inches(14), Inches(1.7))
    res_frame = res_box.text_frame
    res_frame.word_wrap = True
    
    p = res_frame.add_paragraph()
    p.text = "📚 拓展学习资源"
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(100, 220, 255)
    p.font.bold = True
    
    for res in resources[:3]:
        p = res_frame.add_paragraph()
        p.text = res
        p.font.size = Pt(8)
        p.font.color.rgb = RGBColor(170, 185, 200)

def create_summary_slide(prs):
    """创建总结页"""
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(10, 15, 28)
    
    shapes = slide.shapes
    
    gradient_oval = shapes.add_shape(MSO_SHAPE.OVAL, Inches(5), Inches(-2), Inches(14), Inches(12))
    gradient_oval.fill.solid()
    gradient_oval.fill.fore_color.rgb = RGBColor(0, 120, 200)
    gradient_oval.fill.transparency = 0.15
    gradient_oval.line.fill.background()
    
    bottom_bar = shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(8.5), Inches(16), Inches(0.5))
    bottom_bar.fill.solid()
    bottom_bar.fill.fore_color.rgb = RGBColor(0, 180, 255)
    bottom_bar.fill.transparency = 0.2
    bottom_bar.line.fill.background()
    
    add_watermark(slide)
    
    title_box = slide.shapes.add_textbox(Inches(3), Inches(0.5), Inches(10), Inches(1))
    title_frame = title_box.text_frame
    p = title_frame.add_paragraph()
    p.text = "🎯 A2A协议核心要点总结"
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    summary_data = [
        ['章节', '核心内容', '关键价值', '技术亮点'],
        ['第一章', '去中心化架构', '高可用弹性', '标准化协议'],
        ['第二章', 'JSON-RPC规范', '跨平台互操作', '异步通信'],
        ['第三章', 'OAuth2+mTLS', '🔴企业级安全', '零信任'],
        ['第四章', 'K8s部署', '🔴99.99%可用', '自动扩缩容'],
        ['第五章', '缓存+批量', '🔴性能提升3x', '低延迟'],
        ['第六章', '实战案例', '业务价值验证', '多Agent协作']
    ]
    
    table = slide.shapes.add_table(7, 4, Inches(0.8), Inches(1.5), Inches(14.5), Inches(5.5)).table
    
    highlight_keywords = ['核心', '关键', '重要', '必须', '最佳', '提升', '降低', '99', 'SOC2', '百万', '🔴']
    
    for i, row in enumerate(summary_data):
        for j, cell in enumerate(row):
            cell_text = cell.replace('🔴', '')
            table.cell(i, j).text = cell_text
            paragraph = table.cell(i, j).text_frame.paragraphs[0]
            paragraph.font.size = Pt(11)
            paragraph.alignment = PP_ALIGN.CENTER
            
            if i == 0:
                paragraph.font.color.rgb = RGBColor(255, 255, 255)
                paragraph.font.bold = True
                table.cell(i, j).fill.solid()
                table.cell(i, j).fill.fore_color.rgb = RGBColor(0, 180, 255)
            else:
                is_highlight = any(keyword in row[j] for keyword in highlight_keywords)
                if is_highlight:
                    paragraph.font.color.rgb = RGBColor(255, 100, 100)
                    paragraph.font.bold = True
                else:
                    paragraph.font.color.rgb = RGBColor(220, 230, 240)
                table.cell(i, j).fill.solid()
                table.cell(i, j).fill.fore_color.rgb = RGBColor(25, 30, 40)
    
    add_copyright(slide)
    
    return slide

def create_thank_you_slide(prs):
    """创建感谢页"""
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(8, 12, 22)
    
    shapes = slide.shapes
    
    large_gradient = shapes.add_shape(MSO_SHAPE.OVAL, Inches(5), Inches(-1), Inches(10), Inches(10))
    large_gradient.fill.solid()
    large_gradient.fill.fore_color.rgb = RGBColor(0, 100, 180)
    large_gradient.fill.transparency = 0.25
    large_gradient.line.fill.background()
    
    hexagon_decor = shapes.add_shape(MSO_SHAPE.HEXAGON, Inches(7.5), Inches(2.5), Inches(2), Inches(2))
    hexagon_decor.fill.solid()
    hexagon_decor.fill.fore_color.rgb = RGBColor(0, 200, 255)
    hexagon_decor.fill.transparency = 0.35
    hexagon_decor.line.fill.background()
    
    create_robot_face(slide, Inches(1), Inches(5), Inches(1.2), {'accent': RGBColor(0, 180, 255), 'secondary': RGBColor(50, 120, 180)})
    create_robot_face(slide, Inches(13.8), Inches(5), Inches(1.2), {'accent': RGBColor(0, 180, 255), 'secondary': RGBColor(50, 120, 180)})
    
    thank_box = slide.shapes.add_textbox(Inches(2), Inches(4), Inches(12), Inches(2.5))
    thank_frame = thank_box.text_frame
    p = thank_frame.add_paragraph()
    p.text = "THANK YOU"
    p.font.size = Pt(52)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    p = thank_frame.add_paragraph()
    p.text = "谢谢"
    p.font.size = Pt(36)
    p.font.color.rgb = RGBColor(0, 200, 255)
    p.alignment = PP_ALIGN.CENTER
    
    presenter_box = slide.shapes.add_textbox(Inches(3), Inches(6.8), Inches(10), Inches(1))
    presenter_frame = presenter_box.text_frame
    p = presenter_frame.add_paragraph()
    p.text = "【主讲人：何晓冬】"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(255, 200, 100)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    add_copyright(slide)
    
    return slide

def parse_markdown_for_tables(md_content):
    """解析Markdown，提取章节表格数据，按正确顺序排列"""
    chapters = {}
    
    lines = md_content.split('\n')
    i = 0
    
    while i < len(lines):
        line = lines[i]
        
        if line.startswith('# 第') and '章：' in line:
            title = line[2:].strip()
            
            chapter_num = None
            for idx, expected_title in enumerate(CHAPTER_TITLES, 1):
                if title.startswith(expected_title[:5]):
                    chapter_num = idx
                    break
            
            if chapter_num is None:
                i += 1
                continue
            
            tables = []
            i += 1
            
            while i < len(lines):
                if lines[i].startswith('|'):
                    table_data = []
                    while i < len(lines) and lines[i].startswith('|'):
                        row = lines[i].strip().split('|')
                        row = [cell.strip() for cell in row if cell.strip()]
                        table_data.append(row)
                        i += 1
                    if len(table_data) > 1:
                        tables.append(table_data)
                elif lines[i].startswith('# 第') and '章：' in lines[i]:
                    break
                elif lines[i].startswith('---'):
                    i += 1
                    continue
                else:
                    i += 1
            
            if not tables and chapter_num in DEFAULT_TABLES:
                tables.append(DEFAULT_TABLES[chapter_num])
            
            if tables:
                chapters[chapter_num] = {
                    'num': chapter_num,
                    'title': title,
                    'tables': tables
                }
        
        i += 1
    
    result = []
    for num in range(1, 12):
        if num in chapters:
            result.append(chapters[num])
        elif num in DEFAULT_TABLES:
            result.append({
                'num': num,
                'title': CHAPTER_TITLES[num-1],
                'tables': [DEFAULT_TABLES[num]]
            })
    
    return result

def create_pptx_compact(chapters, output_path):
    """创建精简版PPTX文件"""
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    create_title_slide(prs)
    
    for chapter in chapters:
        create_chapter_slide(prs, chapter)
    
    create_summary_slide(prs)
    create_thank_you_slide(prs)
    
    prs.save(output_path)
    print(f"PPTX文件已生成: {output_path}")
    print(f"共生成 {len(chapters) + 3} 页幻灯片")

if __name__ == '__main__':
    md_file = r'f:\个人作品\legal-rag-qa-system\A2A_ENTERPRISE_PPT.md'
    output_file = r'f:\个人作品\legal-rag-qa-system\A2A_ENTERPRISE_PPT_V4.pptx'
    
    with open(md_file, 'r', encoding='utf-8') as f:
        md_content = f.read()
    
    chapters = parse_markdown_for_tables(md_content)
    print(f"检测到 {len(chapters)} 个章节")
    for ch in chapters:
        print(f"  第{ch['num']}章: {ch['title']}")
    
    create_pptx_compact(chapters, output_file)
