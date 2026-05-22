from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import re

CORPORATE_COLORS = {
    'primary': RGBColor(30, 58, 138),
    'secondary': RGBColor(59, 130, 246),
    'accent': RGBColor(139, 92, 246),
    'dark': RGBColor(15, 23, 42),
    'light': RGBColor(241, 245, 249),
    'highlight': RGBColor(239, 68, 68),
    'gold': RGBColor(251, 191, 36),
    'gray': RGBColor(100, 116, 139)
}

CHAPTER_COLORS = {
    1: {'bg': RGBColor(15, 23, 42), 'accent': RGBColor(59, 130, 246), 'text': RGBColor(241, 245, 249), 'highlight': RGBColor(239, 68, 68), 'secondary': RGBColor(30, 64, 175)},
    2: {'bg': RGBColor(12, 18, 33), 'accent': RGBColor(139, 92, 246), 'text': RGBColor(241, 245, 249), 'highlight': RGBColor(239, 68, 68), 'secondary': RGBColor(79, 70, 229)},
    3: {'bg': RGBColor(18, 12, 33), 'accent': RGBColor(236, 72, 153), 'text': RGBColor(241, 245, 249), 'highlight': RGBColor(239, 68, 68), 'secondary': RGBColor(192, 132, 252)},
    4: {'bg': RGBColor(18, 22, 12), 'accent': RGBColor(251, 191, 36), 'text': RGBColor(254, 249, 195), 'highlight': RGBColor(239, 68, 68), 'secondary': RGBColor(234, 179, 8)},
    5: {'bg': RGBColor(33, 4, 4), 'accent': RGBColor(239, 68, 68), 'text': RGBColor(254, 226, 226), 'highlight': RGBColor(251, 191, 36), 'secondary': RGBColor(220, 38, 38)},
    6: {'bg': RGBColor(28, 18, 15), 'accent': RGBColor(249, 115, 22), 'text': RGBColor(254, 243, 199), 'highlight': RGBColor(239, 68, 68), 'secondary': RGBColor(234, 88, 12)},
    7: {'bg': RGBColor(12, 28, 22), 'accent': RGBColor(16, 185, 129), 'text': RGBColor(220, 252, 231), 'highlight': RGBColor(239, 68, 68), 'secondary': RGBColor(5, 150, 105)},
    8: {'bg': RGBColor(25, 15, 33), 'accent': RGBColor(139, 92, 246), 'text': RGBColor(241, 245, 249), 'highlight': RGBColor(239, 68, 68), 'secondary': RGBColor(99, 102, 241)},
    9: {'bg': RGBColor(15, 25, 12), 'accent': RGBColor(34, 197, 94), 'text': RGBColor(220, 252, 231), 'highlight': RGBColor(239, 68, 68), 'secondary': RGBColor(22, 163, 74)},
    10: {'bg': RGBColor(33, 15, 33), 'accent': RGBColor(236, 72, 153), 'text': RGBColor(253, 224, 239), 'highlight': RGBColor(239, 68, 68), 'secondary': RGBColor(219, 39, 119)}
}

def read_markdown(file_path):
    with open(file_path, 'r', encoding='utf-8') as f:
        content = f.read()
    return content

def parse_markdown(content):
    chapters = []
    lines = content.split('\n')
    
    current_chapter = None
    current_section_lines = []
    
    for line in lines:
        if line.startswith('## '):
            if current_chapter:
                if current_section_lines:
                    current_chapter['sections'].append('\n'.join(current_section_lines))
                chapters.append(current_chapter)
            
            chapter_title = line[3:].strip()
            num_match = re.search(r'项目(\d+)', chapter_title)
            if num_match:
                chapter_num = int(num_match.group(1))
                title_text = re.sub(r'^[📍\s]*项目\d+[:：]\s*', '', chapter_title)
                current_chapter = {
                    'num': chapter_num,
                    'title': title_text,
                    'sections': []
                }
                current_section_lines = []
            else:
                current_chapter = None
                current_section_lines = []
        
        elif current_chapter and line.strip() == '---':
            if current_section_lines:
                current_chapter['sections'].append('\n'.join(current_section_lines))
            current_section_lines = []
        
        elif current_chapter:
            current_section_lines.append(line)
    
    if current_chapter:
        if current_section_lines:
            current_chapter['sections'].append('\n'.join(current_section_lines))
        chapters.append(current_chapter)
    
    return chapters

def create_network_node(slide, left, top, size, colors):
    shapes = slide.shapes
    main_circle = shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    main_circle.fill.solid()
    main_circle.fill.fore_color.rgb = colors['accent']
    main_circle.fill.transparency = 0.2
    main_circle.line.color.rgb = colors['accent']
    main_circle.line.width = Pt(1.5)
    
    inner_ring = shapes.add_shape(MSO_SHAPE.OVAL, left + size * 0.2, top + size * 0.2, size * 0.6, size * 0.6)
    inner_ring.fill.solid()
    inner_ring.fill.fore_color.rgb = colors['secondary']
    inner_ring.fill.transparency = 0.5
    inner_ring.line.fill.background()
    
    center_dot = shapes.add_shape(MSO_SHAPE.OVAL, left + size * 0.35, top + size * 0.35, size * 0.3, size * 0.3)
    center_dot.fill.solid()
    center_dot.fill.fore_color.rgb = RGBColor(255, 255, 255)
    center_dot.line.fill.background()

def create_corporate_logo_element(slide, left, top, size, colors):
    shapes = slide.shapes
    outer_hex = shapes.add_shape(MSO_SHAPE.HEXAGON, left, top, size, size)
    outer_hex.fill.solid()
    outer_hex.fill.fore_color.rgb = colors['accent']
    outer_hex.fill.transparency = 0.25
    outer_hex.line.fill.background()
    
    inner_hex = shapes.add_shape(MSO_SHAPE.HEXAGON, left + size * 0.15, top + size * 0.15, size * 0.7, size * 0.7)
    inner_hex.fill.solid()
    inner_hex.fill.fore_color.rgb = colors['secondary']
    inner_hex.fill.transparency = 0.4
    inner_hex.line.fill.background()
    
    core_dot = shapes.add_shape(MSO_SHAPE.OVAL, left + size * 0.35, top + size * 0.35, size * 0.3, size * 0.3)
    core_dot.fill.solid()
    core_dot.fill.fore_color.rgb = RGBColor(255, 255, 255)
    core_dot.fill.transparency = 0.3
    core_dot.line.fill.background()

def create_title_slide(prs):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(8, 12, 22)
    shapes = slide.shapes
    
    large_gradient = shapes.add_shape(MSO_SHAPE.OVAL, Inches(5), Inches(-4), Inches(16), Inches(14))
    large_gradient.fill.solid()
    large_gradient.fill.fore_color.rgb = RGBColor(30, 58, 138)
    large_gradient.fill.transparency = 0.25
    large_gradient.line.fill.background()
    
    for i in range(4):
        hex_size = Inches(1.2 + i * 0.3)
        hex_left = Inches(12.5 - i * 0.5)
        hex_top = Inches(0.5 + i * 0.35)
        create_corporate_logo_element(slide, hex_left, hex_top, hex_size, {'accent': RGBColor(59, 130, 246), 'secondary': RGBColor(139, 92, 246)})
    
    for i in range(6):
        node_left = Inches(0.5 + i * 2.5)
        node_top = Inches(7.2 + (i % 2) * 0.6)
        create_network_node(slide, node_left, node_top, Inches(0.35), {'accent': RGBColor(59, 130, 246), 'secondary': RGBColor(139, 92, 246)})
    
    title_box = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(13), Inches(1.8))
    title_frame = title_box.text_frame
    p = title_frame.add_paragraph()
    p.text = "A2A协议与AI智能体生态"
    p.font.size = Pt(48)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide.shapes.add_textbox(Inches(1.5), Inches(3.6), Inches(13), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    p = subtitle_frame.add_paragraph()
    p.text = "2026年AI技术趋势 · 智能体协作 · 市场展望"
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(139, 92, 246)
    p.alignment = PP_ALIGN.CENTER
    
    info_box = slide.shapes.add_textbox(Inches(3), Inches(5.2), Inches(10), Inches(1.2))
    info_frame = info_box.text_frame
    p = info_frame.add_paragraph()
    p.text = "2026年AI智能体浪潮"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(150, 170, 190)
    p.alignment = PP_ALIGN.CENTER
    
    p = info_frame.add_paragraph()
    p.text = "版本：V14.0 | 更新日期：2026年5月"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(120, 140, 160)
    p.alignment = PP_ALIGN.CENTER
    
    copyright_box = slide.shapes.add_textbox(Inches(0.5), Inches(8.4), Inches(15), Inches(0.5))
    copyright_frame = copyright_box.text_frame
    p = copyright_frame.add_paragraph()
    p.text = "©2026 A2A协议研究团队 | 非商业学习用途"
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(100, 120, 140)
    p.alignment = PP_ALIGN.LEFT
    
    return slide

def parse_table_content(text):
    lines = text.strip().split('\n')
    table_data = []
    in_table = False
    
    for line in lines:
        if line.startswith('|') and line.endswith('|'):
            in_table = True
            cells = [cell.strip() for cell in line.split('|')[1:-1]]
            if cells and not all(c == '-' for c in cells):
                table_data.append(cells)
        elif in_table and not line.startswith('|'):
            break
    
    return table_data if len(table_data) > 1 else None

def parse_code_block(text):
    lines = text.strip().split('\n')
    code_lines = []
    in_code = False
    
    for line in lines:
        if line.startswith('```'):
            in_code = not in_code
            continue
        if in_code:
            code_lines.append(line)
    
    return '\n'.join(code_lines) if code_lines else None

def create_chapter_slide(prs, chapter):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    colors = CHAPTER_COLORS.get(chapter['num'], CHAPTER_COLORS[1])
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors['bg']
    
    shapes = slide.shapes
    gradient_oval = shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(-2), Inches(12), Inches(10))
    gradient_oval.fill.solid()
    gradient_oval.fill.fore_color.rgb = colors['accent']
    gradient_oval.fill.transparency = 0.12
    gradient_oval.line.fill.background()
    
    bottom_bar = shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(8.5), Inches(16), Inches(0.5))
    bottom_bar.fill.solid()
    bottom_bar.fill.fore_color.rgb = colors['accent']
    bottom_bar.fill.transparency = 0.2
    bottom_bar.line.fill.background()
    
    create_corporate_logo_element(slide, Inches(13.5), Inches(0.3), Inches(1.2), colors)
    create_network_node(slide, Inches(13.8), Inches(2), Inches(0.4), colors)
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(12), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.add_paragraph()
    p.text = f"项目{chapter['num']}: {chapter['title']}"
    p.font.size = Pt(24)
    p.font.color.rgb = colors['accent']
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    y_pos = Inches(1.2)
    max_y = Inches(7.5)
    
    for section in chapter['sections']:
        if y_pos >= max_y:
            break
        
        table_data = parse_table_content(section)
        if table_data and len(table_data) > 1:
            create_table(slide, table_data, colors, y_pos)
            y_pos += Inches(4.5)
            continue
        
        code_content = parse_code_block(section)
        if code_content:
            code_box = slide.shapes.add_textbox(Inches(1), y_pos, Inches(14), Inches(3))
            code_frame = code_box.text_frame
            code_frame.word_wrap = True
            p = code_frame.add_paragraph()
            p.text = code_content[:800]
            p.font.size = Pt(9)
            p.font.color.rgb = colors['text']
            code_box.fill.solid()
            code_box.fill.fore_color.rgb = RGBColor(20, 25, 35)
            code_box.fill.transparency = 0.5
            y_pos += Inches(3.2)
            continue
        
        lines = section.strip().split('\n')
        content_text = []
        for line in lines:
            line = line.strip()
            if line.startswith('###'):
                content_text.append(f"【{line[4:]}】")
            elif line.startswith('- ') or line.startswith('* '):
                content_text.append(f"• {line[2:]}")
            elif line.startswith('>'):
                content_text.append(f"「{line[2:]}」")
            elif line and not line.startswith('#') and not line.startswith('**') and line:
                content_text.append(line)
        
        if content_text:
            content_box = slide.shapes.add_textbox(Inches(1), y_pos, Inches(14), Inches(2.5))
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            
            for text in content_text[:8]:
                clean_text = text.replace('🔥', '').replace('⭐', '').replace('🔴', '').replace('🤝', '').replace('🧠', '').replace('📊', '').replace('📍', '')
                p = content_frame.add_paragraph()
                p.text = clean_text[:60]
                p.font.size = Pt(11)
                p.font.color.rgb = colors['text']
            
            y_pos += Inches(2.8)
    
    copyright_box = slide.shapes.add_textbox(Inches(0.5), Inches(8.4), Inches(15), Inches(0.5))
    copyright_frame = copyright_box.text_frame
    p = copyright_frame.add_paragraph()
    p.text = "©2026 A2A协议研究团队 | 非商业学习用途"
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(100, 120, 140)
    p.alignment = PP_ALIGN.LEFT
    
    return slide

def create_table(slide, table_data, colors, top):
    num_rows = len(table_data)
    num_cols = len(table_data[0]) if num_rows > 0 else 3
    
    if num_rows > 10:
        num_rows = 10
    
    left = Inches(0.5)
    width = Inches(15)
    height = Inches(4)
    
    table = slide.shapes.add_table(num_rows, num_cols, left, top, width, height).table
    highlight_keywords = ['核心', '关键', '重要', '必须', '最佳', '提升', '降低', '99', '显著']
    
    for i, row in enumerate(table_data[:num_rows]):
        for j, cell in enumerate(row[:num_cols]):
            cell_text = re.sub(r'[🔥⭐🔴🤝🧠📊📍]', '', str(cell))
            table.cell(i, j).text = cell_text[:20]
            paragraph = table.cell(i, j).text_frame.paragraphs[0]
            
            if num_cols >= 7:
                paragraph.font.size = Pt(6.5)
            elif num_cols >= 5:
                paragraph.font.size = Pt(7.5)
            else:
                paragraph.font.size = Pt(8.5)
            
            paragraph.alignment = PP_ALIGN.CENTER
            
            if i == 0:
                paragraph.font.color.rgb = RGBColor(255, 255, 255)
                paragraph.font.bold = True
                table.cell(i, j).fill.solid()
                table.cell(i, j).fill.fore_color.rgb = colors['accent']
            else:
                is_highlight = any(keyword in str(cell) for keyword in highlight_keywords)
                if is_highlight:
                    paragraph.font.color.rgb = colors['highlight']
                    paragraph.font.bold = True
                else:
                    paragraph.font.color.rgb = colors['text']
                table.cell(i, j).fill.solid()
                table.cell(i, j).fill.fore_color.rgb = RGBColor(20, 25, 35)
                table.cell(i, j).fill.transparency = 0.4

def create_summary_slide(prs):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(8, 12, 22)
    shapes = slide.shapes
    
    large_gradient = shapes.add_shape(MSO_SHAPE.OVAL, Inches(4), Inches(-3), Inches(16), Inches(14))
    large_gradient.fill.solid()
    large_gradient.fill.fore_color.rgb = RGBColor(30, 58, 138)
    large_gradient.fill.transparency = 0.2
    large_gradient.line.fill.background()
    
    for i in range(3):
        hex_left = Inches(0.5 + i * 0.8)
        hex_top = Inches(0.5 + i * 0.4)
        create_corporate_logo_element(slide, hex_left, hex_top, Inches(1), {'accent': RGBColor(59, 130, 246), 'secondary': RGBColor(139, 92, 246)})
    
    title_box = slide.shapes.add_textbox(Inches(2), Inches(0.4), Inches(12), Inches(1))
    title_frame = title_box.text_frame
    p = title_frame.add_paragraph()
    p.text = "🎯 A2A协议与AI智能体核心要点"
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    summary_data = [
        ['项目', '核心内容', '关键价值'],
        ['1', '封面升级', '科技感设计'],
        ['2', '目录优化', '清晰导航'],
        ['3', '协议核心', '标准化通信'],
        ['4', '能力矩阵', '分层架构'],
        ['5', 'AI趋势', '2026新动态'],
        ['6', '应用场景', '四大领域'],
        ['7', '案例分析', '真实验证'],
        ['8', '技术对比', '优势突出'],
        ['9', '市场展望', '投资热点'],
        ['10', '结语', '行动建议']
    ]
    
    table = slide.shapes.add_table(11, 3, Inches(1), Inches(1.4), Inches(14), Inches(6)).table
    
    for i, row in enumerate(summary_data):
        for j, cell in enumerate(row):
            table.cell(i, j).text = cell
            paragraph = table.cell(i, j).text_frame.paragraphs[0]
            paragraph.font.size = Pt(11)
            paragraph.alignment = PP_ALIGN.CENTER
            
            if i == 0:
                paragraph.font.color.rgb = RGBColor(255, 255, 255)
                paragraph.font.bold = True
                table.cell(i, j).fill.solid()
                table.cell(i, j).fill.fore_color.rgb = RGBColor(59, 130, 246)
            else:
                paragraph.font.color.rgb = RGBColor(220, 230, 240)
                table.cell(i, j).fill.solid()
                table.cell(i, j).fill.fore_color.rgb = RGBColor(20, 25, 35)
    
    copyright_box = slide.shapes.add_textbox(Inches(0.5), Inches(8.4), Inches(15), Inches(0.5))
    copyright_frame = copyright_box.text_frame
    p = copyright_frame.add_paragraph()
    p.text = "©2026 A2A协议研究团队 | 非商业学习用途"
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(100, 120, 140)
    p.alignment = PP_ALIGN.LEFT
    
    return slide

def create_pptx_from_md(md_path, output_path):
    content = read_markdown(md_path)
    chapters = parse_markdown(content)
    print(f'解析到 {len(chapters)} 个章节')
    
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    create_title_slide(prs)
    
    for chapter in chapters:
        print(f'生成章节 {chapter["num"]}: {chapter["title"]}')
        create_chapter_slide(prs, chapter)
    
    create_summary_slide(prs)
    
    prs.save(output_path)
    print(f"PPTX文件已生成: {output_path}")

if __name__ == '__main__':
    md_file = r'f:\个人作品\legal-rag-qa-system\A2A_PPT完整升级内容.md'
    output_file = r'f:\个人作品\legal-rag-qa-system\A2A_PROTOCOL_AI_AGENT_V14.pptx'
    create_pptx_from_md(md_file, output_file)