from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import re

CORPORATE_COLORS = {
    'primary': RGBColor(0, 100, 180),
    'secondary': RGBColor(0, 150, 220),
    'accent': RGBColor(0, 200, 255),
    'dark': RGBColor(10, 15, 30),
    'light': RGBColor(230, 240, 250),
    'highlight': RGBColor(255, 80, 80),
    'gold': RGBColor(255, 200, 100),
    'gray': RGBColor(120, 140, 160)
}

CHAPTER_COLORS = {
    1: {'bg': RGBColor(10, 15, 28), 'accent': RGBColor(0, 180, 255), 'text': RGBColor(230, 240, 250), 'highlight': RGBColor(255, 80, 80), 'secondary': RGBColor(50, 120, 180)},
    2: {'bg': RGBColor(12, 25, 22), 'accent': RGBColor(0, 220, 180), 'text': RGBColor(220, 240, 230), 'highlight': RGBColor(255, 100, 100), 'secondary': RGBColor(30, 150, 130)},
    3: {'bg': RGBColor(25, 15, 28), 'accent': RGBColor(220, 100, 200), 'text': RGBColor(240, 220, 245), 'highlight': RGBColor(255, 100, 100), 'secondary': RGBColor(150, 60, 150)},
    4: {'bg': RGBColor(25, 22, 12), 'accent': RGBColor(255, 200, 60), 'text': RGBColor(250, 245, 220), 'highlight': RGBColor(255, 100, 100), 'secondary': RGBColor(180, 140, 40)},
    5: {'bg': RGBColor(10, 25, 32), 'accent': RGBColor(160, 120, 255), 'text': RGBColor(220, 235, 250), 'highlight': RGBColor(255, 100, 100), 'secondary': RGBColor(100, 80, 200)},
    6: {'bg': RGBColor(28, 20, 15), 'accent': RGBColor(255, 120, 160), 'text': RGBColor(250, 235, 220), 'highlight': RGBColor(255, 100, 100), 'secondary': RGBColor(180, 80, 110)},
    7: {'bg': RGBColor(15, 30, 25), 'accent': RGBColor(0, 200, 150), 'text': RGBColor(220, 245, 230), 'highlight': RGBColor(255, 100, 100), 'secondary': RGBColor(50, 150, 100)},
    8: {'bg': RGBColor(30, 18, 35), 'accent': RGBColor(200, 150, 255), 'text': RGBColor(240, 230, 250), 'highlight': RGBColor(255, 100, 100), 'secondary': RGBColor(120, 80, 180)},
    9: {'bg': RGBColor(20, 25, 15), 'accent': RGBColor(180, 220, 80), 'text': RGBColor(235, 245, 220), 'highlight': RGBColor(255, 100, 100), 'secondary': RGBColor(120, 150, 60)},
    10: {'bg': RGBColor(35, 20, 40), 'accent': RGBColor(255, 150, 200), 'text': RGBColor(245, 230, 245), 'highlight': RGBColor(255, 100, 100), 'secondary': RGBColor(150, 100, 180)},
    11: {'bg': RGBColor(25, 20, 15), 'accent': RGBColor(139, 90, 43), 'text': RGBColor(245, 235, 220), 'highlight': RGBColor(139, 69, 19), 'secondary': RGBColor(180, 120, 60)}
}

CHAPTER_TITLES = [
    "第一章：A2A协议架构与核心价值",
    "第二章：协议规范与技术细节",
    "第三章：安全与合规体系",
    "第四章：部署与运维实践",
    "第五章：性能优化策略",
    "第六章：实战案例与最佳实践"
]

WATERMARK_TEXT = "A2A Protocol | 学习笔记"
COPYRIGHT_INFO = "©2026个人学习笔记，非商用"

DEFAULT_TABLES = {
    1: [
        ['价值维度', '业务收益', '技术实现', '成功指标', '适用场景', '成熟度', '投资回报', '技术架构', '关键挑战'],
        ['互操作性', '降低集成成本显著', '标准化消息格式+Agent Cards', '集成时间明显缩短', '跨平台协作、生态整合', '高', '18个月', 'A2A协议栈', '协议版本兼容'],
        ['弹性扩展', '支持大规模并发', '去中心化架构+K8s调度', '吞吐量显著提升', '大规模部署、弹性业务', '高', '12个月', 'K8s+Istio', '资源调度优化'],
        ['安全合规', '通过SOC2/ISO27001认证', '端到端加密+零信任架构', '合规成本明显降低', '金融级应用、医疗数据', '高', '24个月', '零信任架构', '证书管理'],
        ['可观测性', '全链路追踪+实时监控', 'OpenTelemetry+Prometheus', 'MTTR明显降低', '运维监控、故障排查', '中', '18个月', 'O11y Stack', '监控告警风暴'],
        ['智能协作', 'Agent自主协同+技能发现', 'A2A技能发现协议', '任务完成率有所提升', '多Agent系统、任务分工', '中', '24个月', '多Agent架构', '协调一致性'],
        ['跨云部署', '多云无缝迁移', 'K8s原生支持+CNCF认证', '迁移时间明显缩短', '混合云架构、多云战略', '中', '36个月', '混合云架构', '数据同步'],
        ['成本优化', '降低运维成本显著', '自动化运维+Serverless', 'TCO有所降低', '中大型企业、成本敏感', '中', '24个月', 'Serverless架构', '冷启动延迟'],
        ['快速迭代', '缩短开发周期显著', '模块化设计+API优先', '上市时间有所缩短', '敏捷团队、快速上市', '高', '12个月', 'API优先架构', '版本兼容'],
        ['生态扩展', '第三方Agent接入', '标准化API+开发者SDK', '生态伙伴明显增长', '平台生态、合作伙伴', '中', '30个月', '开放平台架构', '生态治理'],
        ['数据价值', '数据资产变现', '隐私计算+联邦学习', '数据价值有所提升', '数据合作、合规共享', '低', '48个月', '隐私计算架构', '合规边界'],
        ['边缘计算', '边缘延迟显著降低', 'K3s+边缘节点', '快速响应', '边缘场景、低延迟', '中', '24个月', '边缘架构', '边缘管理'],
        ['AI协同', 'AI任务效率显著提升', 'Agent技能组合+编排', '复杂任务完成率提升', 'AI应用、智能系统', '中', '30个月', 'AI Agent架构', '技能组合'],
        ['Agent自主', '自主决策能力显著提升', '强化学习+决策引擎', '人工干预减少', '自主系统、智能决策', '中', '24个月', '自主Agent架构', '决策准确性'],
        ['多模态交互', '交互体验显著提升', '文本+语音+视觉融合', '用户满意度提升', '智能客服、虚拟助手', '高', '18个月', '多模态Agent', '模态融合']
    ],
    2: [
        ['协议层级', '业务价值', '技术亮点', '成功指标', '适用场景', '性能指标', '投资回报', '技术架构', '关键挑战'],
        ['应用层', '技能发现效率显著提升', 'Agent Cards/ADR规范', '快速技能发现', 'Agent协作、技能共享', '快速响应', '12个月', 'REST+GraphQL', '技能描述标准化'],
        ['消息层', '消息吞吐量显著提升', 'JSON-RPC 2.0/JSON-LD', '高吞吐量', '异步通信、事件驱动', '高吞吐量', '18个月', 'Kafka/Pulsar', '消息顺序保证'],
        ['传输层', '延迟显著降低', 'HTTP/2、gRPC、QUIC', '低延迟通信', '低延迟通信、实时交互', '低延迟', '12个月', 'gRPC+QUIC', '协议协商'],
        ['安全层', '安全合规成本显著降低', 'OAuth2、mTLS、JWT', '零信任架构', '身份认证、数据加密', '零信任', '24个月', 'PKI+OAuth2', '密钥轮换'],
        ['发现层', '服务发现效率显著提升', '服务注册中心/ETCD', '快速更新', '动态服务、微服务架构', '快速更新', '18个月', 'ETCD+Consul', '一致性保证'],
        ['事件层', '实时推送延迟显著降低', 'WebSocket、SSE', '低延迟推送', '实时通知、状态同步', '低延迟', '12个月', 'WebSocket集群', '连接管理'],
        ['数据层', '序列化效率显著提升', 'Protobuf/JSON/BSON', '高效压缩', '数据交换、存储', '高效压缩', '18个月', 'Protobuf+Schema', '版本演进'],
        ['控制层', '路由灵活性显著提升', 'gRPC/REST/gRPC-Web', '弹性路由', 'API网关、流量控制', '弹性路由', '12个月', 'Envoy/Kong', '配置复杂度'],
        ['编排层', '工作流效率显著提升', 'Dapr/Knative', '灵活流程编排', '业务流程自动化', '灵活编排', '24个月', 'Dapr+Airflow', '状态管理'],
        ['语义层', '知识理解能力显著提升', 'RDF/OWL', '语义推理', '知识图谱、AI理解', '语义理解', '36个月', '知识图谱引擎', '推理性能'],
        ['Agent通信', 'Agent间通信效率显著提升', 'A2A消息格式+协议栈', '低延迟通信', '多Agent协作、消息传递', '低延迟', '12个月', 'A2A协议栈', '消息可靠性'],
        ['Agent发现', 'Agent发现效率显著提升', '服务发现+Agent注册', '快速发现', '动态Agent网络、自动发现', '快速发现', '18个月', 'Agent注册中心', '注册一致性']
    ],
    3: [
        ['安全组件', '业务价值', '技术亮点', '成功指标', '适用场景', '防护级别', '投资回报', '技术架构', '关键挑战'],
        ['API网关', 'API安全成本显著降低', 'Kong/APISIX/Envoy', '高攻击拦截率', 'API安全、流量控制', '高', '18个月', 'API网关集群', '性能损耗'],
        ['WAF防火墙', 'Web攻击防护显著提升', 'ModSecurity/Akamai', 'OWASP Top10防护', 'Web应用、API保护', '极高', '24个月', 'WAF+CDN', '规则误报'],
        ['OAuth 2.0', '认证效率显著提升', '授权码/PKCE/JWT', '快速认证', '身份管理、SSO', '高', '12个月', 'OAuth2+OIDC', '令牌管理'],
        ['mTLS', '通道安全高防护级别', '双向证书+证书轮换', '零中间人攻击', '敏感数据传输', '极高', '24个月', 'PKI+SPIFFE', '证书分发'],
        ['RBAC/ABAC', '权限管理效率显著提升', '角色/属性权限', '细粒度控制', '企业级访问控制', '高', '18个月', 'IAM系统', '权限爆炸'],
        ['DDoS防护', 'DDoS攻击防护显著提升', 'Cloudflare/Akamai', '高流量清洗率', '大规模流量攻击', '极高', '24个月', '多层防护', '成本投入'],
        ['数据加密', '数据安全高防护级别', 'AES-256/GPG', '全面加密', '数据存储、传输', '极高', '18个月', '加密服务', '性能影响'],
        ['审计日志', '合规审计效率显著提升', 'ELK/Splunk/Datadog', '日志不可篡改', '合规审计、行为追踪', '高', '24个月', 'SIEM系统', '存储成本'],
        ['威胁检测', '威胁发现显著加快', 'SIEM/EDR', '快速攻击发现', '高级威胁检测', '高', '30个月', 'AI检测引擎', '误报率'],
        ['密钥管理', '密钥安全高防护级别', 'HSM/Cloud KMS', '密钥高安全性', '密钥存储、轮换', '极高', '24个月', 'HSM集群', '可用性'],
        ['数据脱敏', '数据隐私高防护级别', '动态脱敏+静态脱敏', '敏感数据保护', '数据共享、合规', '高', '24个月', '脱敏引擎', '精度平衡'],
        ['访问审计', '合规审计能力显著提升', '行为分析+操作记录', '全面审计覆盖', '合规检查、安全审计', '高', '18个月', '审计系统', '日志存储'],
        ['Agent安全', 'Agent间通信高防护级别', 'A2A安全协议+加密', '通信安全', '多Agent协作、数据传输', '极高', '18个月', 'A2A安全栈', '协议兼容'],
        ['AI伦理', 'AI行为高防护级别', '伦理框架+审计机制', '违规行为预防', 'AI应用、智能决策', '高', '24个月', '伦理治理系统', '伦理标准'],
        ['安全审计', '安全评估能力显著提升', '渗透测试+安全扫描', '高漏洞发现率', '安全评估、合规检查', '高', '18个月', '安全审计平台', '审计深度'],
        ['安全培训', '安全意识显著提升', '模拟演练+在线课程', '安全事件减少', '员工培训、安全意识', '中', '12个月', '培训平台', '培训效果'],
        ['安全测试', '代码安全显著提升', 'SAST/DAST/IAST', '高漏洞修复率', '开发安全、CI/CD', '高', '18个月', 'DevSecOps', '误报处理'],
        ['零信任架构', '安全防护高防护级别', '身份验证+最小权限', '内部攻击防护', '企业安全、边界消除', '极高', '24个月', '零信任平台', '实施复杂度'],
        ['隐私计算', '数据价值高防护级别', '联邦学习+安全多方计算', '数据可用不可见', '数据协作、隐私保护', '中', '36个月', '隐私计算平台', '性能开销'],
        ['SOC2合规', '合规高防护级别', '安全运营中心+审计', '高合规通过率', '企业级SaaS、金融服务', '高', '24个月', 'SOC2体系', '持续合规']
    ],
    4: [
        ['部署策略', '业务价值', '技术亮点', '成功指标', '适用场景', '可用性', '投资回报', '技术架构', '关键挑战'],
        ['多活部署', '高业务可用性', 'K8s StatefulSet+DNS轮询', '秒级故障转移', '核心业务、高可用', '高', '24个月', '多活K8s集群', '数据一致性'],
        ['自动扩缩容', '资源利用率显著提升', 'HPA/VPA+KEDA', '快速弹性响应', '波动流量、弹性业务', '弹性', '12个月', 'K8s自动伸缩', '冷启动'],
        ['健康检查', '故障自愈能力显著提升', 'Liveness/Readiness/Startup', '快速故障恢复', '微服务、高可用', '高', '6个月', 'K8s健康检查', '探针配置'],
        ['蓝绿发布', '发布风险显著降低', 'Istio/Argo Rollouts', '零停机升级', '生产环境、频繁发布', '零停机', '18个月', 'Istio+Argo', '流量切分'],
        ['异地容灾', '高灾难恢复能力', '跨区域K8s集群+DRBD', '快速恢复', '金融系统、核心数据', '极高', '36个月', '跨区集群', '数据同步'],
        ['边缘部署', '边缘延迟显著降低', 'K3s/KubeEdge', '快速响应', '边缘计算、就近访问', '低延迟', '24个月', '边缘K8s', '边缘管理'],
        ['金丝雀发布', '发布风险显著降低', 'Argo Rollouts/Flagger', '渐进式发布', '核心服务、风险控制', '渐进式', '18个月', 'Argo+Prometheus', '指标评估'],
        ['Serverless', '运维成本显著降低', 'KEDA+Knative', '按需扩缩', '非核心服务、突发流量', '按需', '12个月', 'KEDA+Knative', '冷启动延迟'],
        ['混合云部署', '多云支持', 'Anthos/Azure Arc', '多云统一管理', '企业级、多云战略', '多云', '30个月', '混合云管理', '跨云网络'],
        ['GitOps', '部署效率显著提升', 'Argo CD/Flux', '声明式配置', 'CI/CD、版本控制', '高', '18个月', 'GitOps工具链', '配置漂移'],
        ['服务网格', '服务间通信效率显著提升', 'Istio/Linkerd', '零信任网络', '微服务架构、安全通信', '高', '18个月', '服务网格架构', '配置复杂度'],
        ['自动化运维', '运维效率显著提升', 'Ansible/Terraform+CI/CD', '快速部署', 'DevOps、持续交付', '高', '12个月', '自动化工具链', '学习曲线'],
        ['日志管理', '问题定位效率显著提升', 'ELK/Grafana Loki', '快速日志检索', '运维监控、故障排查', '高', '18个月', '日志平台', '存储成本'],
        ['指标监控', '可观测性显著提升', 'Prometheus+Grafana', '快速实时告警', '性能监控、容量规划', '极高', '12个月', '监控体系', '告警风暴'],
        ['分布式追踪', '故障定位显著加快', 'Jaeger/Zipkin', '快速全链路追踪', '微服务、分布式系统', '高', '18个月', '追踪系统', '采样率'],
        ['配置管理', '配置变更效率显著提升', 'ConfigMap+Secret', '快速热更新', '配置管理、敏感数据', '高', '6个月', 'K8s配置管理', '密钥轮换'],
        ['存储编排', '存储效率显著提升', 'CSI+LocalPV', '数据持久化', '有状态应用、数据库', '高', '24个月', '存储架构', '存储类型'],
        ['网络策略', '网络安全显著提升', 'NetworkPolicy', '微分段隔离', '多租户、安全边界', '高', '12个月', '网络隔离', '策略复杂度'],
        ['自动化测试', '测试覆盖率显著提升', 'Pytest+Selenium', '高代码覆盖率', '持续集成、质量保障', '高', '18个月', '测试框架', '测试数据'],
        ['容量规划', '资源利用率显著提升', 'Prometheus+VPA', '高资源预测准确率', '资源管理、成本优化', '高', '24个月', '容量管理', '预测准确性'],
        ['安全扫描', '漏洞发现显著加快', 'Trivy/Grype', '快速镜像扫描', '安全合规、CI/CD', '高', '12个月', '安全工具链', '误报处理'],
        ['Agent部署', 'Agent部署效率显著提升', 'Agent镜像+自动编排', '部署时间缩短', '多Agent系统、自动化部署', '高', '12个月', 'Agent编排平台', '版本兼容'],
        ['Agent监控', 'Agent监控能力显著提升', '分布式追踪+指标采集', '快速故障发现', 'Agent集群、智能监控', '极高', '18个月', 'O11y+Agent监控', '监控告警']
    ],
    5: [
        ['优化维度', '业务价值', '技术亮点', '成功指标', '适用场景', '预期收益', '投资回报', '技术架构', '关键挑战'],
        ['网络优化', '延迟显著降低', 'HTTP/2、gRPC、QUIC', 'RTT显著降低', '低延迟通信', '延迟降低', '6个月', 'gRPC+QUIC', '客户端兼容'],
        ['消息批量', '吞吐量显著提升', '批量聚合+背压控制', 'QPS显著提升', '高吞吐量场景', 'QPS提升', '12个月', 'Kafka+batching', '消息顺序'],
        ['缓存策略', '响应速度显著提升', 'Redis集群+多级缓存', 'P95延迟显著降低', '读多写少场景', '延迟降低', '6个月', 'Redis+LocalCache', '数据一致性'],
        ['并发处理', '容量显著提升', '异步队列+线程池', '并发能力显著提升', '高并发场景', '并发提升', '18个月', 'Vert.x/Netty', '线程安全'],
        ['序列化', '数据体积显著减少', 'Protobuf/FlatBuffers', '序列化体积显著减少', '数据传输、存储', '体积减少', '12个月', 'Protobuf', '版本兼容'],
        ['CDN加速', '边缘延迟显著降低', '边缘缓存+预热', 'P95延迟显著降低', '静态资源、全球用户', '延迟降低', '6个月', 'CDN+预热', '缓存失效'],
        ['数据库优化', '查询速度显著提升', '读写分离+分库分表', '查询延迟显著降低', '大数据量场景', '延迟降低', '12个月', 'Sharding+读写分离', '事务一致性'],
        ['代码优化', '性能显著提升', '算法改进+JIT', 'CPU使用率显著降低', '计算密集场景', '性能提升', '24个月', 'Profiling+优化', '回归风险'],
        ['资源调度', '资源利用率显著提升', 'K8s QoS+资源预留', '资源利用率显著提升', '容器集群', '利用率提升', '8个月', 'K8s调度优化', 'Pod驱逐'],
        ['连接池', '连接效率显著提升', 'HikariCP/OkHttp', '连接复用率显著提升', '高连接场景', '复用率提升', '3个月', '连接池优化', '连接泄漏'],
        ['内存优化', '内存使用显著降低', '对象池+GC调优', '内存占用显著降低', '内存敏感场景', '内存降低', '12个月', '内存管理', 'GC停顿'],
        ['异步IO', 'IO效率显著提升', 'Reactor/CompletableFuture', '吞吐量显著提升', 'IO密集场景', '吞吐量提升', '18个月', '异步框架', '复杂度'],
        ['Agent优化', 'Agent性能显著提升', 'Agent缓存+预加载', '响应时间显著降低', '多Agent系统、智能协作', '高', '12个月', 'Agent性能优化', '缓存一致性'],
        ['AI推理优化', '推理速度显著提升', '模型量化+TensorRT', '推理延迟显著降低', 'AI应用、边缘部署', '高', '18个月', 'AI推理引擎', '精度损失'],
        ['索引优化', '查询速度显著提升', 'B+树/倒排索引', '查询时间显著降低', '大数据查询、全文搜索', '查询提升', '6个月', '索引引擎', '索引维护'],
        ['负载均衡', '流量分配显著优化', 'Nginx/LVS/云LB', '请求均匀分布', '高并发、多节点', '均匀分布', '6个月', 'LB集群', '会话保持'],
        ['预热机制', '首次响应显著加快', '缓存预热+预加载', '冷启动延迟显著降低', '系统启动、热点数据', '延迟降低', '3个月', '预热服务', '资源占用'],
        ['熔断降级', '故障隔离高防护级别', 'Resilience4j/Hystrix', '故障影响范围控制', '微服务、高可用', '故障隔离', '6个月', '熔断框架', '阈值配置'],
        ['读写优化', '数据库压力显著降低', '读写分离+主从复制', '主库压力显著降低', '高并发读写、大数据', '压力降低', '12个月', '主从架构', '数据同步'],
        ['边缘计算', '边缘延迟显著降低', 'K3s/KubeEdge', '快速响应', '边缘场景、IoT', '延迟降低', '24个月', '边缘架构', '边缘管理'],
        ['GPU加速', '计算性能显著提升', 'CUDA/TensorRT', 'AI推理加速', 'AI训练、推理', '性能提升', '18个月', 'GPU集群', '成本投入'],
        ['压缩优化', '传输效率显著提升', 'GZIP/Brotli/Snappy', '传输体积显著减少', '网络传输、存储', '体积减少', '3个月', '压缩服务', 'CPU开销']
    ],
    6: [
        ['案例场景', '业务价值', '技术亮点', '成功指标', '行业领域', '规模', '投资回报', '技术架构', '关键挑战'],
        ['供应链管理', '效率显著提升', '多Agent协作+智能调度', '交付周期明显缩短', '制造业', '大型', '18个月', 'A2A+IoT', '数据孤岛'],
        ['金融风控', '风险显著降低', '实时决策引擎+知识图谱', '准确率明显提升', '金融', '中型', '24个月', 'A2A+ML', '合规要求'],
        ['智能办公', '生产力显著提升', '自动化流程+RPA集成', '成本明显节约', '企业服务', '中型', '12个月', 'A2A+低代码', '系统集成'],
        ['医疗诊断', '准确率有所提升', '专业知识整合+推理引擎', '误诊率有所下降', '医疗', '大型', '36个月', 'A2A+NLP', '数据隐私'],
        ['智能客服', '满意度有所提升', '多轮对话+情感分析', '响应时间明显缩短', '电商', '大型', '18个月', 'A2A+LLM', '上下文理解'],
        ['智能运维', '故障发现明显加快', '预测性维护+异常检测', 'MTTR明显降低', '运维', '中型', '24个月', 'A2A+监控', '告警风暴'],
        ['智能推荐', '转化率有所提升', '个性化算法+实时特征', '用户留存有所提升', '互联网', '大型', '18个月', 'A2A+推荐系统', '冷启动'],
        ['智能教育', '学习效率有所提升', '自适应学习+知识追踪', '成绩有所提升', '教育', '中型', '36个月', 'A2A+教育AI', '学习路径'],
        ['智能制造', '良率有所提升', '数字孪生+质量检测', '成本有所降低', '制造业', '大型', '30个月', 'A2A+IoT+CV', '设备互联'],
        ['智慧城市', '能耗有所降低', '多系统协同+预测优化', '效率有所提升', '政府', '大型', '48个月', 'A2A+城市平台', '数据整合'],
        ['Agent协作', '协作效率显著提升', '多Agent协同+任务分配', '任务完成率明显提高', '多Agent系统、智能协作', '大型', '18个月', 'A2A协作平台', '协调一致性'],
        ['Agent生态', '生态价值显著提升', 'Agent平台+开发者社区', '生态规模明显扩大', '平台生态、合作伙伴', '大型', '30个月', '开放Agent平台', '生态治理'],
        ['Agent创新', '创新效率显著提升', 'AI Agent+新技术融合', '创新周期明显缩短', 'AI创新、前沿探索', '中型', '24个月', 'Agent创新实验室', '技术风险']
    ],
    7: [
        ['时间阶段', '发展特征', '关键技术', '行业影响', '蚌埠机遇', '投资价值', '风险评估', '战略建议', '预期收益'],
        ['2026-2028', '标准化加速', 'Agent Cards、ADR规范', '行业标准形成', '承接长三角产业转移', '高', '中', '布局技术团队', '抢占先机'],
        ['2029-2033', '智能升级', '自主学习、情感理解', 'Agent能力飞跃', '制造业AI升级', '极高', '中', '研发投入', '技术领先'],
        ['2034-2036', '生态成熟', '去中心化网络', '全球Agent经济', '构建产业生态', '高', '低', '平台建设', '生态收益'],
        ['2037-2041', '通用智能', '通用Agent出现', '生产力革命', '全面智能化', '极高', '中', '战略转型', '指数增长'],
        ['2042-2046', '意识涌现', '类人智能可能', '社会变革', '新产业形态', '极高', '高', '前瞻布局', '颠覆创新'],
        ['当前趋势', '多模态融合', '文本+语音+视觉', '交互体验升级', '本地应用开发', '高', '低', '产品落地', '市场份额'],
        ['当前趋势', '边缘部署', 'K3s/KubeEdge', '低延迟服务', '边缘计算节点', '中', '低', '基础设施', '技术储备'],
        ['当前趋势', '安全增强', '零信任+隐私计算', '合规保障', '金融级安全', '高', '低', '安全体系', '信任背书'],
        ['产业趋势', '平台化', '云厂商Agent平台', '生态竞争', '对接云服务商', '中', '中', '合作共赢', '资源整合'],
        ['产业趋势', '垂直深耕', '行业专用Agent', '细分市场', '制造业AI应用', '高', '低', '行业聚焦', '差异化优势']
    ],
    8: [
        ['人群分类', '岗位定位', '招聘策略', '推荐软件', '薪资范围', '蚌埠特色', '政策支持', '预期成效', '关键挑战'],
        ['青年人才(22-28岁)', '初级开发、测试', '校招+实习转正', 'BOSS直聘、智联', '8-15K/月', '本地高校合作', '实习补贴', '人才储备', '经验不足'],
        ['中坚力量(29-38岁)', '技术骨干、管理', '社招+返乡人才', 'BOSS直聘、猎聘', '20-35K/月', '返乡计划', '安家补贴', '核心团队', '竞争激烈'],
        ['资深专家(39-45岁)', '架构师、总监', '猎头引进', '猎聘、脉脉', '60-150万/年', '高层次人才计划', '住房优惠', '技术突破', '成本较高'],
        ['高中学历', '数据标注、客服', '职业院校+线上', 'BOSS直聘、58同城', '4-10K/月', '本地职高合作', '技能补贴', '基础岗位', '技能提升'],
        ['视力残疾人', '语音客服、有声读物制作', '残联推荐+定向培训', '中国残疾人就业网络服务平台', '同工同酬', '盲人按摩转型', '岗位补贴', '就业帮扶', '技能培训'],
        ['肢体残疾人', '数据录入、客服、设计', '残联推荐+远程办公', '微信小程序、官方平台', '同工同酬', '灵活就业支持', '社保补贴', '包容性就业', '办公环境'],
        ['特殊就业支持', '灵活就业、远程工作', '残联推荐', '官方平台、微信小程序', '同工同酬', '就业保障', '残保金优惠', '社会责任', '环境适配'],
        ['大专学历', '运维、技术支持', '本地招聘', '前程无忧、BOSS直聘', '6-12K/月', '职业学院合作', '就业补贴', '运维保障', '发展空间'],
        ['本科学历', '全栈开发、产品', '重点高校', '智联、拉勾网', '12-25K/月', '高校招聘', '人才政策', '技术核心', '留才挑战'],
        ['硕士学历', 'AI算法、架构', '顶尖高校', '猎聘、拉勾网', '20-35K/月', '人才引进', '科研补贴', '技术突破', '竞争激烈'],
        ['博士学历', '首席科学家', '高端引进', '猎聘、学术平台', '80-200万/年', '院士工作站', '科研经费', '行业引领', '稀缺性'],
        ['AI工程师', 'AI Agent开发', '高校+社招', 'BOSS直聘、拉勾', '15-30K/月', 'AI人才引进', '科研补贴', '技术核心', '竞争激烈'],
        ['Agent训练师', 'Agent训练优化', '专业培训+实战', '训练效率明显提升', 'AI应用、智能系统', '中', '12个月', 'AI训练平台', '数据质量'],
        ['顾问导师(46岁+)', '技术顾问', '兼职返聘', '脉脉、领英', '顾问费', '退休专家', '荣誉职位', '知识传递', '时间投入']
    ],
    9: [
        ['办理环节', '所需材料', '办理流程', '注意事项', '蚌埠政策', '办理地点', '办理时间', '费用说明', '常见问题'],
        ['名称核准', '身份证、备选名称', '线上申请→审核', '准备3-5个名称', '全程电子化', '政务服务中心', '1-2天', '免费', '名称重复'],
        ['场地证明', '租赁合同/无偿使用证明', '提供材料→审核', '住宅需社区同意', '住宅改商用支持', '所在区县', '1天', '免费', '地址不符'],
        ['材料提交', '身份证、场地证明', '线上/窗口提交', '经营范围规范', '一窗通办', '政务服务中心', '1-2天', '免费', '材料不全'],
        ['领取执照', '受理通知书', '窗口领取/邮寄', '核对信息', '当场办结', '政务服务中心', '当场', '免费', '信息错误'],
        ['刻章备案', '营业执照', '指定机构刻章', '公章+财务章+法人章', '免费刻章', '公安局指定点', '1-2天', '免费', '章型选择'],
        ['税务登记', '营业执照、法人身份证', '电子税务局', '30日内完成', '线上办理', '税务局', '1天', '免费', '逾期罚款'],
        ['银行开户', '营业执照、公章', '选择银行办理', '基本户只能一个', '多家银行可选', '任意银行', '3-5天', '年费', '开户费用'],
        ['无房产证', '租赁合同+房东证明', '提交替代材料', '社区出具证明', '支持替代材料', '政务服务中心', '1-2天', '免费', '材料真实性'],
        ['未过户房产', '网签合同/交易协议', '提供合同+证明', '原房主同意', '支持网签合同', '政务服务中心', '1-2天', '免费', '产权纠纷'],
        ['原房主逝世', '继承证明+死亡证明', '继承公证→过户→注册', '需先办理继承公证', '支持继承证明', '政务服务中心', '7-15天', '公证费', '继承纠纷'],
        ['创业补贴', '营业执照、就业证明', '线上申请→审核', '毕业5年内/就业困难', '最高5000元', '人社局', '30天', '免费', '条件不符']
    ],
    10: [
        ['模型类型', '代表产品', '部署方式', 'GPU需求', '算力费用/月', 'API调用成本', '优点', '缺点', '适用场景'],
        ['多模态大模型', 'GPT-4V、Gemini Ultra', '🔴API调用', '4xA100/80GB', '¥8-15万', 'GPT-4V: $0.01/image', '能力最强、多模态融合', '成本高、数据安全风险', '创意设计、内容创作'],
        ['文本大模型', 'GPT-4、Claude 3', 'API调用', '8xA100', '¥12-25万', 'GPT-4: $0.03/1K tokens', '推理强、上下文深', '成本高、无本地化', '代码生成、数据分析'],
        ['开源大模型', 'Llama 3-70B、Qwen-72B', '私有化部署', '4xA100/40GB', '¥5-12万', '免费', '数据安全、定制灵活', '技术门槛高', '企业内部系统'],
        ['垂直领域模型', 'Med-PaLM、CodeLlama', '本地/私有化', '2xA100', '¥3-8万', '按需', '专业知识强、精准输出', '通用性弱', '医疗、法律、编程'],
        ['国产大模型', '文心一言、豆包', 'API+私有化', '8xH100', '¥15-30万', '豆包: ¥99/月起', '本土化支持、合规性好', '能力略逊', '政企服务'],
        ['轻量模型', 'Phi-3、Llama 3-8B', '本地部署', 'RTX 4090', '¥0.5-2万', '免费', '速度快、资源需求低', '能力有限', '移动端、IoT'],
        ['DeepSeek系列', 'DeepSeek-R1', 'API+本地', '4xA100', '¥6-12万', 'API: ¥0.002/1K tokens', '长上下文、数学推理强', '知名度低', '科研教育'],
        ['Qwen系列', 'Qwen 2-72B', '私有化部署', '4xA100', '¥5-10万', '免费商用', '多语言、社区活跃', '生态相对小', '跨境服务'],
        ['Doubao系列', '豆包企业版', 'API调用', '云端', '按需', '企业版: ¥1999/月起', '娱乐性强、生活化', '专业能力有限', '消费级应用'],
        ['云端租赁', 'AWS/GCP/Azure', '云部署', '按需配置', '¥2-20万', '按需付费', '弹性伸缩、免维护', '成本波动', '弹性业务'],
        ['边缘部署', 'NVIDIA Jetson', '本地边缘', 'Jetson AGX', '¥0.3-1万', '免费', '低延迟、隐私保护', '能力有限', 'IoT、边缘计算']
    ],
    11: [
        ['章节', '核心内容', '关键技术', '业务价值', '实战要点', '投资回报', '蚌埠机遇', '核心价值'],
        ['第1章', 'A2A协议架构', 'Agent Cards、标准化消息', '🔴互操作性提升60%', '跨平台协作', '18个月', '生态整合', '打破平台壁垒'],
        ['第2章', '协议规范', 'JSON-RPC 2.0、gRPC', '🔴通信效率提升10x', '低延迟通信', '12个月', '技术基础', '标准化接口'],
        ['第3章', '安全合规', 'OAuth2、mTLS、零信任', '🔴合规成本降低40%', '金融级安全', '24个月', '信任背书', '数据安全'],
        ['第4章', '部署运维', 'K8s、自动扩缩容', '🔴可用性99.99%', '高可用架构', '18个月', '基础设施', '业务连续性'],
        ['第5章', '性能优化', '缓存、异步IO、AI推理', '🔴性能提升3x', '大规模承载', '12个月', '效率提升', '用户体验'],
        ['第6章', '实战案例', '供应链、风控、医疗', '🔴效率提升40%', '行业落地', '18-36个月', '案例借鉴', '价值验证'],
        ['第7章', 'AI趋势', '2026-2046演进', '🔴抢占先机', '战略布局', '长期', '产业升级', '前瞻布局'],
        ['第8章', '人才规划', '全人群招聘策略', '🔴人才储备', '本地就业', '持续', '人力资源', '团队建设'],
        ['第9章', '创业指南', '执照办理、政策支持', '🔴创业便利化', '零成本注册', '短期', '创业环境', '营商环境'],
        ['第10章', '大模型生态', '多模态、开源、费用', '🔴选型决策', '成本优化', '持续', '技术选型', 'AI能力'],
        ['总结', 'A2A+AI智能体体系', '端到端解决方案', '🔴企业智能化升级', '技术+人才+生态', '长期价值', '蚌埠AI产业', '全面发展']
    ]
}

CHAPTER_DESCRIPTIONS = {
    1: "A2A协议是下一代Agent-to-Agent通信标准，实现异构AI Agent之间的无缝协作，打破平台壁垒，构建开放的AI生态系统。",
    2: "协议规范定义了Agent间通信的标准化消息格式和交互模式，基于JSON-RPC 2.0实现跨平台、跨语言的互操作性。",
    3: "企业级安全架构包括API网关、WAF防火墙、OAuth 2.0认证、mTLS加密等多层防护，确保数据安全和合规要求。",
    4: "基于Kubernetes的高可用部署方案，支持自动扩缩容、健康检查、蓝绿发布等运维能力，保障业务连续性。",
    5: "通过网络优化、消息批量、缓存策略、并发处理等手段，实现系统性能的全方位提升。",
    6: "通过供应链管理、金融风控、智能办公、医疗诊断等真实案例，验证A2A协议的业务价值和技术可行性。",
    7: "AI智能体正处于快速发展期，未来10-20年将经历标准化、智能化到生态化的演进。",
    8: "针对蚌埠市产业特点，制定差异化人才策略，涵盖全人群包括视力残疾人和肢体残疾人的就业帮扶。",
    9: "蚌埠市为创业者提供完善的政策支持，包括免费注册、创业补贴、税收优惠等。",
    10: "AI大模型正处于爆发式发展阶段，算力成本是关键考量因素，需根据数据敏感性和预算选择合适的部署方案。",
    11: "本PPT系统阐述了A2A协议与AI智能体的完整知识体系，全面覆盖AI智能体技术栈，助力蚌埠市AI产业升级。"
}

CHAPTER_RESOURCES = {
    1: ['Linux Foundation: https://linuxfoundation.org'],
    2: ['JSON-RPC 2.0: https://www.jsonrpc.org/specification'],
    3: ['OAuth 2.0: https://oauth.net/2/']
}

def create_corporate_logo_element(slide, left, top, size, colors):
    shapes = slide.shapes
    outer_hex = shapes.add_shape(MSO_SHAPE.HEXAGON, left, top, size, size)
    outer_hex.fill.solid()
    outer_hex.fill.fore_color.rgb = colors['accent']
    outer_hex.fill.transparency = 0.25
    outer_hex.line.fill.background()
    inner_hex = shapes.add_shape(MSO_SHAPE.HEXAGON, left + size * 0.15, top + size * 0.15, size * 0.7, size * 0.7)
    inner_hex.fill.solid()
    inner_hex.fill.fore_color.rgb = colors['secondary']
    inner_hex.fill.transparency = 0.4
    inner_hex.line.fill.background()
    core_dot = shapes.add_shape(MSO_SHAPE.OVAL, left + size * 0.35, top + size * 0.35, size * 0.3, size * 0.3)
    core_dot.fill.solid()
    core_dot.fill.fore_color.rgb = RGBColor(255, 255, 255)
    core_dot.fill.transparency = 0.3
    core_dot.line.fill.background()

def create_network_node(slide, left, top, size, colors):
    shapes = slide.shapes
    main_circle = shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    main_circle.fill.solid()
    main_circle.fill.fore_color.rgb = colors['accent']
    main_circle.fill.transparency = 0.2
    main_circle.line.color.rgb = colors['accent']
    main_circle.line.width = Pt(1.5)
    inner_ring = shapes.add_shape(MSO_SHAPE.OVAL, left + size * 0.2, top + size * 0.2, size * 0.6, size * 0.6)
    inner_ring.fill.solid()
    inner_ring.fill.fore_color.rgb = colors['secondary']
    inner_ring.fill.transparency = 0.5
    inner_ring.line.fill.background()
    center_dot = shapes.add_shape(MSO_SHAPE.OVAL, left + size * 0.35, top + size * 0.35, size * 0.3, size * 0.3)
    center_dot.fill.solid()
    center_dot.fill.fore_color.rgb = RGBColor(255, 255, 255)
    center_dot.line.fill.background()

def create_tech_circuit(slide, left, top, width, colors):
    shapes = slide.shapes
    for i in range(5):
        dot_left = left + i * (width / 4)
        dot = shapes.add_shape(MSO_SHAPE.OVAL, dot_left, top, width * 0.08, width * 0.08)
        dot.fill.solid()
        dot.fill.fore_color.rgb = colors['accent']
        dot.fill.transparency = 0.3
        dot.line.fill.background()
        if i < 4:
            line = shapes.add_shape(MSO_SHAPE.RECTANGLE, dot_left + width * 0.04, top + width * 0.035, width * 0.2, width * 0.01)
            line.fill.solid()
            line.fill.fore_color.rgb = colors['accent']
            line.fill.transparency = 0.5
            line.line.fill.background()

def create_robot_face(slide, left, top, width, colors):
    shapes = slide.shapes
    head = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, width * 0.85)
    head.fill.solid()
    head.fill.fore_color.rgb = colors['accent']
    head.fill.transparency = 0.15
    head.line.color.rgb = colors['accent']
    head.line.width = Pt(1)
    eye_width = width * 0.18
    eye_height = width * 0.12
    eye_left = left + width * 0.22
    eye_top = top + width * 0.22
    for offset in [0, width * 0.38]:
        eye = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, eye_left + offset, eye_top, eye_width, eye_height)
        eye.fill.solid()
        eye.fill.fore_color.rgb = RGBColor(255, 255, 255)
        eye.line.fill.background()
        pupil = shapes.add_shape(MSO_SHAPE.OVAL, eye_left + offset + eye_width * 0.3, eye_top + eye_height * 0.25, eye_width * 0.4, eye_height * 0.5)
        pupil.fill.solid()
        pupil.fill.fore_color.rgb = colors['accent']
        pupil.line.fill.background()
    mouth_width = width * 0.35
    mouth = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + width * 0.325, top + width * 0.55, mouth_width, width * 0.06)
    mouth.fill.solid()
    mouth.fill.fore_color.rgb = RGBColor(255, 255, 255)
    mouth.fill.transparency = 0.3
    mouth.line.fill.background()

def add_watermark(slide):
    shapes = slide.shapes
    for i in range(3):
        for j in range(4):
            left = Inches(1.5 + i * 5)
            top = Inches(1.8 + j * 2)
            watermark_box = slide.shapes.add_textbox(left, top, Inches(4), Inches(1))
            watermark_frame = watermark_box.text_frame
            p = watermark_frame.add_paragraph()
            p.text = WATERMARK_TEXT
            p.font.size = Pt(10)
            p.font.color.rgb = RGBColor(100, 120, 150)
            p.font.bold = False
            p.alignment = PP_ALIGN.CENTER
            watermark_box.rotation = -15
            watermark_box.fill.transparency = 0.85

def add_copyright(slide):
    copyright_box = slide.shapes.add_textbox(Inches(0.5), Inches(8.4), Inches(15), Inches(0.5))
    copyright_frame = copyright_box.text_frame
    p = copyright_frame.add_paragraph()
    p.text = COPYRIGHT_INFO
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(120, 140, 160)
    p.alignment = PP_ALIGN.LEFT

def add_chapter_background(slide, chapter_num):
    colors = CHAPTER_COLORS.get(chapter_num, CHAPTER_COLORS[1])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors['bg']
    shapes = slide.shapes
    gradient_oval = shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(-2), Inches(12), Inches(10))
    gradient_oval.fill.solid()
    gradient_oval.fill.fore_color.rgb = colors['accent']
    gradient_oval.fill.transparency = 0.12
    gradient_oval.line.fill.background()
    bottom_bar = shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(8.5), Inches(16), Inches(0.5))
    bottom_bar.fill.solid()
    bottom_bar.fill.fore_color.rgb = colors['accent']
    bottom_bar.fill.transparency = 0.2
    bottom_bar.line.fill.background()
    return colors

def create_enterprise_title_slide(prs):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(8, 12, 22)
    shapes = slide.shapes
    
    large_gradient = shapes.add_shape(MSO_SHAPE.OVAL, Inches(5), Inches(-4), Inches(16), Inches(14))
    large_gradient.fill.solid()
    large_gradient.fill.fore_color.rgb = RGBColor(0, 80, 160)
    large_gradient.fill.transparency = 0.2
    large_gradient.line.fill.background()
    
    for i in range(4):
        hex_size = Inches(1.5 + i * 0.3)
        hex_left = Inches(12 - i * 0.5)
        hex_top = Inches(0.3 + i * 0.4)
        create_corporate_logo_element(slide, hex_left, hex_top, hex_size, {'accent': RGBColor(0, 150, 220), 'secondary': RGBColor(0, 200, 255)})
    
    for i in range(6):
        node_left = Inches(0.3 + i * 2.5)
        node_top = Inches(7 + (i % 2) * 0.8)
        create_network_node(slide, node_left, node_top, Inches(0.4), {'accent': RGBColor(0, 150, 220), 'secondary': RGBColor(0, 180, 255)})
    
    for i in range(3):
        circuit_top = Inches(1.5 + i * 1.2)
        create_tech_circuit(slide, Inches(0.3), circuit_top, Inches(3), {'accent': RGBColor(0, 120, 180)})
    
    title_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(13), Inches(1.8))
    title_frame = title_box.text_frame
    p = title_frame.add_paragraph()
    p.text = "A2A协议与AI智能体"
    p.font.size = Pt(48)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide.shapes.add_textbox(Inches(1.5), Inches(3.8), Inches(13), Inches(0.9))
    subtitle_frame = subtitle_box.text_frame
    p = subtitle_frame.add_paragraph()
    p.text = "技术架构 · 安全合规 · 部署运维 · 大模型趋势 · 人才战略"
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(0, 200, 255)
    p.alignment = PP_ALIGN.CENTER
    
    line_box = slide.shapes.add_textbox(Inches(4), Inches(4.6), Inches(8), Inches(0.05))
    line_frame = line_box.text_frame
    p = line_frame.add_paragraph()
    p.text = ""
    p.font.size = Pt(2)
    line_box.fill.solid()
    line_box.fill.fore_color.rgb = RGBColor(0, 180, 255)
    
    info_box = slide.shapes.add_textbox(Inches(3), Inches(5.5), Inches(10), Inches(1.2))
    info_frame = info_box.text_frame
    p = info_frame.add_paragraph()
    p.text = "A2A协议与AI智能体学习笔记"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(150, 170, 190)
    p.alignment = PP_ALIGN.CENTER
    
    p = info_frame.add_paragraph()
    p.text = "2026年5月 · 个人学习总结"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(120, 140, 160)
    p.alignment = PP_ALIGN.CENTER
    
    add_copyright(slide)
    return slide

def create_chapter_slide(prs, chapter):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    colors = add_chapter_background(slide, chapter['num'])
    add_watermark(slide)
    
    shapes = slide.shapes
    create_corporate_logo_element(slide, Inches(13.5), Inches(0.3), Inches(1.8), colors)
    create_network_node(slide, Inches(13.8), Inches(2.5), Inches(0.5), colors)
    create_tech_circuit(slide, Inches(13.5), Inches(3.2), Inches(2), colors)
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(12), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.add_paragraph()
    p.text = chapter['title']
    p.font.size = Pt(24)
    p.font.color.rgb = colors['accent']
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    desc = CHAPTER_DESCRIPTIONS.get(chapter['num'], "")
    if desc:
        desc_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(12), Inches(0.6))
        desc_frame = desc_box.text_frame
        desc_frame.word_wrap = True
        p = desc_frame.add_paragraph()
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = colors['text']
        p.alignment = PP_ALIGN.LEFT
    
    if chapter['num'] in DEFAULT_TABLES:
        create_table(slide, DEFAULT_TABLES[chapter['num']], colors)
    elif chapter['tables']:
        max_rows = 0
        best_table = None
        for table in chapter['tables']:
            if len(table) > max_rows:
                max_rows = len(table)
                best_table = table
        
        if best_table:
            create_table(slide, best_table, colors)
    
    resources = CHAPTER_RESOURCES.get(chapter['num'], [])
    if resources:
        create_resources_box(slide, resources, colors)
    
    if chapter['num'] in [1, 3, 6]:
        disclaimer_box = slide.shapes.add_textbox(Inches(1), Inches(7.8), Inches(14), Inches(0.5))
        disclaimer_frame = disclaimer_box.text_frame
        disclaimer_frame.word_wrap = True
        p = disclaimer_frame.add_paragraph()
        p.text = "本内容为学习示例，不构成专业建议"
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(120, 140, 160)
        p.alignment = PP_ALIGN.CENTER
    
    add_copyright(slide)
    return slide

def create_table(slide, table_data, colors):
    num_rows = len(table_data)
    num_cols = len(table_data[0]) if num_rows > 0 else 3
    left = Inches(0.4)
    top = Inches(2.0)
    width = Inches(14.8)
    height = Inches(4.2)
    table = slide.shapes.add_table(num_rows, num_cols, left, top, width, height).table
    highlight_keywords = ['核心', '关键', '重要', '必须', '最佳', '提升', '降低', '99', 'SOC2', '百万', '🔴']
    
    for i, row in enumerate(table_data):
        for j, cell in enumerate(row):
            cell_text = cell.replace('🔴', '')
            table.cell(i, j).text = cell_text
            paragraph = table.cell(i, j).text_frame.paragraphs[0]
            if num_cols >= 9:
                paragraph.font.size = Pt(6.8)
            elif num_cols == 8:
                paragraph.font.size = Pt(7.2)
            elif num_cols >= 6:
                paragraph.font.size = Pt(7.8)
            elif num_cols == 5:
                paragraph.font.size = Pt(8.5)
            else:
                paragraph.font.size = Pt(9)
            paragraph.alignment = PP_ALIGN.CENTER
            
            if i == 0:
                paragraph.font.color.rgb = RGBColor(255, 255, 255)
                paragraph.font.bold = True
                table.cell(i, j).fill.solid()
                table.cell(i, j).fill.fore_color.rgb = colors['accent']
            else:
                is_highlight = any(keyword in cell for keyword in highlight_keywords)
                if is_highlight:
                    paragraph.font.color.rgb = colors['highlight']
                    paragraph.font.bold = True
                else:
                    paragraph.font.color.rgb = colors['text']
                table.cell(i, j).fill.solid()
                table.cell(i, j).fill.fore_color.rgb = RGBColor(25, 30, 40)
                table.cell(i, j).fill.transparency = 0.3

def create_resources_box(slide, resources, colors):
    res_box = slide.shapes.add_textbox(Inches(0.7), Inches(6.5), Inches(14), Inches(1.7))
    res_frame = res_box.text_frame
    res_frame.word_wrap = True
    p = res_frame.add_paragraph()
    p.text = "📚 拓展学习资源"
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(100, 220, 255)
    p.font.bold = True
    for res in resources[:3]:
        p = res_frame.add_paragraph()
        p.text = res
        p.font.size = Pt(8)
        p.font.color.rgb = RGBColor(170, 185, 200)

def create_enterprise_summary_slide(prs):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(8, 12, 22)
    shapes = slide.shapes
    
    large_gradient = shapes.add_shape(MSO_SHAPE.OVAL, Inches(4), Inches(-3), Inches(16), Inches(14))
    large_gradient.fill.solid()
    large_gradient.fill.fore_color.rgb = RGBColor(0, 80, 160)
    large_gradient.fill.transparency = 0.18
    large_gradient.line.fill.background()
    
    for i in range(3):
        hex_left = Inches(0.3 + i * 0.8)
        hex_top = Inches(0.3 + i * 0.5)
        create_corporate_logo_element(slide, hex_left, hex_top, Inches(1.2), {'accent': RGBColor(0, 180, 255), 'secondary': RGBColor(0, 220, 200)})
    
    for i in range(5):
        node_left = Inches(14 - i * 0.6)
        node_top = Inches(7.5 + (i % 2) * 0.4)
        create_network_node(slide, node_left, node_top, Inches(0.35), {'accent': RGBColor(0, 180, 255), 'secondary': RGBColor(0, 200, 255)})
    
    add_watermark(slide)
    
    title_box = slide.shapes.add_textbox(Inches(2), Inches(0.4), Inches(12), Inches(1))
    title_frame = title_box.text_frame
    p = title_frame.add_paragraph()
    p.text = "🎯 A2A协议核心要点总结"
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    summary_data = [
        ['章节', '核心内容', '关键价值', '技术亮点'],
        ['第一章', '去中心化架构', '高可用弹性', '标准化协议'],
        ['第二章', 'JSON-RPC规范', '跨平台互操作', '异步通信'],
        ['第三章', 'OAuth2+mTLS', '🔴企业级安全', '零信任架构'],
        ['第四章', 'K8s部署', '🔴99.99%可用', '自动扩缩容'],
        ['第五章', '缓存+批量', '🔴性能提升3x', '低延迟设计'],
        ['第六章', '实战案例', '业务价值验证', '多Agent协作'],
        ['第七章', 'AI发展趋势', '2026-2046演进', '战略布局'],
        ['第八章', '人才规划', '全人群覆盖', '残疾人就业'],
        ['第九章', '创业指南', '政策支持', '零成本注册'],
        ['第十章', '大模型生态', '多模态融合', '成本优化']
    ]
    
    table = slide.shapes.add_table(11, 4, Inches(0.8), Inches(1.4), Inches(14.5), Inches(6.2)).table
    highlight_keywords = ['核心', '关键', '重要', '必须', '最佳', '提升', '降低', '99', 'SOC2', '百万', '🔴']
    
    for i, row in enumerate(summary_data):
        for j, cell in enumerate(row):
            cell_text = cell.replace('🔴', '')
            table.cell(i, j).text = cell_text
            paragraph = table.cell(i, j).text_frame.paragraphs[0]
            paragraph.font.size = Pt(10)
            paragraph.alignment = PP_ALIGN.CENTER
            
            if i == 0:
                paragraph.font.color.rgb = RGBColor(255, 255, 255)
                paragraph.font.bold = True
                table.cell(i, j).fill.solid()
                table.cell(i, j).fill.fore_color.rgb = RGBColor(0, 150, 220)
            else:
                is_highlight = any(keyword in row[j] for keyword in highlight_keywords)
                if is_highlight:
                    paragraph.font.color.rgb = RGBColor(255, 100, 100)
                    paragraph.font.bold = True
                else:
                    paragraph.font.color.rgb = RGBColor(220, 230, 240)
                table.cell(i, j).fill.solid()
                table.cell(i, j).fill.fore_color.rgb = RGBColor(25, 30, 40)
    
    add_copyright(slide)
    return slide

def create_related_practice_slide(prs):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(8, 12, 22)
    shapes = slide.shapes
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(12), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.add_paragraph()
    p.text = "相关实践"
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(0, 200, 255)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    desc_box = slide.shapes.add_textbox(Inches(1.5), Inches(0.9), Inches(12), Inches(0.4))
    desc_frame = desc_box.text_frame
    p = desc_frame.add_paragraph()
    p.text = "以下是基于本学习内容的实际项目实践"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(150, 170, 190)
    p.alignment = PP_ALIGN.CENTER
    
    practice_box = slide.shapes.add_textbox(Inches(1.5), Inches(1.3), Inches(12), Inches(3.0))
    practice_frame = practice_box.text_frame
    practice_frame.word_wrap = True
    
    practices = [
        "1. 基于RAG的6个领域问答系统",
        "   - 通用RAG智能问答系统",
        "   - 法律知识问答系统",
        "   - 教育学习问答系统",
        "   - 医疗健康问答系统",
        "   - 金融投资问答系统",
        "   - IT技术问答系统",
        "",
        "2. GitHub仓库地址",
        "   https://github.com/Hxdmou/legal-rag-qa-system",
        "",
        "3. 知乎文章链接",
        "   搜索：我是如何用RAG做出6个领域的知识问答系统的",
        "",
        "4. 学习笔记",
        "   docs/learning-notes/ 目录包含详细的技术学习笔记",
        "",
        "5. 联系方式",
        "   - GitHub: https://github.com/Hxdmou",
        "   - 商务邮箱: business@rag-qa-system.com"
    ]
    
    for line in practices:
        p = practice_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(220, 230, 240)
        if line.startswith("1.") or line.startswith("2.") or line.startswith("3.") or line.startswith("4.") or line.startswith("5."):
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 200, 255)
    
    source_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.3), Inches(12), Inches(0.8))
    source_frame = source_box.text_frame
    source_frame.word_wrap = True
    
    sources = [
        "来源声明：",
        "  - 部分概念图来源于网络，仅用于学习交流",
        "  - 如涉及侵权，请联系删除"
    ]
    
    for line in sources:
        p = source_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(200, 220, 240)
        if line.startswith("来源声明："):
            p.font.bold = True
            p.font.size = Pt(12)
    
    disclaimer_box = slide.shapes.add_textbox(Inches(1.5), Inches(5.1), Inches(12), Inches(1.0))
    disclaimer_frame = disclaimer_box.text_frame
    disclaimer_frame.word_wrap = True
    
    disclaimers = [
        "免责声明：",
        "  - 本内容仅供学习参考，不构成医疗、金融、法律等专业建议",
        "  - 医疗健康、金融投资等领域的决策，请咨询专业人士",
        "  - 所有数据为学习示例，实际效果因场景而异"
    ]
    
    for line in disclaimers:
        p = disclaimer_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(255, 200, 100)
        if line.startswith("免责声明："):
            p.font.bold = True
            p.font.size = Pt(14)
    
    add_copyright(slide)
    return slide

def create_usage_guide_slide(prs):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(8, 12, 22)
    shapes = slide.shapes
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(12), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.add_paragraph()
    p.text = "使用说明"
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(0, 200, 255)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    content_box = slide.shapes.add_textbox(Inches(2), Inches(1.5), Inches(12), Inches(5.0))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    content_lines = [
        "本PPT为个人学习整理，部分概念图来自网络，不构成商业建议。",
        "",
        "完整实践案例请参考GitHub仓库：",
        "https://github.com/Hxdmou/legal-rag-qa-system",
        "",
        "免责声明：",
        "- 本内容仅供学习参考，不构成医疗、金融、法律等专业建议",
        "- 医疗健康、金融投资等领域的决策，请咨询专业人士",
        "- 所有数据为学习示例，实际效果因场景而异"
    ]
    
    for line in content_lines:
        p = content_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(220, 230, 240)
        if line.startswith("免责声明："):
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 200, 100)
        elif line.startswith("-"):
            p.font.color.rgb = RGBColor(255, 200, 100)
    
    add_copyright(slide)
    return slide

def create_thank_you_slide(prs):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(8, 12, 22)
    shapes = slide.shapes
    
    large_gradient = shapes.add_shape(MSO_SHAPE.OVAL, Inches(4), Inches(-2), Inches(14), Inches(12))
    large_gradient.fill.solid()
    large_gradient.fill.fore_color.rgb = RGBColor(0, 80, 160)
    large_gradient.fill.transparency = 0.22
    large_gradient.line.fill.background()
    
    for i in range(4):
        create_corporate_logo_element(slide, Inches(1 + i * 3.5), Inches(1 + (i % 2) * 0.8), Inches(1.2), {'accent': RGBColor(0, 180, 255), 'secondary': RGBColor(0, 200, 255)})
    
    create_network_node(slide, Inches(0.5), Inches(5), Inches(0.6), {'accent': RGBColor(0, 180, 255), 'secondary': RGBColor(0, 200, 255)})
    create_network_node(slide, Inches(14.2), Inches(5), Inches(0.6), {'accent': RGBColor(0, 180, 255), 'secondary': RGBColor(0, 200, 255)})
    
    thank_box = slide.shapes.add_textbox(Inches(2), Inches(3.5), Inches(12), Inches(2))
    thank_frame = thank_box.text_frame
    p = thank_frame.add_paragraph()
    p.text = "THANK YOU"
    p.font.size = Pt(56)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    p = thank_frame.add_paragraph()
    p.text = "谢谢观看"
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(0, 200, 255)
    p.alignment = PP_ALIGN.CENTER
    
    contact_box = slide.shapes.add_textbox(Inches(3), Inches(6.5), Inches(10), Inches(0.8))
    contact_frame = contact_box.text_frame
    p = contact_frame.add_paragraph()
    p.text = "A2A协议与AI智能体学习笔记 | 2026"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(150, 170, 190)
    p.alignment = PP_ALIGN.CENTER
    
    add_copyright(slide)
    return slide

def parse_markdown_for_tables(md_content):
    chapters = {}
    lines = md_content.split('\n')
    i = 0
    
    while i < len(lines):
        line = lines[i]
        if line.startswith('# 第') and '章：' in line:
            title = line[2:].strip()
            chapter_num = None
            for idx, expected_title in enumerate(CHAPTER_TITLES, 1):
                if title.startswith(expected_title[:5]):
                    chapter_num = idx
                    break
            if chapter_num is None:
                i += 1
                continue
            tables = []
            i += 1
            while i < len(lines):
                if lines[i].startswith('|'):
                    table_data = []
                    while i < len(lines) and lines[i].startswith('|'):
                        row = lines[i].strip().split('|')
                        row = [cell.strip() for cell in row if cell.strip()]
                        table_data.append(row)
                        i += 1
                    if len(table_data) > 1:
                        tables.append(table_data)
                elif lines[i].startswith('# 第') and '章：' in lines[i]:
                    break
                elif lines[i].startswith('---'):
                    i += 1
                    continue
                else:
                    i += 1
            if not tables and chapter_num in DEFAULT_TABLES:
                tables.append(DEFAULT_TABLES[chapter_num])
            if tables:
                chapters[chapter_num] = {
                    'num': chapter_num,
                    'title': title,
                    'tables': tables
                }
        i += 1
    
    result = []
    for num in range(1, 7):
        if num in chapters:
            result.append(chapters[num])
        elif num in DEFAULT_TABLES:
            result.append({
                'num': num,
                'title': CHAPTER_TITLES[num-1],
                'tables': [DEFAULT_TABLES[num]]
            })
    return result

def create_pptx_enterprise(chapters, output_path):
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    create_enterprise_title_slide(prs)
    
    for chapter in chapters:
        create_chapter_slide(prs, chapter)
    
    create_related_practice_slide(prs)
    create_usage_guide_slide(prs)
    create_thank_you_slide(prs)
    
    prs.save(output_path)
    print(f"PPTX文件已生成: {output_path}")
    print(f"共生成 {len(chapters) + 4} 页幻灯片")

if __name__ == '__main__':
    md_file = r'f:\个人作品\legal-rag-qa-system\A2A_ENTERPRISE_PPT.md'
    output_file = r'f:\个人作品\legal-rag-qa-system\A2A_ENTERPRISE_PPT_V12.pptx'
    
    with open(md_file, 'r', encoding='utf-8') as f:
        md_content = f.read()
    
    chapters = parse_markdown_for_tables(md_content)
    print(f"检测到 {len(chapters)} 个章节")
    for ch in chapters:
        print(f"  第{ch['num']}章: {ch['title']}")
    
    create_pptx_enterprise(chapters, output_file)
