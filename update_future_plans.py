from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

def fill_cell(cell, text, font_size=9, bold=False):
    cell.text_frame.clear()
    para = cell.text_frame.add_paragraph()
    run = para.add_run()
    run.text = str(text) if text else "---"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(255, 255, 255)
    run.font.name = '微软雅黑'
    para.alignment = PP_ALIGN.CENTER
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

def update_all_slides(prs):
    # 更新项目5 - AI趋势热点（8行17列）
    if len(prs.slides) > 4:
        slide = prs.slides[4]
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                # 行1: 表头
                # 行2-8: 8大热点
                data = [
                    ["趋势", "成熟度", "市场影响", "代表技术", "商业价值", "技术难度", "投资热度", "竞争格局", "预期收益", "风险等级", "时间窗口", "人才需求", "落地建议", "投资回报", "10年规划", "优先级", "推荐度"],
                    ["新模型爆发", "高", "颠覆性", "GPT 5.0/DeepSeek V5.0", "极高", "高", "极热", "激烈", "8倍", "中", "2026", "紧缺", "抢先布局", "12个月", "技术领先", "P0", "10"],
                    ["具身智能", "中", "深远", "机器人+LLM融合", "极高", "极高", "很热", "中等", "20倍", "高", "2027-2028", "极缺", "战略投入", "战略价值", "生态闭环", "P1", "9.5"],
                    ["AI安全对齐", "中高", "关键", "RLHF+Constitutional", "极高", "高", "很热", "中等", "5倍", "低", "2026-2027", "紧缺", "合规必投", "合规保障", "安全体系", "P0", "10"],
                    ["多模态融合", "高", "广泛", "GPT 5.0/Gemini 2.5", "高", "中高", "很热", "激烈", "3倍", "中", "2026-2027", "紧缺", "垂直落地", "年回报30%", "行业深耕", "P0", "9.5"],
                    ["边缘AI", "中", "普及", "端侧推理+隐私计算", "中高", "中", "较热", "中等", "2倍", "低", "2027-2028", "中", "逐步推进", "稳健增长", "边缘协同", "P1", "9"],
                    ["AI机器人", "中", "爆发", "人形机器人+AGI", "极高", "极高", "极热", "激烈", "15倍", "高", "2028-2030", "极缺", "战略布局", "长期价值", "产业变革", "P0", "10"],
                    ["量子AI", "低", "颠覆", "量子计算+AI融合", "极高", "极高", "较热", "早期", "50倍", "极高", "2030-2035", "极缺", "前沿探索", "战略期权", "突破创新", "P2", "9"],
                    ["脑机接口", "低", "革命", "BCI+神经连接", "极高", "极高", "较热", "早期", "30倍", "极高", "2032-2035", "极缺", "医疗切入", "长期投资", "人机共生", "P2", "9"],
                ]
                for i in range(len(data)):
                    for j in range(len(data[i])):
                        if i < len(table.rows) and j < len(table.columns):
                            fill_cell(table.cell(i, j), data[i][j], 12 if i == 0 else 9, i == 0)
                break

    # 更新项目10 - 市场预测（7行13列）
    if len(prs.slides) > 9:
        slide = prs.slides[9]
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                data = [
                    ["年份", "市场规模", "增长率", "主要驱动力", "投资规模", "技术突破", "应用场景", "竞争格局", "政策环境", "人才供给", "技术成熟度", "商业化程度", "投资建议"],
                    ["2026", "2000亿", "85%", "GPT 5.0发布", "600亿", "GPT 5.0/DeepSeek V5.0", "企业服务爆发", "格局初定", "政策完善", "紧缺", "高", "高", "重仓"],
                    ["2027", "3800亿", "90%", "具身智能崛起", "1200亿", "GPT 5.5/人形机器人", "智能硬件普及", "巨头主导", "监管落地", "紧缺", "高", "极高", "长期"],
                    ["2028", "6500亿", "71%", "AI机器人爆发", "2000亿", "GPT 6.0/Optimus 2.0", "服务机器人普及", "生态竞争", "全球标准", "平衡", "极高", "极高", "稳健"],
                    ["2029", "10000亿", "54%", "脑机接口突破", "3000亿", "DeepSeek V6.0/BCI", "人机融合", "寡头垄断", "伦理规范", "充足", "极高", "极高", "价值"],
                    ["2030", "15000亿", "50%", "AGI原型落地", "4500亿", "GPT 7.0/AGI", "全面渗透", "全球格局", "国际治理", "充足", "极高", "极高", "战略"],
                    ["2031", "22000亿", "47%", "量子AI商用", "6000亿", "量子AI融合", "智能经济", "生态竞争", "标准统一", "充足", "极高", "极高", "价值"],
                    ["2035", "85000亿", "37%", "AGI全面普及", "15000亿", "通用AGI", "智能时代", "全球一体", "AI治理", "充足", "极高", "极高", "长期价值"],
                ]
                for i in range(len(data)):
                    for j in range(len(data[i])):
                        if i < len(table.rows) and j < len(table.columns):
                            fill_cell(table.cell(i, j), data[i][j], 12 if i == 0 else 9, i == 0)
                break

    # 更新项目11 - 技术评估（6行11列）
    if len(prs.slides) > 10:
        slide = prs.slides[10]
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                data = [
                    ["维度", "价值描述", "量化指标", "测试环境", "数据规模", "准确率", "响应时间", "资源消耗", "可扩展性", "安全性", "商业价值"],
                    ["推理能力", "复杂任务+逻辑推理", "96%", "GPT 5.0", "500万token", "93%", "<150ms", "中", "高", "高", "极高"],
                    ["多模态", "图文音视频+3D", "92%", "Gemini 2.5", "50万样本", "90%", "<400ms", "高", "中", "中", "高"],
                    ["长上下文", "百万Token+记忆", "88%", "Claude 3.5", "1M+ token", "87%", "<800ms", "高", "中", "中", "极高"],
                    ["工具调用", "API+代码执行", "99%", "GPT 5.0", "10000次", "99.5%", "<80ms", "低", "高", "高", "高"],
                    ["自我反思", "错误纠正+优化", "85%", "DeepSeek V5.0", "10万样本", "82%", "<1.5s", "高", "中", "中", "极高"],
                    ["综合评分", "全面能力评估", "91%", "GPT 5.0/DeepSeek", "综合测试", "90%", "<300ms", "中", "高", "高", "极高"],
                ]
                for i in range(len(data)):
                    for j in range(len(data[i])):
                        if i < len(table.rows) and j < len(table.columns):
                            fill_cell(table.cell(i, j), data[i][j], 12 if i == 0 else 9, i == 0)
                break

    # 更新项目12 - 总结（11行3列）
    if len(prs.slides) > 11:
        slide = prs.slides[11]
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                data = [
                    ["项目", "核心内容", "未来10年规划"],
                    ["1", "封面升级", "季度更新迭代"],
                    ["2", "目录优化", "AI智能导航"],
                    ["3", "协议架构", "A2A 3.0升级(2027)"],
                    ["4", "能力矩阵", "扩展至10层(2028)"],
                    ["5", "16大热点", "年度更新跟踪"],
                    ["6", "应用场景", "扩展至20领域(2029)"],
                    ["7", "案例分析", "扩展至20案例(2027)"],
                    ["8", "技术对比", "季度对比报告"],
                    ["9", "市场预测", "年度滚动预测"],
                    ["10", "10大维度", "扩展至15维度(2030)"],
                    ["合计", "完整AI方案", "分阶段落地"],
                ]
                for i in range(len(data)):
                    for j in range(len(data[i])):
                        if i < len(table.rows) and j < len(table.columns):
                            fill_cell(table.cell(i, j), data[i][j], 12 if i == 0 else 9, i == 0)
                break

def main():
    ppt_path = r"f:\个人作品\legal-rag-qa-system\A2A_PROTOCOL_AI_AGENT_V14_终极完整版.pptx"
    prs = Presentation(ppt_path)
    
    print("🔄 正在补充最新大模型和AI机器人未来10年规划...")
    update_all_slides(prs)
    
    output_path = r"f:\个人作品\legal-rag-qa-system\A2A_PROTOCOL_AI_AGENT_V14_终极完整版.pptx"
    prs.save(output_path)
    print(f"✅ 已完成！保存到: {output_path}")
    
    print("\n📌 补充内容:")
    print("   🤖 大模型规划:")
    print("      2026年: GPT 5.0 / DeepSeek V5.0")
    print("      2027年: GPT 5.5 / 人形机器人")
    print("      2028年: GPT 6.0 / Optimus 2.0")
    print("      2029年: DeepSeek V6.0 / BCI")
    print("      2030年: GPT 7.0 / AGI")
    print("      2031年: 量子AI融合")
    print("      2035年: 通用AGI")
    print("   🤖 AI机器人应用场景:")
    print("      2026-2027: 企业服务、智能硬件")
    print("      2028-2030: 服务机器人普及、人机融合")
    print("      2030-2035: 智能经济、人机共生")

if __name__ == "__main__":
    main()