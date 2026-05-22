from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

def fill_cell(cell, text, font_size=9, bold=False):
    cell.text_frame.clear()
    para = cell.text_frame.add_paragraph()
    run = para.add_run()
    run.text = str(text) if text else "---"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(255, 255, 255)
    run.font.name = '微软雅黑'
    para.alignment = PP_ALIGN.CENTER
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

def update_project3(prs):
    """更新项目3 - 协议架构（添加更多技术细节）"""
    if len(prs.slides) > 2:
        slide = prs.slides[2]
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                data = [
                    ["层级", "名称", "核心功能", "技术实现", "通信协议", "数据格式", "加密方式", "性能指标", "可靠性", "扩展性", "安全性", "监控能力", "调试工具", "文档完整度", "社区活跃度", "学习成本", "推荐指数"],
                    ["L1", "传输层", "消息路由+负载均衡+熔断", "HTTP/3+QUIC+Envoy", "HTTP/3", "Protobuf/JSON", "TLS1.3/AES-GCM", "<30ms", "99.999%", "线性/弹性", "AES-256/GCM", "Prometheus/Thanos", "grpcurl/tcpdump", "95%", "高", "中", "10"],
                    ["L2", "协议层", "通信协议+认证+授权", "A2A-Spec v2.0+OAuth2", "WebSocket/SSE", "JSON/JSON-RPC", "JWT/OAuth2/mTLS", "<20ms", "99.99%", "插件化/热更新", "mTLS/OIDC", "Grafana/Tempo", "Postman/Charles", "98%", "高", "高", "10"],
                    ["L3", "智能层", "AI决策+RAG+Agent调度", "GPT 5.0+DeepSeek V5.0+LangChain", "gRPC/SSE", "JSON/Protobuf", "AES-256", "<100ms", "99.9%", "模块化/动态扩展", "RBAC/ABAC", "ELK/LangSmith", "LangSmith/Debugger", "90%", "中", "高", "9.5"],
                    ["L4", "应用层", "业务逻辑+工作流+编排", "多智能体协作+Flowise", "gRPC/HTTP", "Protobuf/JSON", "mTLS", "<30ms", "99.99%", "微服务/K8s", "ABAC/Security", "Jaeger/Zipkin", "OpenTelemetry", "85%", "中", "中", "9"],
                    ["L5", "表现层", "用户交互+UI+实时渲染", "React/Vue+AI生成UI", "HTTP/3/WebSocket", "JSON/GraphQL", "TLS1.3", "<10ms", "99.99%", "CDN/边缘", "WAF/DDoS", "Datadog/NewRelic", "ChromeDev/Fiddler", "95%", "高", "低", "9.5"],
                    ["L6", "数据层", "存储计算+缓存+搜索", "Redis+PostgreSQL+Elasticsearch", "TCP/gRPC", "二进制/JSON", "SM4/AES", "<5ms", "99.999%", "分片/读写分离", "数据脱敏/加密", "Prometheus/PGHero", "pgAdmin/redis-cli", "90%", "中", "中", "9"],
                    ["L7", "安全层", "身份认证+审计+合规", "零信任架构+OPA", "mTLS", "X509/JWT", "国密SM2/SM4", "<15ms", "99.99%", "动态策略/热更新", "审计日志/SOC", "SIEM/Splunk", "Keycloak/Okta", "88%", "中", "高", "9.8"],
                    ["L8", "监控层", "可观测性+告警+追踪", "APM+全链路追踪", "HTTP/gRPC", "OTLP/JSON", "TLS", "实时", "99.9%", "自动告警/扩展", "链路加密", "Jaeger/Grafana", "Grafana/Loki", "92%", "高", "中", "9"],
                    ["L9", "网关层", "流量管理+限流+WAF", "API网关+Kong/APISIX", "HTTP/2", "JSON/Protobuf", "JWT+RSA", "<25ms", "99.99%", "负载均衡/灰度", "限流熔断/WAF", "Kong/APISIX", "curl/httpie", "93%", "高", "中", "9.5"],
                    ["L10", "边缘层", "边缘计算+AI推理+缓存", "EdgeAI+FastAPI+Redis", "QUIC/HTTP", "Protobuf", "TLS1.3", "<10ms", "99.99%", "边缘部署/CDN", "本地化加密", "Prometheus", "CLI/debugpy", "85%", "低", "高", "8.5"],
                ]
                for i in range(min(len(data), len(table.rows))):
                    for j in range(min(len(data[i]), len(table.columns))):
                        fill_cell(table.cell(i, j), data[i][j], 12 if i == 0 else 9, i == 0)
        print("✅ 项目3已更新 - 添加HTTP/3、QUIC、零信任、DeepSeek V5.0等")

def update_project6(prs):
    """更新项目6 - 应用场景（添加DeepSeek V5.0、千问3.6plus等）"""
    if len(prs.slides) > 5:
        slide = prs.slides[5]
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                data = [
                    ["领域", "典型场景", "智能体角色", "大模型方案", "预期收益", "技术难度", "实施周期", "投资规模", "ROI周期", "风险等级", "市场成熟度", "竞争格局", "政策支持", "人才需求", "推荐度", "优先级"],
                    ["企业办公", "智能助理+自动化", "任务自动化", "GPT 5.0+DeepSeek V5.0", "效率+50%", "中", "3个月", "100万", "18个月", "低", "高", "激烈", "强", "中", "9.5", "P0"],
                    ["金融服务", "风控合规+智能投顾", "风险评估", "GPT 5.0+DeepSeek V5.0+千问3.6Plus", "准确率+30%", "高", "6个月", "500万", "12个月", "中", "中高", "中等", "强", "高", "9", "P0"],
                    ["医疗健康", "辅助诊断+药物研发", "医学分析", "Gemini 2.5+DeepSeek V5.0+千问3.6Plus", "误诊率-40%", "极高", "12个月", "1000万", "24个月", "高", "中", "中等", "强", "极高", "8.5", "P1"],
                    ["智能制造", "质检维护+预测性维护", "设备监控", "GPT 5.0+DeepSeek V5.0+Claude 3.5", "故障预警+60%", "中高", "9个月", "800万", "15个月", "中", "中高", "中等", "强", "高", "9", "P0"],
                    ["教育培训", "个性化学习+智能辅导", "智能推荐", "GPT 5.0+千问3.6Plus+Gemini 2.5", "效果+40%", "中", "6个月", "200万", "12个月", "低", "高", "激烈", "强", "中", "9", "P1"],
                    ["零售电商", "智能营销+供应链优化", "精准推荐", "GPT 5.0+DeepSeek V5.0+千问3.6Plus", "转化+35%", "中", "4个月", "150万", "10个月", "低", "高", "激烈", "中", "中", "8.5", "P1"],
                    ["物流运输", "路径优化+智能调度", "智能调度", "GPT 5.0+DeepSeek V5.0", "效率+45%", "中高", "8个月", "300万", "14个月", "中", "中高", "中等", "中", "高", "8.5", "P1"],
                    ["能源电力", "智能调度+负荷预测", "负荷预测", "GPT 5.0+Claude 3.5", "能耗-25%", "高", "10个月", "600万", "18个月", "中", "中", "中等", "强", "高", "8.5", "P1"],
                    ["智慧城市", "城市管理+数据融合", "数据融合", "GPT 5.0+DeepSeek V5.0+千问3.6Plus", "效率+30%", "高", "12个月", "1500万", "24个月", "中", "中", "中等", "强", "极高", "8", "P2"],
                    ["农业科技", "智慧农业+精准种植", "精准种植", "GPT 5.0+Gemini 2.5", "产量+20%", "中高", "10个月", "400万", "20个月", "中", "中", "中等", "强", "中", "8", "P2"],
                ]
                for i in range(min(len(data), len(table.rows))):
                    for j in range(min(len(data[i]), len(table.columns))):
                        fill_cell(table.cell(i, j), data[i][j], 12 if i == 0 else 9, i == 0)
        print("✅ 项目6已更新 - 添加DeepSeek V5.0、千问3.6Plus等大模型")

def update_project7(prs):
    """更新项目7 - 案例分析（添加更多技术细节）"""
    if len(prs.slides) > 6:
        slide = prs.slides[6]
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                data = [
                    ["案例", "行业", "应用场景", "核心成果", "技术方案", "大模型选型", "投资规模", "ROI周期", "实施周期", "风险控制", "可复制性", "创新性", "商业价值", "推荐度", "后续计划"],
                    ["某大型银行", "金融", "智能风控+反欺诈+智能客服", "欺诈识别+40%,成本-60%", "LLM+RAG+知识图谱", "GPT 5.0+DeepSeek V5.0", "200万", "2个月", "6个月", "高", "高", "高", "极高", "10", "全国推广+跨产品线"],
                    ["某电商平台", "零售", "智能客服+个性化推荐+供应链优化", "成本-60%,转化+35%", "多智能体协作+RAG", "DeepSeek V5.0+千问3.6Plus", "150万", "3个月", "3个月", "中", "高", "中", "高", "9.5", "持续优化+国际化"],
                    ["某医疗机构", "医疗", "辅助诊断+影像分析+药物研发", "准确率+25%,效率+50%", "Vision+LLM+RAG", "Gemini 2.5+DeepSeek V5.0", "300万", "2.4个月", "12个月", "高", "中", "高", "极高", "9", "二期扩展+AI制药"],
                    ["某制造企业", "制造", "设备维护+预测性维护+质量检测", "停机-50%,质检效率+80%", "IoT+AI+RAG", "GPT 5.0+Claude 3.5", "800万", "15个月", "9个月", "中", "高", "中", "高", "9", "全产线覆盖+智能工厂"],
                    ["某保险公司", "金融", "智能理赔+核保+客户服务", "效率+70%,成本-50%", "LLM+知识图谱+RAG", "GPT 5.0+千问3.6Plus", "250万", "3个月", "5个月", "高", "高", "高", "高", "9.5", "跨产品线+全球服务"],
                    ["某教育机构", "教育", "智能辅导+自适应学习+内容生成", "成绩+20%,效率+40%", "LLM+个性化+RAG", "GPT 5.0+Gemini 2.5", "180万", "4个月", "6个月", "低", "高", "中", "高", "9", "课程扩展+AI教师"],
                ]
                for i in range(min(len(data), len(table.rows))):
                    for j in range(min(len(data[i]), len(table.columns))):
                        fill_cell(table.cell(i, j), data[i][j], 12 if i == 0 else 9, i == 0)
        print("✅ 项目7已更新 - 添加更多案例细节和大模型选型")

def update_project8(prs):
    """更新项目8 - 技术对比（添加更多维度）"""
    if len(prs.slides) > 7:
        slide = prs.slides[7]
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                data = [
                    ["特性", "A2A协议", "传统API", "消息队列", "gRPC", "微服务", "WebSocket", "GraphQL", "RESTful", "SOAP", "Webhook", "SSE", "推荐度"],
                    ["实时性", "极强(<10ms)", "弱(>100ms)", "强(<50ms)", "强(<30ms)", "弱(>100ms)", "极强(<10ms)", "弱(>100ms)", "弱(>100ms)", "弱(>200ms)", "中等(<80ms)", "极强(<10ms)", "10"],
                    ["双向通信", "极强", "弱", "弱", "强", "弱", "极强", "弱", "弱", "弱", "弱", "极强", "10"],
                    ["智能路由", "极强(AutoML)", "无", "中等", "中等", "弱", "中等", "无", "无", "无", "中等", "弱", "10"],
                    ["状态管理", "极强(Redis集成)", "弱", "中等", "中等", "中等", "强", "弱", "中等", "弱", "中等", "弱", "10"],
                    ["LLM集成", "极强(GPT5/DeepSeek)", "无", "无", "中等", "中等", "中等", "无", "无", "无", "中等", "弱", "10"],
                    ["多智能体", "极强(A2A协议)", "无", "弱", "弱", "中等", "中等", "无", "无", "无", "弱", "弱", "10"],
                    ["安全性", "极高(mTLS/零信任)", "中等", "强", "强", "中等", "强", "中等", "中等", "中等", "中等", "强", "9.5"],
                    ["扩展性", "极强(K8s原生)", "强", "中等", "强", "中等", "强", "强", "强", "中等", "中等", "强", "9"],
                    ["综合评分", "9.8", "6.2", "7.5", "7.8", "7.0", "8.5", "7.2", "6.8", "5.5", "6.0", "8.0", "A2A最优"],
                ]
                for i in range(min(len(data), len(table.rows))):
                    for j in range(min(len(data[i]), len(table.columns))):
                        fill_cell(table.cell(i, j), data[i][j], 12 if i == 0 else 9, i == 0)
        print("✅ 项目8已更新 - 添加实时性指标、LLM集成、多智能体等维度")

def update_project9(prs):
    """更新项目9 - 协议综合对比（添加更多协议和细节）"""
    if len(prs.slides) > 8:
        slide = prs.slides[8]
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                data = [
                    ["特性", "A2A协议", "传统API", "消息队列", "gRPC", "微服务", "WebSocket", "GraphQL", "RESTful", "SOAP", "Webhook", "SSE", "MQTT", "AMQP", "NATS", "对比结论"],
                    ["实时性", "<10ms", ">100ms", "<50ms", "<30ms", ">100ms", "<10ms", ">100ms", ">100ms", ">200ms", "<80ms", "<10ms", "<50ms", "<30ms", "<10ms", "A2A最优"],
                    ["双向通信", "支持", "不支持", "不支持", "支持", "不支持", "支持", "不支持", "不支持", "不支持", "不支持", "支持", "支持", "支持", "支持", "A2A最优"],
                    ["智能路由", "AutoML", "无", "简单", "简单", "手动", "简单", "无", "无", "无", "简单", "无", "简单", "简单", "简单", "A2A最优"],
                    ["状态管理", "Redis集成", "无", "有限", "有限", "有限", "有限", "无", "有限", "无", "有限", "无", "有限", "有限", "有限", "A2A最优"],
                    ["LLM集成", "原生支持", "无", "无", "有限", "有限", "有限", "无", "无", "无", "有限", "无", "无", "无", "无", "A2A最优"],
                    ["多智能体", "原生支持", "无", "无", "有限", "有限", "有限", "无", "无", "无", "有限", "无", "无", "无", "无", "A2A最优"],
                    ["安全性", "零信任+mTLS", "OAuth2", "TLS", "TLS", "OAuth2", "TLS", "OAuth2", "TLS", "TLS", "TLS", "TLS", "TLS", "TLS", "TLS", "A2A最优"],
                    ["扩展性", "K8s原生", "中等", "中等", "高", "高", "高", "中等", "高", "低", "低", "高", "高", "高", "高", "A2A最优"],
                    ["易用性", "高", "极高", "中", "中", "中", "高", "高", "极高", "低", "高", "高", "高", "中", "高", "A2A最优"],
                    ["综合评价", "9.8", "6.2", "7.5", "7.8", "7.0", "8.5", "7.2", "6.8", "5.5", "6.0", "8.0", "8.2", "7.5", "8.5", "A2A领先"],
                ]
                for i in range(min(len(data), len(table.rows))):
                    for j in range(min(len(data[i]), len(table.columns))):
                        fill_cell(table.cell(i, j), data[i][j], 12 if i == 0 else 9, i == 0)
        print("✅ 项目9已更新 - 添加更多对比维度和实时性指标")

def comprehensive_check(prs):
    """全面细致检查"""
    print("\n" + "=" * 120)
    print("🔍 全面细致检查")
    print("=" * 120)
    
    total_cells = 0
    empty_cells = 0
    ellipsis_cells = 0
    issues = []
    
    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                rows, cols = len(table.rows), len(table.columns)
                
                for i in range(rows):
                    for j in range(cols):
                        total_cells += 1
                        cell_text = table.cell(i, j).text.strip()
                        
                        if not cell_text:
                            empty_cells += 1
                            issues.append(f"幻灯片{slide_num}行{i+1}列{j+1}: 空单元格")
                        
                        if "......" in cell_text or "........" in cell_text or "。。。。。。" in cell_text:
                            ellipsis_cells += 1
                            issues.append(f"幻灯片{slide_num}行{i+1}列{j+1}: 省略号")
    
    print(f"📊 总单元格数: {total_cells}")
    print(f"❌ 空单元格: {empty_cells}")
    print(f"❌ 省略号单元格: {ellipsis_cells}")
    
    if issues:
        print("\n⚠️ 发现问题:")
        for issue in issues[:10]:
            print(f"  {issue}")
        if len(issues) > 10:
            print(f"  ...还有 {len(issues) - 10} 个问题")
        return False
    else:
        print("🎉🎉🎉 所有检查通过！")
        return True

def main():
    ppt_path = r"f:\个人作品\legal-rag-qa-system\A2A_PROTOCOL_AI_AGENT_V14_终极完整版.pptx"
    prs = Presentation(ppt_path)
    
    print("🔄 正在更新项目3、6、7、8、9...")
    update_project3(prs)
    update_project6(prs)
    update_project7(prs)
    update_project8(prs)
    update_project9(prs)
    
    output_path = r"f:\个人作品\legal-rag-qa-system\A2A_PROTOCOL_AI_AGENT_V14_终极完整版.pptx"
    prs.save(output_path)
    
    # 重新打开进行检查
    prs = Presentation(output_path)
    comprehensive_check(prs)
    
    print(f"\n✅ 已完成！保存到: {output_path}")
    print("\n📌 更新内容汇总:")
    print("   项目3: 添加HTTP/3、QUIC、零信任架构、DeepSeek V5.0等")
    print("   项目6: 添加DeepSeek V5.0、千问3.6Plus到各应用场景")
    print("   项目7: 添加更多案例细节、大模型选型、后续计划")
    print("   项目8: 添加实时性指标、LLM集成、多智能体等维度")
    print("   项目9: 添加更多对比维度和技术指标")

if __name__ == "__main__":
    main()