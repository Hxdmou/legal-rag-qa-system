from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

def fill_cell(cell, text, font_size=9, bold=False):
    cell.text_frame.clear()
    para = cell.text_frame.add_paragraph()
    run = para.add_run()
    run.text = str(text) if text else "---"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(255, 255, 255)
    run.font.name = '微软雅黑'
    para.alignment = PP_ALIGN.CENTER
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

def update_robot_scenarios(prs):
    # 更新项目5 - AI趋势热点中的AI机器人部分
    if len(prs.slides) > 4:
        slide = prs.slides[4]
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                # 更新AI机器人行（第6行）
                fill_cell(table.cell(5, 0), "AI机器人")
                fill_cell(table.cell(5, 1), "中")
                fill_cell(table.cell(5, 2), "爆发")
                fill_cell(table.cell(5, 3), "人形机器人+AGI")
                fill_cell(table.cell(5, 4), "极高")
                fill_cell(table.cell(5, 5), "极高")
                fill_cell(table.cell(5, 6), "极热")
                fill_cell(table.cell(5, 7), "激烈")
                fill_cell(table.cell(5, 8), "15倍")
                fill_cell(table.cell(5, 9), "高")
                fill_cell(table.cell(5, 10), "2028-2035")
                fill_cell(table.cell(5, 11), "极缺")
                fill_cell(table.cell(5, 12), "战略布局")
                fill_cell(table.cell(5, 13), "长期价值")
                fill_cell(table.cell(5, 14), "产业变革")
                fill_cell(table.cell(5, 15), "P0")
                fill_cell(table.cell(5, 16), "10")
                print("✅ 已更新项目5 AI机器人内容")
                break
    
    # 更新项目10 - 市场预测，细化每个年份的应用场景
    if len(prs.slides) > 9:
        slide = prs.slides[9]
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                # 详细的年度应用场景
                scenarios = [
                    "企业服务爆发",  # 2026
                    "智能硬件普及+工业质检",  # 2027
                    "服务机器人普及+物流配送",  # 2028
                    "家庭助手+医疗护理",  # 2029
                    "人机融合+教育陪伴",  # 2030
                    "智能经济+农业机器人",  # 2031
                    "万物互联+城市服务",  # 2032
                    "人机共生+养老服务",  # 2033
                    "智能文明+太空探索",  # 2034
                    "智能时代+全面智能化",  # 2035
                ]
                # 更新应用场景列
                for i in range(1, min(len(scenarios)+1, len(table.rows))):
                    if i < len(table.rows) and 6 < len(table.columns):
                        fill_cell(table.cell(i, 6), scenarios[i-1])
                print("✅ 已更新项目10年度应用场景")
                break

    # 添加专门的AI机器人发展规划到项目12总结
    if len(prs.slides) > 11:
        slide = prs.slides[11]
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                # 更新项目总结中的AI机器人规划
                fill_cell(table.cell(6, 1), "AI机器人案例")
                fill_cell(table.cell(6, 2), "2028年服务机器人普及")
                fill_cell(table.cell(7, 1), "技术对比")
                fill_cell(table.cell(7, 2), "季度对比报告")
                fill_cell(table.cell(8, 1), "市场预测")
                fill_cell(table.cell(8, 2), "2026-2035年滚动预测")
                fill_cell(table.cell(9, 1), "技术评估")
                fill_cell(table.cell(9, 2), "扩展至15维度(2030)")
                print("✅ 已更新项目12总结")
                break

def main():
    ppt_path = r"f:\个人作品\legal-rag-qa-system\A2A_PROTOCOL_AI_AGENT_V14_终极完整版.pptx"
    prs = Presentation(ppt_path)
    
    print("🔄 正在更新AI机器人应用场景（按年份详细展开）...")
    update_robot_scenarios(prs)
    
    output_path = r"f:\个人作品\legal-rag-qa-system\A2A_PROTOCOL_AI_AGENT_V14_终极完整版.pptx"
    prs.save(output_path)
    print(f"✅ 已完成！保存到: {output_path}")
    
    print("\n📌 AI机器人应用场景详细规划（按年份）:")
    print("   2026年: 企业服务爆发 - 智能客服、办公自动化、数据处理")
    print("   2027年: 智能硬件普及、工业质检 - 生产线上料、质量检测、设备巡检")
    print("   2028年: 服务机器人普及、物流配送 - 仓储分拣、末端配送、餐厅服务")
    print("   2029年: 家庭助手、医疗护理 - 陪伴机器人、康复辅助、健康监测")
    print("   2030年: 人机融合、教育陪伴 - 个性化学习、技能培训、情感陪伴")
    print("   2031年: 智能经济、农业机器人 - 精准种植、畜牧养殖、农产品加工")
    print("   2032年: 万物互联、城市服务 - 智慧城市管理、公共安全、环境监测")
    print("   2033年: 人机共生、养老服务 - 老年护理、生活辅助、健康管理")
    print("   2034年: 智能文明、太空探索 - 太空机器人、深海探测、极端环境作业")
    print("   2035年: 智能时代、全面智能化 - AGI机器人、脑机接口应用、超级智能")

if __name__ == "__main__":
    main()