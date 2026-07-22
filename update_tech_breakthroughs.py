from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

def fill_cell(cell, text, font_size=9, bold=False):
    cell.text_frame.clear()
    para = cell.text_frame.add_paragraph()
    run = para.add_run()
    run.text = str(text) if text else "---"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(255, 255, 255)
    run.font.name = '微软雅黑'
    para.alignment = PP_ALIGN.CENTER
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

def update_tech_breakthroughs(prs):
    # 更新项目10 - 市场预测
    if len(prs.slides) > 9:
        slide = prs.slides[9]
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                # 完整数据 - 更新2027-2035年技术突破
                data = [
                    ["年份", "市场规模", "增长率", "主要驱动力", "投资规模", "技术突破", "应用场景", "竞争格局", "政策环境", "人才供给", "技术成熟度", "商业化程度", "投资建议"],
                    ["2026", "2000亿", "85%", "GPT 5.0发布", "600亿", "GPT 5.0/DeepSeek V5.0", "企业服务爆发", "格局初定", "政策完善", "紧缺", "高", "高", "重仓"],
                    ["2027", "3800亿", "90%", "具身智能崛起", "1200亿", "GPT 5.5/人形机器人", "智能硬件普及", "巨头主导", "监管落地", "紧缺", "高", "极高", "长期"],
                    ["2028", "6500亿", "71%", "AI机器人爆发", "2000亿", "GPT 6.0/Optimus 2.0", "服务机器人普及", "生态竞争", "全球标准", "平衡", "极高", "极高", "稳健"],
                    ["2029", "10000亿", "54%", "脑机接口突破", "3000亿", "DeepSeek V6.0/BCI神经连接", "人机融合", "寡头垄断", "伦理规范", "充足", "极高", "极高", "价值"],
                    ["2030", "15000亿", "50%", "AGI原型落地", "4500亿", "GPT 7.0/通用AGI原型", "全面渗透", "全球格局", "国际治理", "充足", "极高", "极高", "战略"],
                    ["2031", "22000亿", "47%", "量子AI商用", "6000亿", "量子AI融合/IBM Q1000", "智能经济", "生态竞争", "标准统一", "充足", "极高", "极高", "价值"],
                    ["2032", "32000亿", "45%", "边缘智能爆发", "8000亿", "GPT 8.0/边缘大脑", "万物互联", "平台竞争", "全球协作", "充足", "极高", "极高", "稳健"],
                    ["2033", "45000亿", "41%", "脑机融合深化", "10000亿", "DeepSeek V7.0/脑机接口2.0", "人机共生", "寡头格局", "伦理框架", "充足", "极高", "极高", "长期"],
                    ["2034", "62000亿", "38%", "量子AGI突破", "12000亿", "量子AGI融合/通用量子", "智能文明", "全球垄断", "国际治理", "充足", "极高", "极高", "战略"],
                    ["2035", "85000亿", "37%", "AGI全面普及", "15000亿", "通用AGI/超级智能", "智能时代", "全球一体", "AI治理", "充足", "极高", "极高", "长期价值"],
                ]
                # 更新表格数据
                for i in range(min(len(data), len(table.rows))):
                    for j in range(min(len(data[i]), len(table.columns))):
                        fill_cell(table.cell(i, j), data[i][j], 12 if i == 0 else 9, i == 0)
                print("✅ 已更新2026-2035年完整技术突破数据")
                break

def main():
    ppt_path = r"f:\个人作品\legal-rag-qa-system\A2A_PROTOCOL_AI_AGENT_V14_终极完整版.pptx"
    prs = Presentation(ppt_path)
    
    print("🔄 正在补充2030-2035年大模型和技术突破...")
    update_tech_breakthroughs(prs)
    
    output_path = r"f:\个人作品\legal-rag-qa-system\A2A_PROTOCOL_AI_AGENT_V14_终极完整版.pptx"
    prs.save(output_path)
    print(f"✅ 已完成！保存到: {output_path}")
    
    print("\n📌 2030-2035年技术突破详细信息:")
    print("   2030年: GPT 7.0 / 通用AGI原型")
    print("   2031年: 量子AI融合 / IBM Q1000")
    print("   2032年: GPT 8.0 / 边缘大脑")
    print("   2033年: DeepSeek V7.0 / 脑机接口2.0")
    print("   2034年: 量子AGI融合 / 通用量子")
    print("   2035年: 通用AGI / 超级智能")

if __name__ == "__main__":
    main()