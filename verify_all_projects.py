from pptx import Presentation

def verify_all_projects():
    print("=" * 120)
    print("✅ 验证所有项目修复结果")
    print("=" * 120)
    
    ppt_path = r"f:\个人作品\legal-rag-qa-system\docs\presentations\A2A_ENTERPRISE_PPT_V12_完整修复版.pptx"
    prs = Presentation(ppt_path)
    
    print(f"幻灯片总数: {len(prs.slides)}")
    
    for slide_idx, slide in enumerate(prs.slides):
        print(f"\n📄 幻灯片 {slide_idx + 1}:")
        
        # 获取幻灯片标题
        title = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text.strip()
                if text:
                    title = text[:30]
                    break
        
        print(f"   标题: {title}")
        
        # 查找表格
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                print(f"   表格: {len(table.rows)}行 × {len(table.columns)}列")
                
                # 打印第一列内容
                first_col = []
                for i in range(len(table.rows)):
                    cell_text = table.cell(i, 0).text.strip()
                    first_col.append(cell_text)
                print(f"   内容: {first_col}")
                break

if __name__ == "__main__":
    verify_all_projects()