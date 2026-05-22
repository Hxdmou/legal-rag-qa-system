from pptx import Presentation
import os

def check_file_info():
    ppt_path = r"f:\个人作品\legal-rag-qa-system\A2A_PROTOCOL_AI_AGENT_V14_终极完整版.pptx"
    
    print("=" * 120)
    print("📊 文件信息检查")
    print("=" * 120)
    
    # 文件是否存在
    if os.path.exists(ppt_path):
        print(f"✅ 文件存在: {ppt_path}")
        file_size = os.path.getsize(ppt_path)
        print(f"📁 文件大小: {file_size} 字节")
        
        # 尝试打开文件
        try:
            prs = Presentation(ppt_path)
            print(f"✅ 成功打开PPT")
            print(f"📄 幻灯片数: {len(prs.slides)}")
            
            # 检查项目3
            if len(prs.slides) > 2:
                slide = prs.slides[2]
                for shape in slide.shapes:
                    if shape.has_table:
                        table = shape.table
                        print(f"\n📌 项目3表格: {len(table.rows)}行 × {len(table.columns)}列")
                        
                        # 列出所有层级
                        layers = []
                        for i in range(1, len(table.rows)):
                            layer_id = table.cell(i, 0).text.strip()
                            layer_name = table.cell(i, 1).text.strip()
                            layers.append(f"{layer_id}: {layer_name}")
                        
                        print(f"\n层级列表:")
                        for layer in layers:
                            print(f"  {layer}")
                        
                        # 检查是否有L10
                        has_l10 = any("L10" in table.cell(i, 0).text for i in range(len(table.rows)))
                        if has_l10:
                            print("\n✅ 包含L10边缘层！")
                        else:
                            print("\n❌ 缺少L10边缘层！")
                        
                        # 检查L9的内容
                        for i in range(len(table.rows)):
                            if "L9" in table.cell(i, 0).text:
                                print(f"\nL9内容:")
                                for j in range(min(5, len(table.columns))):
                                    print(f"  列{j+1}: {table.cell(i, j).text}")
                        
                        break
            
        except Exception as e:
            print(f"❌ 打开文件失败: {e}")
    else:
        print(f"❌ 文件不存在: {ppt_path}")

def check_file_modification():
    """检查文件修改时间"""
    ppt_path = r"f:\个人作品\legal-rag-qa-system\A2A_PROTOCOL_AI_AGENT_V14_终极完整版.pptx"
    
    if os.path.exists(ppt_path):
        import time
        mod_time = os.path.getmtime(ppt_path)
        print(f"\n🕐 修改时间: {time.ctime(mod_time)}")
        
        # 检查文件是否可写
        if os.access(ppt_path, os.W_OK):
            print("✅ 文件可写")
        else:
            print("❌ 文件不可写")

if __name__ == "__main__":
    check_file_info()
    check_file_modification()