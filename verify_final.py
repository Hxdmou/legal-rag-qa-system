from pptx import Presentation
import os
import time

def verify_final_version():
    print("=" * 120)
    print("✅ 验证最终版PPT")
    print("=" * 120)
    
    ppt_path = r"f:\个人作品\legal-rag-qa-system\docs\presentations\A2A_ENTERPRISE_PPT_V12_最终版.pptx"
    
    # 文件信息
    if os.path.exists(ppt_path):
        mod_time = os.path.getmtime(ppt_path)
        file_size = os.path.getsize(ppt_path)
        print(f"📁 文件路径: {ppt_path}")
        print(f"🕐 修改时间: {time.ctime(mod_time)}")
        print(f"📊 文件大小: {file_size} 字节")
    else:
        print(f"❌ 文件不存在")
        return
    
    # 内容检查
    prs = Presentation(ppt_path)
    print(f"\n📄 幻灯片总数: {len(prs.slides)}")
    
    for slide_idx, slide in enumerate(prs.slides):
        title = slide.shapes.title.text if slide.shapes.title else "无标题"
        print(f"\n📌 {title}")
        
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                print(f"   表格: {len(table.rows)}行 × {len(table.columns)}列")
                # 打印第一行标题
                headers = [table.cell(0, j).text.strip() for j in range(min(5, len(table.columns)))]
                print(f"   表头: {headers}...")
                break

def list_all_ppts():
    print("\n" + "=" * 120)
    print("📁 所有PPT文件列表")
    print("=" * 120)
    
    import glob
    ppt_files = glob.glob(r"f:\个人作品\legal-rag-qa-system\**\*.pptx", recursive=True)
    
    for ppt in ppt_files:
        mod_time = time.ctime(os.path.getmtime(ppt))
        size = os.path.getsize(ppt)
        print(f"\n✅ {os.path.basename(ppt)}")
        print(f"   路径: {ppt}")
        print(f"   修改时间: {mod_time}")
        print(f"   大小: {size} 字节")

if __name__ == "__main__":
    verify_final_version()
    list_all_ppts()